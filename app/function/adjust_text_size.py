"""
PPT文本框自适应大小调整模块
支持多种方法和跨平台兼容性，优先使用LibreOffice触发渲染
"""
import os
import logging
import platform
import subprocess
from typing import Optional, Dict, Any

# 尝试导入各种依赖
try:
    from pptx import Presentation
    from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Inches, Pt, Cm
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.error("python-pptx不可用")

try:
    import win32com.client
    import pythoncom
    COM_AVAILABLE = True
except ImportError:
    COM_AVAILABLE = False
    logging.warning("win32com不可用，将使用python-pptx作为替代方案")

# 配置日志
logger = logging.getLogger(__name__)


def set_file_permissions(file_path):
    """设置文件权限（仅Windows）"""
    if platform.system() != "Windows":
        return {"status": "success", "message": "非Windows系统，跳过权限设置"}
    
    if not os.path.exists(file_path):
        return {"status": "error", "message": f"文件 {file_path} 不存在。"}

    try:
        cmd = f'icacls "{file_path}" /grant Everyone:F'
        subprocess.check_call(cmd, shell=True)
        return {"status": "success", "message": f"已授予 {file_path} 的完全控制权限。"}
    except subprocess.CalledProcessError as e:
        return {"status": "error", "message": f"设置权限时出错：{e}"}


def set_textbox_autofit_pptx(ppt_path: str) -> bool:
    """使用python-pptx设置文本框自适应"""
    if not PPTX_AVAILABLE:
        logging.error("python-pptx库不可用")
        return False
    
    try:
        prs = Presentation(ppt_path)
        processed_count = 0
        
        logging.info(f"开始处理PPT文件: {os.path.basename(ppt_path)}")
        logging.info(f"总共 {len(prs.slides)} 张幻灯片")
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    text_frame.word_wrap = True
                    processed_count += 1
                elif shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            text_frame = cell.text_frame
                            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                            text_frame.word_wrap = True
                            processed_count += 1
        
        prs.save(ppt_path)
        logging.info(f"处理完成: {len(prs.slides)} 张幻灯片, {processed_count}/{processed_count} 个文本框")
        return True
        
    except Exception as e:
        logging.error(f"python-pptx处理失败: {e}")
        return False


import os
import gc
import logging
import platform

try:
    import pythoncom
    import win32com.client as win32
    from win32com.client import constants
except Exception:
    pass  # 维持你原来的 COM_AVAILABLE 逻辑

def _iter_shapes(shape_or_shapes):
    """递归遍历所有 shape（包含组合里的成员）"""
    try:
        for sh in shape_or_shapes:
            try:
                if getattr(sh, "Type", None) == getattr(constants, "msoGroup", 6):
                    # 组合内递归
                    for inner in _iter_shapes(sh.GroupItems):
                        yield inner
                    continue
            except Exception:
                pass
            yield sh
    except TypeError:
        # 单个 shape
        sh = shape_or_shapes
        try:
            if getattr(sh, "Type", None) == getattr(constants, "msoGroup", 6):
                for inner in _iter_shapes(sh.GroupItems):
                    yield inner
            else:
                yield sh
        except Exception:
            yield sh

def _is_textual_shape(shape) -> bool:
    """严格判断是否是值得处理的文本形状"""
    try:
        if getattr(shape, "HasTextFrame", 0) != -1:
            return False
        tf2 = shape.TextFrame2
        if getattr(tf2, "HasText", 0) == 0:
            return False
        t = getattr(shape, "Type", None)
        msoTable = getattr(constants, "msoTable", 19)
        msoSmartArt = getattr(constants, "msoSmartArt", 24)
        msoPlaceholder = getattr(constants, "msoPlaceholder", 14)
        if t in (msoTable, msoSmartArt, msoPlaceholder):
            return False
        return True
    except Exception:
        return False

def set_textbox_autofit_com(ppt_path: str) -> bool:
    """使用COM接口设置文本框自适应（Windows专用，包含轻微扰动触发渲染）"""
    if not COM_AVAILABLE:
        logging.error("COM接口不可用")
        return False

    if platform.system() != "Windows":
        logging.error("COM接口仅在Windows上可用")
        return False

    # 设置文件权限（保留你的逻辑）
    permission_result = set_file_permissions(ppt_path)
    if permission_result.get("status") == "error":
        logging.warning(f"权限设置失败: {permission_result.get('message')}")

    app = None
    presentation = None

    # --- COM 初始化（优先 STA）---
    try:
        pythoncom.CoInitializeEx(pythoncom.COINIT_APARTMENTTHREADED)
    except Exception:
        pythoncom.CoInitialize()

    try:
        # 连接或创建PowerPoint实例
        try:
            app = win32.GetActiveObject("PowerPoint.Application")
            logging.debug("连接到现有PowerPoint实例")
        except Exception:
            app = win32.Dispatch("PowerPoint.Application")
            logging.debug("创建新的PowerPoint实例")

        # 静默、屏蔽弹窗
        try:
            app.Visible = True
        except Exception as e:
            logging.warning(f"设置Visible属性失败: {e}")
        try:
            app.DisplayAlerts = constants.ppAlertsNone
        except Exception:
            pass

        # 打开演示文稿（后台窗口）
        presentation = app.Presentations.Open(ppt_path, WithWindow=True)
        logging.info(f"成功打开演示文稿: {os.path.basename(ppt_path)}")

        slide_count = 0
        handled_count = 0
        skipped_count = 0

        for slide in presentation.Slides:
            slide_count += 1
            for shape in _iter_shapes(slide.Shapes):
                if not _is_textual_shape(shape):
                    skipped_count += 1
                    continue
                try:
                    tf2 = shape.TextFrame2

                    # 记录原值，必要时可回滚
                    orig_wrap = getattr(tf2, "WordWrap", None)

                    # 设置自适应：文字大小适应文本框（只缩小不放大）
                    tf2.AutoSize = constants.msoAutoSizeTextToFitShape  # 2

                    # 换行策略：一般保留原值即可；若想强制换行可改为 True(-1)
                    if orig_wrap is None:
                        tf2.WordWrap = -1  # True

                    # 轻微扰动：移动位置而非改尺寸，避免布局抖动
                    ol, ot = shape.Left, shape.Top
                    shape.Left = ol + 0.1
                    shape.Top = ot + 0.1
                    shape.Left = ol
                    shape.Top = ot

                    handled_count += 1
                except Exception as e:
                    logging.debug(f"处理形状时出错: {e}")
                    skipped_count += 1

        logging.info(
            f"COM处理完成: {slide_count} 张幻灯片, 处理 {handled_count} 个文本框, 跳过 {skipped_count} 个"
        )

        # 保存并关闭
        presentation.Save()
        presentation.Close()
        presentation = None

        # 智能退出PowerPoint
        try:
            if app.Presentations.Count == 0:
                app.Quit()
                app = None
                logging.debug("PowerPoint实例已退出")
        except Exception as e:
            logging.debug(f"退出PowerPoint时出错: {e}")

        return True

    except Exception as e:
        logging.error(f"COM处理失败: {e}")

        # 清理资源
        try:
            if presentation:
                presentation.Close()
        except Exception:
            pass

        try:
            if app and app.Presentations.Count == 0:
                app.Quit()
        except Exception:
            pass

        return False

    finally:
        # 彻底释放 COM 对象，避免残留进程
        try:
            del presentation
        except Exception:
            pass
        try:
            del app
        except Exception:
            pass
        gc.collect()
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass


def set_textbox_autofit(ppt_path):
    """
    调整PPT文本框自适应大小（带颜色保护）

    Args:
        ppt_path: PPT文件的绝对路径

    Returns:
        bool: 调整是否成功
    """
    if not os.path.exists(ppt_path):
        logging.error(f"文件路径 {ppt_path} 不存在")
        return False

    logging.info(f"开始调整PPT文本框自适应: {os.path.basename(ppt_path)}")

    # 方法1: 使用带颜色保护的LibreOffice渲染（推荐）
    try:
        from .color_backup_restore import render_with_color_protection
        logging.info("使用带颜色保护的LibreOffice渲染")
        result = render_with_color_protection(ppt_path)
        if result:
            logging.info("✅ 带颜色保护的LibreOffice渲染成功")
            return True
        else:
            logging.warning("带颜色保护的LibreOffice渲染失败，尝试基础方法")
    except ImportError:
        logging.warning("颜色保护模块不可用，使用基础LibreOffice渲染")
    except Exception as e:
        logging.warning(f"颜色保护渲染出错: {e}，尝试基础方法")

    # 方法2: 基础LibreOffice渲染（降级方案）
    try:
        from .libreoffice_autofit import libreoffice_ppt_autofit
        logging.info("使用基础LibreOffice PDF转换触发渲染")
        result = libreoffice_ppt_autofit(ppt_path)
        if result:
            logging.info("基础LibreOffice渲染触发成功")
            return True
        else:
            logging.warning("基础LibreOffice渲染触发失败，尝试其他方法")
    except ImportError as e:
        logging.warning(e)
        logging.warning("LibreOffice渲染触发器不可用，尝试其他方法")
    except Exception as e:
        logging.warning(f"基础LibreOffice渲染触发出错: {e}，尝试其他方法")
    
    # 方法2: Windows环境 - 使用COM接口（包含扰动触发渲染）
    if platform.system() == "Windows" and COM_AVAILABLE:
        logging.info("Windows环境: 使用COM接口方法（模拟PPT客户端扰动）")
        result = set_textbox_autofit_com(ppt_path)
        if result:
            logging.info("COM接口方法处理成功")
            return True
        else:
            logging.warning("COM接口方法处理失败，尝试其他方法")
    
    # 方法3: 跨平台纯python-pptx方法（备用方案）
    try:
        from .pure_pptx_autofit import pure_pptx_set_textbox_autofit
        logging.info("使用纯python-pptx自适应方法处理")
        result = pure_pptx_set_textbox_autofit(ppt_path)
        if result:
            logging.info("纯python-pptx方法处理成功")
            return True
        else:
            logging.warning("纯python-pptx方法处理失败，尝试其他方法")
    except ImportError:
        logging.warning("纯python-pptx方法不可用，尝试其他方法")
    except Exception as e:
        logging.warning(f"纯python-pptx方法出错: {e}，尝试其他方法")
    
    # 方法4: 标准python-pptx方法（最后的保障）
    if PPTX_AVAILABLE:
        logging.info("使用标准python-pptx方法处理")
        result = set_textbox_autofit_pptx(ppt_path)
        if result:
            logging.info("标准python-pptx方法处理成功")
            return True
        else:
            logging.warning("标准python-pptx方法处理失败")
    
    logging.error("所有方法都失败了")
    return False


# 为了兼容性，保留原有的函数名
def set_textbox_autofit_legacy(ppt_path):
    """兼容性函数，调用新的set_textbox_autofit"""
    return set_textbox_autofit(ppt_path)
