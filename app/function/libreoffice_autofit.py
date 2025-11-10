"""
使用LibreOffice命令行触发PPT文本框自适应渲染
通过转换PPT为PDF的过程触发完整渲染，使文本框自适应设置真正生效
"""
import os
import logging
import shutil
import subprocess
import tempfile
import platform
try:
    from pptx import Presentation
    from pptx.enum.text import MSO_AUTO_SIZE
    from pptx.util import Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.error("python-pptx不可用")

logger = logging.getLogger(__name__)


def find_libreoffice_command():
    """直接返回已知的 LibreOffice 安装路径"""
    if platform.system() == "Windows":
        return r"C:\Program Files\LibreOffice\program\soffice.exe"
    else:
        # Linux/macOS 可继续使用自动路径探测
        possible_paths = [
            "libreoffice",
            "soffice",
            "/usr/bin/libreoffice",
            "/opt/libreoffice/program/soffice",
            "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        ]
        for cmd_path in possible_paths:
            if os.path.exists(cmd_path):
                return cmd_path

    return None


def set_ppt_autofit_properties(ppt_path: str) -> bool:
    """设置PPT文本框自适应属性"""
    if not PPTX_AVAILABLE:
        logger.error("python-pptx库不可用")
        return False
    
    try:
        prs = Presentation(ppt_path)
        processed_count = 0
        
        logger.info(f"设置文本框自适应属性: {os.path.basename(ppt_path)}")
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    
                    # 检查是否有实际内容
                    has_content = any(
                        paragraph.text.strip() 
                        for paragraph in text_frame.paragraphs
                    )
                    
                    if has_content:
                        # 使用颜色保护器保存原始格式
                        from .color_protection import save_textframe_colors
                        color_info = save_textframe_colors(text_frame, f"shape_{processed_count}")

                        # 设置文本大小适应文本框
                        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                        text_frame.word_wrap = True

                        # 优化边距
                        text_frame.margin_left = Pt(3)
                        text_frame.margin_right = Pt(3)
                        text_frame.margin_top = Pt(2)
                        text_frame.margin_bottom = Pt(2)

                        # 使用颜色保护器恢复原始格式
                        from .color_protection import restore_textframe_colors
                        restore_textframe_colors(text_frame, color_info)
                        
                        processed_count += 1
                
                elif shape.has_table:
                    # 处理表格中的文本框
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            text_frame = cell.text_frame
                            
                            has_content = any(
                                paragraph.text.strip() 
                                for paragraph in text_frame.paragraphs
                            )
                            
                            if has_content:
                                # 使用颜色保护器保存表格单元格格式
                                from .color_protection import save_textframe_colors, restore_textframe_colors
                                color_info = save_textframe_colors(text_frame, f"table_cell_{processed_count}")

                                # 设置自适应
                                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                                text_frame.word_wrap = True

                                # 恢复原始字体格式
                                restore_textframe_colors(text_frame, color_info)

                                processed_count += 1
        
        # 保存设置
        prs.save(ppt_path)
        logger.info(f"设置了 {processed_count} 个文本框的自适应属性")
        return True
        
    except Exception as e:
        logger.error(f"设置自适应属性时出错: {e}")
        return False


def trigger_ppt_render_with_libreoffice(ppt_path: str) -> bool:
    """
    使用LibreOffice渲染并导出渲染后的PPTX，覆盖原文件
    使用.odp中转转换并显式指定导出格式

    Args:
        ppt_path: PPT文件路径

    Returns:
        bool: 处理是否成功
    """
    if not os.path.exists(ppt_path):
        logger.error(f"文件不存在: {ppt_path}")
        return False

    libreoffice_cmd = find_libreoffice_command()
    if not libreoffice_cmd:
        return False

    logger.info(f"开始LibreOffice渲染PPTX: {os.path.basename(ppt_path)}")

    try:
        # 步骤1: 设置文本框自适应属性
        if not set_ppt_autofit_properties(ppt_path):
            logger.error("设置自适应属性失败")
            return False

        # 步骤2: 使用LibreOffice渲染并重新导出PPTX
        # logger.info("使用LibreOffice通过.odp中转转换并重新保存PPTX以触发文本框渲染...")

        # with tempfile.TemporaryDirectory() as temp_dir:
        #     temp_dir = os.path.abspath(temp_dir)
        #     ppt_path_abs = os.path.abspath(ppt_path)
        #     ppt_filename = os.path.basename(ppt_path)
        #     ppt_name = os.path.splitext(ppt_filename)[0]

        #     # 步骤2.1: 先转换为ODP格式
        #     if platform.system() == "Windows":
        #         cmd_to_odp = [
        #             libreoffice_cmd,
        #             "--headless", "--convert-to", "odp",
        #             "--outdir", temp_dir,
        #             ppt_path_abs
        #         ]
        #         creation_flags = subprocess.CREATE_NO_WINDOW
        #     else:
        #         cmd_to_odp = [
        #             libreoffice_cmd,
        #             "--headless", "--convert-to", "odp",
        #             "--outdir", temp_dir,
        #             ppt_path_abs
        #         ]
        #         creation_flags = 0

        #     logger.debug(f"执行ODP转换命令: {' '.join(cmd_to_odp)}")

        #     result_odp = subprocess.run(
        #         cmd_to_odp,
        #         capture_output=True,
        #         text=True,
        #         timeout=120,
        #         creationflags=creation_flags
        #     )

        #     if result_odp.returncode != 0:
        #         logger.error(f"转换为ODP失败: {result_odp.stderr}")
        #         return False

        #     # 检查ODP文件是否生成
        #     odp_file = os.path.join(temp_dir, f"{ppt_name}.odp")
        #     if not os.path.exists(odp_file):
        #         logger.error("未能生成ODP文件")
        #         return False

        #     logger.info(f"成功转换为ODP格式: {odp_file}")

        #     # 步骤2.2: 从ODP转换回PPTX格式
        #     if platform.system() == "Windows":
        #         cmd_to_pptx = [
        #             libreoffice_cmd,
        #             "--headless", "--convert-to", "pptx:Impress MS PowerPoint 2007 XML",
        #             "--outdir", temp_dir,
        #             odp_file
        #         ]
        #     else:
        #         cmd_to_pptx = [
        #             libreoffice_cmd,
        #             "--headless", "--convert-to", "pptx:Impress MS PowerPoint 2007 XML",
        #             "--outdir", temp_dir,
        #             odp_file
        #         ]

        #     logger.debug(f"执行PPTX转换命令: {' '.join(cmd_to_pptx)}")

        #     result_pptx = subprocess.run(
        #         cmd_to_pptx,
        #         capture_output=True,
        #         text=True,
        #         timeout=120,
        #         creationflags=creation_flags
        #     )

        #     logger.debug(f"PPTX转换命令返回码: {result_pptx.returncode}")
        #     if result_pptx.stdout:
        #         logger.debug(f"标准输出: {result_pptx.stdout}")
        #     if result_pptx.stderr:
        #         logger.debug(f"标准错误: {result_pptx.stderr}")

        #     if result_pptx.returncode == 0:
        #         # 查找生成的PPTX文件
        #         pptx_file = os.path.join(temp_dir, f"{ppt_name}.pptx")
        #         # print(pptx_file)
        #         if os.path.exists(pptx_file):
        #             try:
        #                 shutil.copyfile(pptx_file, ppt_path)
        #                 logger.info(f" 渲染后的PPTX已复制并覆盖原文件: {ppt_path}")
        #                 return True
        #             except Exception as copy_err:
        #                 logger.error(f" 无法覆盖源文件: {copy_err}")
        #                 return False
        #         else:
        #             logger.warning("LibreOffice成功执行但未找到输出PPTX文件")
        #             logger.debug(f"临时目录内容: {os.listdir(temp_dir)}")
        #             return False
        #     else:
        #         logger.error("LibreOffice转换失败")
        #         return False

    except subprocess.TimeoutExpired:
        logger.error("LibreOffice转换超时")
        return False
    except Exception as e:
        logger.error(f"处理过程中出错: {e}")
        import traceback
        logger.error(f"详细错误: {traceback.format_exc()}")
        return False

# 简化的主函数接口
def libreoffice_ppt_autofit(ppt_path: str) -> bool:
    """
    使用LibreOffice实现PPT文本框自适应的主函数
    
    这个函数会：
    1. 使用python-pptx设置文本框自适应属性
    2. 使用LibreOffice命令行转换PPT为PDF触发渲染
    3. 自动删除临时PDF文件
    4. 返回处理结果
    
    Args:
        ppt_path: PPT文件路径
        
    Returns:
        bool: 处理是否成功
    """
    return trigger_ppt_render_with_libreoffice(ppt_path)


# 测试函数
def test_libreoffice_autofit():
    """测试LibreOffice自适应功能"""
    # 检查LibreOffice是否可用
    libreoffice_cmd = find_libreoffice_command()
    if libreoffice_cmd:
        print(f" LibreOffice可用: {libreoffice_cmd}")
        return True
    else:
        print(" LibreOffice不可用")
        print("\n安装说明:")
        if platform.system() == "Windows":
            print("Windows: 下载并安装LibreOffice到默认路径")
            print("https://www.libreoffice.org/download/")
        else:
            print("Linux: sudo apt-get install libreoffice")
            print("或者: sudo yum install libreoffice-headless")
        return False


if __name__ == "__main__":
    # 测试LibreOffice可用性
    test_libreoffice_autofit()
