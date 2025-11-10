#!/usr/bin/env python3
"""
基于页面的翻译机制
每页PPT调用一次API，按段落匹配翻译结果
"""
import re
import logging
import difflib
import asyncio
from typing import Dict, List, Tuple, Optional, Any
from dataclasses import dataclass
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE

logger = logging.getLogger(__name__)

@dataclass
class ParagraphInfo:
    """段落信息数据结构"""
    text: str                    # 段落文本
    slide_index: int            # 幻灯片索引
    shape_index: int            # 形状索引
    paragraph_index: int        # 段落索引
    shape_type: str             # 形状类型 ('textbox' 或 'table')
    table_position: Optional[Tuple[int, int]] = None  # 表格位置 (row, col)
    length: int = 0             # 文本长度
    is_translatable: bool = True # 是否需要翻译

    def __post_init__(self):
        self.length = len(self.text)
def clean_text_for_ppt(text):
    # 删除所有不可见的控制字符（ASCII < 32 且不是常见换行）
    text = re.sub(r'[\x00-\x08\x0b-\x1f\x7f]', '', text)
    return text
def clean_brackets(text):
    """
    去除文本中的【和】符号
    """
    return text.replace('【', '').replace('】', '')
class PageBasedTranslator:
    """基于页面的翻译器"""

    def __init__(self):
        self.current_slide_paragraphs: List[ParagraphInfo] = []
        self.translation_results: List[str] = []

    def is_translatable_text(self, text: str) -> bool:
        """判断文本是否需要翻译"""
        if not text or len(text.strip()) == 0:
            return False

        text = text.strip()

        # 跳过纯数字（包括小数点、百分号、连字符等）
        if re.match(r'^[\d\s\.,\-%]+$', text):
            logger.debug(f"跳过纯数字: '{text}'")
            return False

        # 跳过纯标点符号
        if re.match(r'^[^\w\s]+$', text):
            logger.debug(f"跳过纯标点: '{text}'")
            return False

        # 跳过单个字符
        if len(text) <= 1:
            logger.debug(f"跳过单字符: '{text}'")
            return False

        # 跳过纯空格或特殊字符
        if re.match(r'^[\s\-_=+\*#@$%^&()]+$', text):
            logger.debug(f"跳过特殊字符: '{text}'")
            return False

        return True

    def collect_slide_paragraphs(self, slide, slide_index: int) -> List[ParagraphInfo]:
        """收集单页幻灯片的所有段落"""
        paragraphs = []

        # 处理普通文本框
        for shape_index, shape in enumerate(slide.shapes):
            #     if shape.has_text_frame:
            #         text_frame = shape.text_frame
            #         for paragraph_index, paragraph in enumerate(text_frame.paragraphs):
            #             text = paragraph.text.strip()
            #             if text:
            #                 is_translatable = self.is_translatable_text(text)
            #                 para_info = ParagraphInfo(
            #                     text=text,
            #                     slide_index=slide_index,
            #                     shape_index=shape_index,
            #                     paragraph_index=paragraph_index,
            #                     shape_type='textbox',
            #                     is_translatable=is_translatable
            #                 )
            #                 paragraphs.append(para_info)

            #                 if is_translatable:
            #                     logger.debug(f"收集可翻译段落: '{text[:50]}...'")
            #                 else:
            #                     logger.debug(f"收集不可翻译段落: '{text[:50]}...'")

            # 处理表格
            if shape.has_table:
                table = shape.table
                for row_index, row in enumerate(table.rows):
                    for col_index, cell in enumerate(row.cells):
                        text_frame = cell.text_frame
                        for paragraph_index, paragraph in enumerate(text_frame.paragraphs):
                            text = paragraph.text.strip()
                            if text:
                                is_translatable = self.is_translatable_text(text)
                                para_info = ParagraphInfo(
                                    text=text,
                                    slide_index=slide_index,
                                    shape_index=shape_index,
                                    paragraph_index=paragraph_index,
                                    shape_type='table',
                                    table_position=(row_index, col_index),
                                    is_translatable=is_translatable
                                )
                                paragraphs.append(para_info)

                                if is_translatable:
                                    logger.debug(f"收集可翻译表格段落: '{text[:50]}...'")
                                else:
                                    logger.debug(f"收集不可翻译表格段落: '{text[:50]}...'")

        return paragraphs

    def prepare_translation_text(self, paragraphs: List[ParagraphInfo]) -> str:
        """准备翻译文本（只包含可翻译的段落）"""
        translatable_texts = []

        for para in paragraphs:
            if para.is_translatable:
                translatable_texts.append("【"+para.text+"】")

        # 用换行符连接所有段落
        combined_text = '\n'.join(translatable_texts)

        logger.info(f"准备翻译文本: {len(translatable_texts)} 个段落，共 {len(combined_text)} 个字符")
        return combined_text



    def calculate_similarity_score(self, text1: str, text2: str) -> float:
        """计算两个文本的相似度分数"""
        # 长度相似度 (权重: 0.3)
        len1, len2 = len(text1), len(text2)
        length_similarity = 1.0 - abs(len1 - len2) / max(len1, len2, 1)

        # 文本相似度 (权重: 0.7)
        text_similarity = difflib.SequenceMatcher(None, text1.lower(), text2.lower()).ratio()

        # 综合相似度
        total_similarity = length_similarity * 0.3 + text_similarity * 0.7

        return total_similarity



    def _normalize_text(self, text: str) -> str:
        """标准化文本用于比较"""
        if not text:
            return ""

        # 转换为小写
        normalized = text.lower()

        # 移除标点符号和空格
        normalized = re.sub(r'[^\w]', '', normalized)

        return normalized



    def match_translations_to_paragraphs_precise(self, paragraphs: List[ParagraphInfo],
                                               translation_dict: Dict[str, str]) -> Dict[int, str]:
        """
        精确匹配翻译到段落
        使用翻译字典中的原文（source_language）与PPT段落匹配，
        匹配成功后使用译文（target_language）作为翻译结果

        Args:
            paragraphs: 段落信息列表
            translation_dict: 翻译字典 {原文: 译文}

        Returns:
            匹配结果字典 {段落索引: 译文}
        """
        # 获取可翻译的段落
        translatable_paragraphs = [para for para in paragraphs if para.is_translatable]

        if not translatable_paragraphs:
            logger.info("没有可翻译的段落")
            return {}

        logger.info(f"开始精确匹配: {len(translatable_paragraphs)} 个段落 vs {len(translation_dict)} 条翻译")

        # 创建匹配结果字典 (段落索引 -> 译文)
        matches = {}
        used_translations = set()  # 已使用的原文

        # 策略1: 精确匹配
        exact_matches = 0
        for para_idx, para in enumerate(translatable_paragraphs):
            original_text = para.text.strip()

            # 直接在翻译字典中查找
            if original_text in translation_dict:
                translation = translation_dict[original_text]
                if original_text not in used_translations:
                    # 找到原始段落在所有段落中的索引
                    original_para_idx = paragraphs.index(para)
                    matches[original_para_idx] = translation
                    used_translations.add(original_text)
                    exact_matches += 1
                    logger.debug(f"✓ 精确匹配: '{original_text[:30]}...' -> '{translation[:30]}...'")

        logger.info(f"精确匹配完成: {exact_matches} 个段落")

        # 策略2: 标准化匹配（去除空格、标点等）
        normalized_matches = 0
        if exact_matches < len(translatable_paragraphs):
            logger.info("开始标准化匹配...")

            # 创建标准化映射
            normalized_translation_dict = {}
            for orig_text, trans_text in translation_dict.items():
                normalized_orig = self._normalize_text(orig_text)
                if normalized_orig and normalized_orig not in normalized_translation_dict:
                    normalized_translation_dict[normalized_orig] = (orig_text, trans_text)

            for para_idx, para in enumerate(translatable_paragraphs):
                # 跳过已匹配的段落
                original_para_idx = paragraphs.index(para)
                if original_para_idx in matches:
                    continue

                original_text = para.text.strip()
                normalized_original = self._normalize_text(original_text)

                if normalized_original in normalized_translation_dict:
                    orig_text, translation = normalized_translation_dict[normalized_original]
                    if orig_text not in used_translations:
                        matches[original_para_idx] = translation
                        used_translations.add(orig_text)
                        normalized_matches += 1
                        logger.debug(f"✓ 标准化匹配: '{original_text[:30]}...' -> '{translation[:30]}...'")

        logger.info(f"标准化匹配完成: {normalized_matches} 个段落")

        # 策略3: 相似度匹配（最后的备选方案）
        similarity_matches = 0
        remaining_paragraphs = []
        for para_idx, para in enumerate(translatable_paragraphs):
            original_para_idx = paragraphs.index(para)
            if original_para_idx not in matches:
                remaining_paragraphs.append((para_idx, para, original_para_idx))

        if remaining_paragraphs and len(used_translations) < len(translation_dict):
            logger.info(f"开始相似度匹配: {len(remaining_paragraphs)} 个剩余段落")

            # 获取未使用的翻译
            unused_translations = []
            for orig_text, trans_text in translation_dict.items():
                if orig_text not in used_translations:
                    unused_translations.append((orig_text, trans_text))

            for para_idx, para, original_para_idx in remaining_paragraphs:
                original_text = para.text.strip()
                best_score = 0.0
                best_translation = None
                best_orig_text = None

                for orig_text, trans_text in unused_translations:
                    if orig_text in used_translations:
                        continue

                    score = self.calculate_similarity_score(original_text, orig_text)
                    if score > best_score and score > 0.3:  # 相似度阈值
                        best_score = score
                        best_translation = trans_text
                        best_orig_text = orig_text

                if best_translation and best_orig_text:
                    matches[original_para_idx] = best_translation
                    used_translations.add(best_orig_text)
                    similarity_matches += 1
                    logger.debug(f"✓ 相似度匹配 (相似度: {best_score:.3f}): '{original_text[:30]}...' -> '{best_translation[:30]}...'")

        logger.info(f"相似度匹配完成: {similarity_matches} 个段落")

        total_matches = exact_matches + normalized_matches + similarity_matches
        logger.info(f"精确匹配总结: {total_matches}/{len(translatable_paragraphs)} 个段落成功匹配")
        logger.info(f"  - 精确匹配: {exact_matches}")
        logger.info(f"  - 标准化匹配: {normalized_matches}")
        logger.info(f"  - 相似度匹配: {similarity_matches}")

        return matches

    async def translate_slide_paragraphs(self, slide, slide_index: int,
                                       source_language: str, target_language: str,
                                       field: str = "通用") -> Dict[int, str]:
        """翻译单页幻灯片的段落"""
        # 1. 收集段落
        paragraphs = self.collect_slide_paragraphs(slide, slide_index)
        if not paragraphs:
            logger.info(f"第 {slide_index + 1} 页没有可处理的段落")
            return {}

        # 2. 准备翻译文本
        translation_text = self.prepare_translation_text(paragraphs)
        if not translation_text:
            logger.info(f"第 {slide_index + 1} 页没有需要翻译的内容")
            return {}

        # 3. 调用翻译API
        try:
            from .local_qwen_async import translate_async

            logger.info(f"开始翻译第 {slide_index + 1} 页内容...")

            # 构造翻译参数，PPT翻译不需要清理Markdown
            translation_result = await translate_async(
                text=translation_text,
                field=field,
                stop_words=[],  # 空的停止词列表
                custom_translations={},  # 空的自定义翻译
                source_language=source_language,
                target_language=target_language,
                clean_markdown=False  # PPT翻译不需要清理Markdown
            )

            if not translation_result:
                logger.error(f"第 {slide_index + 1} 页翻译失败：返回空结果")
                return {}

            # translate_async 返回的是字典格式 {原文: 译文}
            # 直接使用这个字典进行匹配
            logger.info(f"第 {slide_index + 1} 页翻译完成，获得 {len(translation_result)} 条翻译")

        except Exception as e:
            logger.error(f"第 {slide_index + 1} 页翻译API调用失败: {str(e)}")
            import traceback
            logger.error(f"错误详情: {traceback.format_exc()}")
            return {}

        # 4. 直接使用翻译字典进行匹配
        if not translation_result:
            logger.error(f"第 {slide_index + 1} 页翻译结果为空")
            return {}

        # 5. 匹配翻译到段落（使用原文匹配，应用译文）
        matches = self.match_translations_to_paragraphs_precise(paragraphs, translation_result)

        # 6. 存储段落信息供后续使用
        self.current_slide_paragraphs = paragraphs

        return matches

    def apply_translations_to_slide(self, slide, slide_index: int,
                                  matches: Dict[int, str], bilingual_translation: str = "1"):
        """将翻译结果应用到幻灯片"""
        if not matches:
            logger.info(f"第 {slide_index + 1} 页没有翻译结果需要应用")
            return 0

        applied_count = 0
        paragraphs = self.current_slide_paragraphs

        for para_idx, translation in matches.items():
            if para_idx >= len(paragraphs):
                logger.warning(f"段落索引超出范围: {para_idx}")
                continue

            para_info = paragraphs[para_idx]
            para_info.text=clean_text_for_ppt(para_info.text)
            translation = clean_text_for_ppt(translation)
            para_info.text = clean_brackets(para_info.text)
            translation = clean_brackets(translation)
            try:
                if para_info.shape_type == 'textbox':
                    # 处理普通文本框
                    shape = slide.shapes[para_info.shape_index]
                    if shape.has_text_frame:


                        text_frame = shape.text_frame
                        paragraph = text_frame.paragraphs[para_info.paragraph_index]

                        # 保存原始字体颜色（简化版）
                        original_color = None
                        if paragraph.runs:
                            try:
                                from ..utils.ppt_utils import get_font_color
                                original_color = get_font_color(paragraph.runs[0])
                            except:
                                original_color = None

                        # 检查相似度，如果相似度过高则跳过翻译
                        from .ppt_translate import should_skip_translation_insertion
                        if should_skip_translation_insertion(para_info.text, translation, threshold=0.85, debug=True):
                            logger.info(f"跳过高相似度翻译: '{para_info.text[:30]}...' -> '{translation[:30]}...'")
                            continue

                        # 应用翻译（使用安全的文本替换）
                        try:
                            from .color_protection import safe_replace_paragraph_text

                            if str(bilingual_translation) == "paragraph_up":
                                # 双语模式：原文再上 + 译文在下
                                new_text = para_info.text + "\n" + translation
                            elif str(bilingual_translation) == "paragraph_down":
                                # 双语模式：原文再下 + 译文在上
                                new_text = translation + "\n" + para_info.text
                            elif str(bilingual_translation) == "translation_only":
                                # 仅翻译模式
                                new_text = translation

                            # 使用安全替换，保持所有格式
                            success = safe_replace_paragraph_text(paragraph, new_text, preserve_formatting=True)
                            if not success:
                                # 如果安全替换失败，使用传统方法
                                paragraph.text = new_text
                                # 恢复字体颜色（简化版）
                                if original_color and paragraph.runs:
                                    try:
                                        from ..utils.ppt_utils import apply_font_color
                                        for run in paragraph.runs:
                                            apply_font_color(run, original_color)
                                    except Exception as color_error:
                                        logger.debug(f"恢复字体颜色失败: {color_error}")
                        except ImportError:
                            # 如果颜色保护模块不可用，使用传统方法
                            if str(bilingual_translation) == "1":
                                paragraph.text = para_info.text + "\n" + translation
                            else:
                                paragraph.text = translation

                            # 恢复字体颜色（简化版）
                            if original_color and paragraph.runs:
                                try:
                                    from ..utils.ppt_utils import apply_font_color
                                    for run in paragraph.runs:
                                        apply_font_color(run, original_color)
                                except Exception as color_error:
                                    logger.debug(f"恢复字体颜色失败: {color_error}")

                        # 设置文字大小适应文本框大小
                        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                        applied_count += 1
                        logger.debug(f"✓ 应用文本框翻译: '{para_info.text[:30]}...' -> '{translation[:30]}...'")
                        logger.debug(f"  自适应: 已设置为TEXT_TO_FIT_SHAPE")

                elif para_info.shape_type == 'table' and para_info.table_position:
                    # 处理表格
                    shape = slide.shapes[para_info.shape_index]
                    if shape.has_table:


                        table = shape.table
                        row_idx, col_idx = para_info.table_position

                        if row_idx < len(table.rows) and col_idx < len(table.rows[row_idx].cells):
                            cell = table.rows[row_idx].cells[col_idx]
                            text_frame = cell.text_frame
                            paragraph = text_frame.paragraphs[para_info.paragraph_index]

                            # 保存原始字体颜色（简化版）
                            original_color = None
                            if paragraph.runs:
                                try:
                                    from ..utils.ppt_utils import get_font_color
                                    original_color = get_font_color(paragraph.runs[0])
                                except:
                                    original_color = None

                            # 检查相似度，如果相似度过高则跳过翻译
                            from .ppt_translate import should_skip_translation_insertion
                            if should_skip_translation_insertion(para_info.text, translation, threshold=0.9, debug=True):
                                logger.info(f"跳过表格高相似度翻译: '{para_info.text[:30]}...' -> '{translation[:30]}...'")
                                continue

                            # 应用翻译（使用安全的文本替换）
                            try:
                                from .color_protection import safe_replace_paragraph_text

                                if str(bilingual_translation) == "paragraph_up":
                                # 双语模式：原文再上 + 译文在下
                                    new_text = para_info.text + "\n" + translation
                                elif str(bilingual_translation) == "paragraph_down":
                                # 双语模式：原文再下 + 译文在上
                                    new_text = translation + "\n" + para_info.text
                                elif str(bilingual_translation) == "translation_only":
                                # 仅翻译模式
                                    new_text = translation

                                # 使用安全替换，保持所有格式
                                success = safe_replace_paragraph_text(paragraph, new_text, preserve_formatting=True)
                                if not success:
                                    # 如果安全替换失败，使用传统方法
                                    paragraph.text = new_text
                                    # 恢复字体颜色（简化版）
                                    if original_color and paragraph.runs:
                                        try:
                                            from ..utils.ppt_utils import apply_font_color
                                            for run in paragraph.runs:
                                                apply_font_color(run, original_color)
                                        except Exception as color_error:
                                            logger.debug(f"恢复表格字体颜色失败: {color_error}")
                            except ImportError:
                                # 如果颜色保护模块不可用，使用传统方法
                                if str(bilingual_translation) == "1":
                                    paragraph.text = para_info.text + "\n" + translation
                                else:
                                    paragraph.text = translation

                                # 恢复字体颜色（简化版）
                                if original_color and paragraph.runs:
                                    try:
                                        from ..utils.ppt_utils import apply_font_color
                                        for run in paragraph.runs:
                                            apply_font_color(run, original_color)
                                    except Exception as color_error:
                                        logger.debug(f"恢复表格字体颜色失败: {color_error}")

                            # 设置文字大小适应文本框大小
                            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                            applied_count += 1
                            logger.debug(f"✓ 应用表格翻译: '{para_info.text[:30]}...' -> '{translation[:30]}...'")
                            logger.debug(f"  自适应: 已设置为TEXT_TO_FIT_SHAPE")

            except Exception as e:
                logger.error(f"应用翻译失败 (段落 {para_idx}): {str(e)}")

        # 确保所有文本框都设置了自动适应
        self.ensure_all_textboxes_autofit(slide)

        logger.info(f"第 {slide_index + 1} 页翻译应用完成: {applied_count} 个段落")
        return applied_count

    def ensure_all_textboxes_autofit(self, slide):
        """确保幻灯片中所有文本框都设置了文字大小适应文本框大小"""
        textbox_count = 0

        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                textbox_count += 1
            elif shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        text_frame = cell.text_frame
                        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                        textbox_count += 1

        logger.debug(f"设置了 {textbox_count} 个文本框为文字大小适应文本框大小")

# 全局翻译器实例
page_translator = PageBasedTranslator()

async def translate_slide_by_page(slide, slide_index: int, source_language: str,
                                target_language: str, bilingual_translation: str = "1",
                                field: str = "通用") -> int:
    """按页翻译幻灯片（外部接口）"""
    try:
        # 翻译段落
        matches = await page_translator.translate_slide_paragraphs(
            slide, slide_index, source_language, target_language, field
        )

        # 应用翻译
        applied_count = page_translator.apply_translations_to_slide(
            slide, slide_index, matches, bilingual_translation
        )

        return applied_count

    except Exception as e:
        logger.error(f"按页翻译失败 (第 {slide_index + 1} 页): {str(e)}")
        return 0

def get_translation_statistics() -> Dict[str, Any]:
    """获取翻译统计信息"""
    return {
        "translator_type": "page_based",
        "current_slide_paragraphs": len(page_translator.current_slide_paragraphs),
        "translatable_paragraphs": len([p for p in page_translator.current_slide_paragraphs if p.is_translatable])
    }
