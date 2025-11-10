"""
PPTX整体调度模块
使用python-pptx库替代PyUNO进行PPT编辑
"""
import os,sys
sys.path.insert(0, os.path.dirname(__file__))
from logger_config import get_logger, log_function_call, log_execution_time
import shutil
from pptx import Presentation
from write_ppt_page_pptx import write_page_with_pptx

def get_slide_count(pptx_path):
    """
    获取PPT总页数
    Args:
        pptx_path: PPTX文件路径
    Returns:
        int: 幻灯片总数
    """
    logger = get_logger("pyuno.main")
    try:
        prs = Presentation(pptx_path)
        return len(prs.slides)
    except Exception as e:
        logger.error(f"获取PPT页数失败: {str(e)}")
        return 0

def validate_page_indices_pptx(pptx_path, page_indices):
    """
    验证页面索引的有效性(针对python-pptx)
    Args:
        pptx_path: PPTX文件路径
        page_indices: 页面索引列表(0-based)
    Returns:
        list: 有效的页面索引列表
    """
    logger = get_logger("pyuno.main")
    total_pages = get_slide_count(pptx_path)
    if total_pages == 0:
        return []
    
    valid_indices = []
    for idx in page_indices:
        if 0 <= idx < total_pages:
            valid_indices.append(idx)
        else:
            logger.warning(f"页面索引 {idx} 超出范围(0-{total_pages-1})，已忽略")
    
    return valid_indices

def backup_pptx_for_editing(source_path, target_path):
    """
    为编辑创建PPTX文件副本
    Args:
        source_path: 源文件路径
        target_path: 目标文件路径
    """
    logger = get_logger("pyuno.main")
    try:
        # 如果目标文件已存在，先删除
        if os.path.exists(target_path):
            os.remove(target_path)
        
        shutil.copy2(source_path, target_path)
        logger.info(f"已创建编辑副本: {target_path}")
    except Exception as e:
        logger.error(f"创建编辑副本失败: {str(e)}")
        raise

def edit_ppt_with_pptx(backup_pptx_path, translated_pages_data, bilingual_translation,
                       processed_page_indices, output_path, progress_callback=None):
    """
    使用python-pptx库编辑PPT的主控制函数
    修改版：兼容旧的数据格式
    
    Args:
        backup_pptx_path: 备份的原始PPTX文件路径
        translated_pages_data: 翻译后的数据，可能是完整PPT数据字典或页面数据列表
        bilingual_translation: 双语翻译模式
        processed_page_indices: 需要处理的页面索引列表(0-based)
        output_path: 输出文件路径
        progress_callback: 进度回调函数
    
    Returns:
        str: 输出文件路径
    """
    logger = get_logger("pyuno.main")

    try:
        # 数据格式兼容处理 - 新增
        if isinstance(translated_pages_data, dict):
            # 如果是完整的PPT数据字典
            pages_data = translated_pages_data.get('pages', [])
            logger.info(f"检测到完整PPT数据字典，提取到 {len(pages_data)} 页数据")
        elif isinstance(translated_pages_data, list):
            # 如果直接是页面数据列表
            pages_data = translated_pages_data
            logger.info(f"检测到页面数据列表，共 {len(pages_data)} 页")
        else:
            # 如果是其他格式，记录错误信息
            logger.error(f"不支持的数据格式: {type(translated_pages_data)}")
            logger.debug(f"数据内容: {str(translated_pages_data)[:200]}...")
            raise Exception(f"不支持的翻译数据格式: {type(translated_pages_data)}")
        
        # 创建输出文件副本
        backup_pptx_for_editing(backup_pptx_path, output_path)
        
        # 打开PPTX文件
        # if progress_callback:
        #     progress_callback("正在打开PPTX文件...", 85)
        
        prs = Presentation(output_path)
        total_slides = len(prs.slides)
        
        logger.info(f"PPTX文件已打开，总共 {total_slides} 页")
        
        # 验证页面索引
        if processed_page_indices:
            valid_indices = validate_page_indices_pptx(output_path, processed_page_indices)
        else:
            # 如果没有指定页面索引，处理所有页面
            valid_indices = list(range(min(len(pages_data), total_slides)))
        
        if not valid_indices:
            raise Exception("没有有效的页面索引可处理")
        
        logger.info(f"将处理页面索引: {valid_indices}")
        
        # 逐页处理翻译内容
        for i, page_data in enumerate(pages_data):
            # 确定要处理的页面索引
            if processed_page_indices:
                # 如果指定了页面索引，使用页面数据中的original_page_index
                if isinstance(page_data, dict):
                    original_page_index = page_data.get('original_page_index', page_data.get('page_index', i))
                else:
                    logger.warning(f"页面数据不是字典格式: {type(page_data)}")
                    continue
            else:
                # 如果没有指定页面索引，按顺序处理
                original_page_index = i
            
            # 检查页面索引是否在有效范围内
            if original_page_index >= total_slides:
                logger.warning(f"页面索引 {original_page_index} 超出范围，跳过")
                continue
            
            # 检查是否需要处理这个页面
            if processed_page_indices and original_page_index not in valid_indices:
                logger.debug(f"页面 {original_page_index} 不在处理范围内，跳过")
                continue
            
            slide = prs.slides[original_page_index]
            
            # if progress_callback:
            #     progress = 85 + (i + 1) / len(pages_data) * 10
            #     progress_callback(f"正在处理第 {original_page_index + 1} 页...", progress)
            
            # 调用页面级处理函数
            write_page_with_pptx(slide, page_data, bilingual_translation)
            
            logger.info(f"第 {original_page_index + 1} 页处理完成")
        
        # 保存文件
        # if progress_callback:
        #     progress_callback("正在保存文件...", 95)
        
        prs.save(output_path)
        logger.info(f"PPTX文件已保存到: {output_path}")
        
        return output_path
        
    except Exception as e:
        error_msg = f"使用python-pptx编辑PPT失败: {str(e)}"
        logger.error(error_msg)
        
        # 添加详细的调试信息
        logger.debug(f"backup_pptx_path: {backup_pptx_path}")
        logger.debug(f"translated_pages_data类型: {type(translated_pages_data)}")
        logger.debug(f"processed_page_indices: {processed_page_indices}")
        logger.debug(f"output_path: {output_path}")
        
        raise Exception(error_msg)

def get_pptx_info(pptx_path):
    """
    获取PPTX文件信息
    Args:
        pptx_path: PPTX文件路径
    Returns:
        dict: 文件信息
    """
    try:
        prs = Presentation(pptx_path)
        
        info = {
            'total_slides': len(prs.slides),
            'slide_width': prs.slide_width,
            'slide_height': prs.slide_height,
            'slides_info': []
        }
        
        # 获取每页的文本框信息
        for i, slide in enumerate(prs.slides):
            slide_info = {
                'slide_index': i,
                'textbox_count': sum(1 for shape in slide.shapes if shape.has_text_frame),
                'shape_count': len(slide.shapes)
            }
            info['slides_info'].append(slide_info)
        
        return info
        
    except Exception as e:
        logger.error(f"获取PPTX文件信息失败: {str(e)}")
        return None

def verify_pptx_integrity(pptx_path):
    """
    验证PPTX文件完整性
    Args:
        pptx_path: PPTX文件路径
    Returns:
        bool: 文件是否完整
    """
    logger = get_logger("pyuno.main")
    try:
        prs = Presentation(pptx_path)
        # 尝试访问所有幻灯片
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    # 简单访问文本内容
                    _ = shape.text
        return True
    except Exception as e:
        logger.error(f"PPTX文件完整性验证失败: {str(e)}")
        return False

# 工具函数
def count_textboxes_in_slide(slide):
    """统计幻灯片中的文本框数量"""
    return sum(1 for shape in slide.shapes if shape.has_text_frame)

def get_slide_textboxes(slide):
    """获取幻灯片中的所有文本框"""
    return [shape for shape in slide.shapes if shape.has_text_frame]

def log_slide_structure(slide, slide_index):
    """记录幻灯片结构信息（调试用）"""
    logger = get_logger("pyuno.main")
    textboxes = get_slide_textboxes(slide)
    logger.debug(f"幻灯片 {slide_index + 1}: {len(textboxes)} 个文本框")
    
    for i, textbox in enumerate(textboxes):
        paragraph_count = len(textbox.text_frame.paragraphs)
        char_count = len(textbox.text)
        logger.debug(f"  文本框 {i + 1}: {paragraph_count} 段落, {char_count} 字符")

# 测试函数
def test_edit_pptx(test_file_path):
    """
    测试PPTX编辑功能
    """
    try:
        # 获取文件信息
        info = get_pptx_info(test_file_path)
        print(f"文件信息: {info}")
        
        # 验证文件完整性
        is_valid = verify_pptx_integrity(test_file_path)
        print(f"文件完整性: {is_valid}")
        
        return True
        
    except Exception as e:
        print(f"测试失败: {str(e)}")
        return False

if __name__ == "__main__":
    # 测试代码
    logger = get_logger("pyuno.main")
    logger.basicConfig(level=logger.DEBUG)
    
    test_file = "test.pptx"
    if os.path.exists(test_file):
        test_edit_pptx(test_file)
    else:
        print(f"测试文件 {test_file} 不存在")