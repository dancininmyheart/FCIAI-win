"""
PPTX页面详细编辑模块（改进版 - 保留段落格式）
使用python-pptx库实现页面级的翻译内容写入
"""
import logging
import difflib
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx_format_utils import (
    rgb_to_pptx_color, 
    apply_superscript_subscript,
    set_font_properties,
    calculate_text_similarity
)

def write_page_with_pptx(slide, page_data, bilingual_translation):
    """
    使用python-pptx写入单个页面的翻译内容
    
    Args:
        slide: python-pptx的Slide对象
        page_data: 页面数据(包含翻译结果)
        bilingual_translation: 双语翻译模式
    """
    try:
        logging.info(f"开始处理页面，翻译模式: {bilingual_translation}")
        
        # 获取页面中的所有文本框
        slide_textboxes = get_slide_textboxes(slide)
        page_textboxes_data = page_data.get('text_boxes', [])
        
        logging.info(f"页面有 {len(slide_textboxes)} 个文本框，数据有 {len(page_textboxes_data)} 个文本框")
        
        # 进行文本框匹配
        textbox_mapping = match_textboxes_pptx(slide_textboxes, page_textboxes_data)
        
        # 逐个处理匹配的文本框
        for data_index, pptx_shape_index in textbox_mapping.items():
            if pptx_shape_index is not None and data_index < len(page_textboxes_data):
                textbox_data = page_textboxes_data[data_index]
                pptx_shape = slide_textboxes[pptx_shape_index]
                
                logging.info(f"处理文本框 {data_index} -> PPTX形状 {pptx_shape_index}")
                process_textbox_pptx(pptx_shape, textbox_data, bilingual_translation)
        
        logging.info("页面处理完成")
        
    except Exception as e:
        logging.error(f"写入页面内容失败: {str(e)}")
        raise

def get_slide_textboxes(slide):
    """获取幻灯片中的所有文本框"""
    return [shape for shape in slide.shapes if shape.has_text_frame]

def match_textboxes_pptx(slide_textboxes, page_textboxes_data):
    """
    文本框智能匹配功能
    
    Args:
        slide_textboxes: PPT中的文本框Shape对象列表
        page_textboxes_data: 从PyUNO提取的文本框数据列表
    
    Returns:
        dict: {data_textbox_index: pptx_shape_index} 的映射关系
    """
    mapping = {}
    used_pptx_indices = set()
    
    # 提取PPT文本框的文本内容
    slide_texts = []
    for shape in slide_textboxes:
        try:
            text = shape.text.strip()
            slide_texts.append(text)
        except:
            slide_texts.append("")
    
    # 提取数据中的文本内容
    data_texts = []
    for textbox_data in page_textboxes_data:
        combined_text = ""
        for paragraph in textbox_data.get('paragraphs', []):
            for fragment in paragraph.get('text_fragments', []):
                combined_text += fragment.get('text', '')
        data_texts.append(combined_text.strip())
    
    logging.debug(f"PPT文本框内容: {slide_texts}")
    logging.debug(f"数据文本框内容: {data_texts}")
    
    # 第一步：精确匹配
    for data_idx, data_text in enumerate(data_texts):
        for pptx_idx, slide_text in enumerate(slide_texts):
            if pptx_idx not in used_pptx_indices and data_text == slide_text:
                mapping[data_idx] = pptx_idx
                used_pptx_indices.add(pptx_idx)
                logging.debug(f"精确匹配: 数据 {data_idx} -> PPT {pptx_idx}")
                break
    
    # 第二步：相似度匹配（对未匹配的文本框）
    for data_idx, data_text in enumerate(data_texts):
        if data_idx in mapping:
            continue
        
        best_match_idx = None
        best_similarity = 0.0
        
        for pptx_idx, slide_text in enumerate(slide_texts):
            if pptx_idx in used_pptx_indices:
                continue
            
            similarity = calculate_text_similarity(data_text, slide_text)
            if similarity > best_similarity and similarity >= 0.7:  # 相似度阈值
                best_similarity = similarity
                best_match_idx = pptx_idx
        
        if best_match_idx is not None:
            mapping[data_idx] = best_match_idx
            used_pptx_indices.add(best_match_idx)
            logging.debug(f"相似度匹配: 数据 {data_idx} -> PPT {best_match_idx} (相似度: {best_similarity:.2f})")
    
    # 第三步：按顺序匹配剩余的文本框
    unmatched_data = [i for i in range(len(data_texts)) if i not in mapping]
    unmatched_pptx = [i for i in range(len(slide_texts)) if i not in used_pptx_indices]
    
    for data_idx, pptx_idx in zip(unmatched_data, unmatched_pptx):
        mapping[data_idx] = pptx_idx
        logging.debug(f"顺序匹配: 数据 {data_idx} -> PPT {pptx_idx}")
    
    logging.info(f"文本框匹配结果: {mapping}")
    return mapping

def copy_paragraph_format(source_paragraph, target_paragraph):
    """
    复制段落格式属性
    
    Args:
        source_paragraph: 源段落对象
        target_paragraph: 目标段落对象
    """
    try:
        # 复制段落对齐方式
        if hasattr(source_paragraph, 'alignment') and source_paragraph.alignment is not None:
            target_paragraph.alignment = source_paragraph.alignment
        
        # 复制段落级格式属性
        if hasattr(source_paragraph, 'paragraph_format') and hasattr(target_paragraph, 'paragraph_format'):
            source_format = source_paragraph.paragraph_format
            target_format = target_paragraph.paragraph_format
            
            # 复制缩进设置
            try:
                if source_format.left_indent is not None:
                    target_format.left_indent = source_format.left_indent
                if source_format.right_indent is not None:
                    target_format.right_indent = source_format.right_indent
                if source_format.first_line_indent is not None:
                    target_format.first_line_indent = source_format.first_line_indent
            except Exception as e:
                logging.debug(f"复制缩进设置失败: {e}")
            
            # 复制行间距设置
            try:
                if source_format.line_spacing is not None:
                    target_format.line_spacing = source_format.line_spacing
                if source_format.space_before is not None:
                    target_format.space_before = source_format.space_before
                if source_format.space_after is not None:
                    target_format.space_after = source_format.space_after
            except Exception as e:
                logging.debug(f"复制行间距设置失败: {e}")
        
        # 复制项目符号/编号设置
        try:
            if hasattr(source_paragraph, 'level') and hasattr(target_paragraph, 'level'):
                target_paragraph.level = source_paragraph.level
        except Exception as e:
            logging.debug(f"复制段落级别失败: {e}")
        
        logging.debug("段落格式复制完成")
        
    except Exception as e:
        logging.warning(f"复制段落格式失败: {e}")

def process_textbox_pptx(pptx_shape, textbox_data, bilingual_translation):
    """
    处理单个文本框的翻译写入（改进版 - 保留段落格式 + 自适应修复）
    
    Args:
        pptx_shape: python-pptx的Shape对象
        textbox_data: 文本框数据
        bilingual_translation: 双语翻译模式
    """
    try:
        text_frame = pptx_shape.text_frame
        paragraphs_data = textbox_data.get('paragraphs', [])
        
        # 获取现有段落信息
        existing_paragraphs = text_frame.paragraphs
        existing_count = len(existing_paragraphs)
        needed_count = len(paragraphs_data)
        
        logging.debug(f"文本框段落处理: 现有 {existing_count} 段，需要 {needed_count} 段")
        
        # 1. 处理现有段落（保留格式，只清空内容）
        for i in range(min(existing_count, needed_count)):
            paragraph = existing_paragraphs[i]
            
            # 记录原段落格式信息（调试用）
            original_alignment = getattr(paragraph, 'alignment', None)
            original_level = getattr(paragraph, 'level', None)
            
            # 只清空内容，保留段落格式
            paragraph.clear()
            
            # 重构段落内容
            rebuild_paragraph_pptx(paragraph, paragraphs_data[i], bilingual_translation)
            
            logging.debug(f"段落 {i} 处理完成（保留格式: 对齐={original_alignment}, 级别={original_level}）")
        
        # 2. 如果需要更多段落，添加新段落并复制格式
        if needed_count > existing_count:
            # 选择格式模板（使用最后一个现有段落的格式）
            template_paragraph = existing_paragraphs[-1] if existing_count > 0 else None
            
            for i in range(existing_count, needed_count):
                # 创建新段落
                new_paragraph = text_frame.add_paragraph()
                
                # 如果有模板段落，复制其格式
                if template_paragraph:
                    copy_paragraph_format(template_paragraph, new_paragraph)
                    logging.debug(f"为新段落 {i} 复制了模板格式")
                
                # 重构段落内容
                rebuild_paragraph_pptx(new_paragraph, paragraphs_data[i], bilingual_translation)
                
                logging.debug(f"新段落 {i} 创建并处理完成")
        
        # 3. 如果段落太多，清空多余的段落内容但保留段落本身
        elif needed_count < existing_count:
            for i in range(needed_count, existing_count):
                existing_paragraphs[i].clear()
                logging.debug(f"清空多余段落 {i} 的内容")
        
        # 4. 应用文本自适应并修复（新增）
        logging.debug("开始应用文本自适应修复...") #无效
        # apply_text_autofit_fix(pptx_shape)
        # text_frame.fit_text()
        
        logging.debug(f"文本框处理完成，共处理 {needed_count} 个段落")
        
    except Exception as e:
        logging.error(f"处理文本框失败: {str(e)}")
        raise

def apply_text_autofit_fix(shape):
    """
    应用文本自适应修复
    没啥用
    
    Args:
        shape: python-pptx的Shape对象
    """
    try:
        if not shape.has_text_frame:
            return
        
        text_frame = shape.text_frame
        from pptx.enum.text import MSO_AUTO_SIZE
        
        # 设置自适应
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        # 方法1: 通过微调尺寸触发重新计算
        try:
            original_width = text_frame.width
            original_height = text_frame.height
            
            # 微调尺寸（很小的变化）
            text_frame.width = original_width + 1
            text_frame.height = original_height + 1
            
            # # 恢复原始尺寸
            text_frame.width = original_width
            text_frame.height = original_height
            
            logging.debug("通过尺寸微调触发文本自适应")
            
        except Exception as e:
            logging.debug(f"尺寸微调方法失败: {e}")
            
            # 方法2: 通过margin微调触发
            try:
                original_margin = text_frame.margin_left
                text_frame.margin_left = original_margin + 100
                text_frame.margin_left = original_margin
                
                logging.debug("通过margin微调触发文本自适应")
                
            except Exception as e2:
                logging.debug(f"margin微调方法也失败: {e2}")
                
                # 方法3: 强制刷新文本内容
                try:
                    if text_frame.paragraphs and text_frame.paragraphs[0].runs:
                        first_run = text_frame.paragraphs[0].runs[0]
                        current_text = first_run.text
                        first_run.text = current_text + " "
                        first_run.text = current_text
                        
                        logging.debug("通过文本刷新触发文本自适应")
                        
                except Exception as e3:
                    logging.warning(f"所有自适应触发方法都失败: {e3}")
    
    except Exception as e:
        logging.error(f"应用文本自适应修复失败: {e}")

def rebuild_paragraph_pptx(paragraph, paragraph_data, bilingual_translation):
    """
    重构单个段落内容
    
    Args:
        paragraph: python-pptx的Paragraph对象
        paragraph_data: 段落数据
        bilingual_translation: 双语翻译模式
    """
    try:
        # 注意：这里不再调用 paragraph.clear()，因为已经在上层处理了
        
        original_fragments = paragraph_data.get('text_fragments', [])
        
        # 根据翻译模式处理内容
        handle_bilingual_modes_pptx(paragraph, original_fragments, bilingual_translation)
        
        logging.debug(f"段落重构完成，共 {len(original_fragments)} 个片段")
        
    except Exception as e:
        logging.error(f"重构段落失败: {str(e)}")
        raise

def handle_bilingual_modes_pptx(paragraph, fragments, bilingual_translation):
    """
    处理各种双语翻译模式（支持软换行）
    
    Args:
        paragraph: python-pptx的Paragraph对象
        fragments: 文本片段列表
        bilingual_translation: 翻译模式
    """
    try:
        if bilingual_translation == "translation_only":
            # 完全替换：只显示译文
            for fragment in fragments:
                translated_text = fragment.get('translated_text', '')
                if translated_text:
                    run = paragraph.add_run()
                    run.text = translated_text
                    apply_fragment_format_pptx(run, fragment)
                    
        elif bilingual_translation == "paragraph_up":
            # 原文在上，译文在下（使用软换行）
            # 首先添加所有原文片段
            for fragment in fragments:
                original_text = fragment.get('text', '')
                if original_text:
                    run = paragraph.add_run()
                    run.text = original_text
                    apply_fragment_format_pptx(run, fragment)
            
            # 添加软换行
            if any(f.get('translated_text', '') for f in fragments):
                insert_soft_line_break(paragraph)
            
            # 然后添加所有译文片段
            for fragment in fragments:
                translated_text = fragment.get('translated_text', '')
                if translated_text:
                    run = paragraph.add_run()
                    run.text = translated_text
                    apply_fragment_format_pptx(run, fragment)
                    
        elif bilingual_translation == "paragraph_down":
            # 译文在上，原文在下（使用软换行）
            # 首先添加所有译文片段
            for fragment in fragments:
                translated_text = fragment.get('translated_text', '')
                if translated_text:
                    run = paragraph.add_run()
                    run.text = translated_text
                    apply_fragment_format_pptx(run, fragment)
            
            # 添加软换行
            if any(f.get('text', '') for f in fragments):
                insert_soft_line_break(paragraph)
            
            # 然后添加所有原文片段
            for fragment in fragments:
                original_text = fragment.get('text', '')
                if original_text:
                    run = paragraph.add_run()
                    run.text = original_text
                    apply_fragment_format_pptx(run, fragment)
        
        else:
            # 默认：只显示原文
            for fragment in fragments:
                original_text = fragment.get('text', '')
                if original_text:
                    run = paragraph.add_run()
                    run.text = original_text
                    apply_fragment_format_pptx(run, fragment)
                    
        logging.debug(f"双语模式 {bilingual_translation} 处理完成")
        
    except Exception as e:
        logging.error(f"处理双语模式失败: {str(e)}")
        raise

def insert_soft_line_break(paragraph):
    """
    在段落中插入软换行
    
    软换行保持在同一个段落内，不会产生新的项目符号
    
    Args:
        paragraph: python-pptx的Paragraph对象
    """
    try:
        # 方法1: 使用Unicode Line Separator (U+2028)
        # 这是标准的软换行字符，在大多数情况下都能正确工作
        paragraph.add_line_break()
        logging.debug("使用Unicode Line Separator插入软换行")
        return True
        
    except Exception as e:
        logging.warning(f"插入软换行失败，使用方法1: {e}")
        
        try:
            # 方法2: 使用python-pptx的内置换行方法（如果存在）
            line_break_run = paragraph.add_run()
            # 尝试使用特殊的换行字符组合
            line_break_run.text = "\r\n"  # Windows风格换行
            
            logging.debug("使用CR+LF插入软换行")
            return True
            
        except Exception as e2:
            logging.warning(f"插入软换行失败，使用方法2: {e2}")
            
            try:
                # 方法3: 回退到普通换行符
                line_break_run = paragraph.add_run()
                line_break_run.text = "\n"
                
                logging.debug("使用普通换行符作为软换行")
                return True
                
            except Exception as e3:
                logging.error(f"所有软换行方法都失败: {e3}")
                return False


def apply_fragment_format_pptx(run, fragment_data):
    """
    应用文本片段格式
    
    Args:
        run: python-pptx的Run对象
        fragment_data: 片段格式数据
    """
    try:
        from pptx_format_utils import validate_format_data
        
        # 验证和标准化格式数据
        validated_data = validate_format_data(fragment_data)
        
        # 应用字体属性
        set_font_properties(
            run,
            validated_data['font_size'],
            validated_data['color'],
            validated_data['bold'],
            validated_data['italic'],
            validated_data['underline']
        )
        
        # 应用上下标
        apply_superscript_subscript(
            run,
            validated_data.get('escapement', 0),
        )
        
        logging.debug(f"片段格式应用完成")
        
    except Exception as e:
        logging.error(f"应用片段格式失败: {str(e)}")
        # 即使格式应用失败，也不应该影响文本写入
        pass

# 工具函数
def count_textboxes_in_slide(slide):
    """统计幻灯片中的文本框数量"""
    return sum(1 for shape in slide.shapes if shape.has_text_frame)

def log_slide_structure(slide, slide_index):
    """记录幻灯片结构信息（调试用）"""
    textboxes = get_slide_textboxes(slide)
    logging.debug(f"幻灯片 {slide_index + 1}: {len(textboxes)} 个文本框")
    
    for i, textbox in enumerate(textboxes):
        paragraph_count = len(textbox.text_frame.paragraphs)
        char_count = len(textbox.text)
        logging.debug(f"  文本框 {i + 1}: {paragraph_count} 段落, {char_count} 字符")

def log_paragraph_format_info(paragraph, index):
    """记录段落格式信息（调试用）"""
    try:
        alignment = getattr(paragraph, 'alignment', None)
        level = getattr(paragraph, 'level', None)
        
        format_info = []
        if hasattr(paragraph, 'paragraph_format'):
            pf = paragraph.paragraph_format
            if pf.left_indent is not None:
                format_info.append(f"左缩进={pf.left_indent}")
            if pf.first_line_indent is not None:
                format_info.append(f"首行缩进={pf.first_line_indent}")
        
        logging.debug(f"段落 {index}: 对齐={alignment}, 级别={level}, {', '.join(format_info)}")
        
    except Exception as e:
        logging.debug(f"记录段落 {index} 格式信息失败: {e}")

# 测试函数
def test_write_page():
    """测试页面写入功能"""
    try:
        logging.basicConfig(level=logging.DEBUG)
        
        # 这里可以添加单元测试代码
        print("write_ppt_page_pptx 模块（改进版）加载成功")
        return True
        
    except Exception as e:
        print(f"测试失败: {str(e)}")
        return False

if __name__ == "__main__":
    test_write_page()