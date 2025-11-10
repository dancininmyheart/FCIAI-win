#!/usr/bin/env python3
"""
确保PPT中所有文本框都设置为自动调整大小的工具函数
解决文本框未全部设置为MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE的问题
"""
import os
import logging
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt

# 配置日志记录器
logger = logging.getLogger(__name__)

def ensure_all_textboxes_autofit(presentation_path: str, verbose: bool = True) -> bool:
    """
    确保PPT中所有文本框都设置为自动调整大小

    Args:
        presentation_path: PPT文件路径
        verbose: 是否输出详细日志

    Returns:
        调整是否成功
    """
    if not os.path.exists(presentation_path):
        logger.error(f"文件不存在: {presentation_path}")
        return False

    try:
        # 加载演示文稿
        prs = Presentation(presentation_path)

        total_shapes = 0
        total_textboxes = 0
        processed_textboxes = 0
        skipped_shapes = 0

        if verbose:
            logger.info(f"开始处理PPT文件: {os.path.basename(presentation_path)}")
            logger.info("确保所有文本框都设置为自动调整大小...")

        # 遍历所有幻灯片
        for slide_index, slide in enumerate(prs.slides, 1):
            if verbose:
                logger.debug(f"检查第 {slide_index} 张幻灯片的所有形状...")

            for shape_index, shape in enumerate(slide.shapes):
                total_shapes += 1

                try:
                    # 处理普通文本框
                    if shape.has_text_frame:
                        total_textboxes += 1
                        text_frame = shape.text_frame

                        # 只设置自动调整，不改变其他格式
                        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                        processed_textboxes += 1
                        if verbose:
                            logger.debug(f"✓ 幻灯片{slide_index}-形状{shape_index+1}: 已设置文本框自动调整")

                    # 处理表格
                    elif shape.has_table:
                        table = shape.table
                        if verbose:
                            logger.debug(f"处理表格: {table.rows} 行 x {table.columns} 列")

                        for row_index, row in enumerate(table.rows):
                            for col_index, cell in enumerate(row.cells):
                                total_textboxes += 1

                                # 表格单元格的文本框 - 调整文本大小以适应单元格（保持单元格大小不变）
                                text_frame = cell.text_frame
                                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                                processed_textboxes += 1
                                if verbose:
                                    logger.debug(f"✓ 幻灯片{slide_index}-表格单元格({row_index+1},{col_index+1}): 已设置自动调整")

                    else:
                        skipped_shapes += 1
                        if verbose:
                            logger.debug(f"跳过非文本形状: 幻灯片{slide_index}-形状{shape_index+1} (类型: {shape.shape_type})")

                except Exception as shape_error:
                    logger.warning(f"处理幻灯片{slide_index}-形状{shape_index+1}时出错: {shape_error}")

        # 保存演示文稿
        prs.save(presentation_path)

        if verbose:
            logger.info(f"文本框自动调整设置完成:")
            logger.info(f"  - 总形状数: {total_shapes}")
            logger.info(f"  - 文本框总数: {total_textboxes}")
            logger.info(f"  - 已处理文本框: {processed_textboxes}")
            logger.info(f"  - 跳过的形状: {skipped_shapes}")
            logger.info(f"  - 成功率: {(processed_textboxes/total_textboxes*100):.1f}%" if total_textboxes > 0 else "  - 成功率: N/A")

        return True

    except Exception as e:
        logger.error(f"确保文本框自动调整失败: {e}")
        import traceback
        logger.error(f"错误详情: {traceback.format_exc()}")
        return False

def check_textbox_autofit_status(presentation_path: str) -> dict:
    """
    检查PPT中文本框的自动调整状态

    Args:
        presentation_path: PPT文件路径

    Returns:
        包含检查结果的字典
    """
    if not os.path.exists(presentation_path):
        return {"error": f"文件不存在: {presentation_path}"}

    try:
        # 加载演示文稿
        prs = Presentation(presentation_path)

        total_textboxes = 0
        autofit_textboxes = 0
        non_autofit_textboxes = 0
        textbox_details = []

        # 遍历所有幻灯片
        for slide_index, slide in enumerate(prs.slides, 1):
            for shape_index, shape in enumerate(slide.shapes):
                # 处理普通文本框
                if shape.has_text_frame:
                    total_textboxes += 1
                    text_frame = shape.text_frame

                    is_autofit = text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                    if is_autofit:
                        autofit_textboxes += 1
                    else:
                        non_autofit_textboxes += 1

                    textbox_details.append({
                        "slide": slide_index,
                        "shape": shape_index + 1,
                        "type": "textbox",
                        "auto_size": text_frame.auto_size,
                        "is_autofit": is_autofit,
                        "word_wrap": text_frame.word_wrap
                    })

                # 处理表格
                elif shape.has_table:
                    table = shape.table
                    for row_index, row in enumerate(table.rows):
                        for col_index, cell in enumerate(row.cells):
                            total_textboxes += 1
                            text_frame = cell.text_frame

                            is_autofit = text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                            if is_autofit:
                                autofit_textboxes += 1
                            else:
                                non_autofit_textboxes += 1

                            textbox_details.append({
                                "slide": slide_index,
                                "shape": f"table-{row_index+1}-{col_index+1}",
                                "type": "table_cell",
                                "auto_size": text_frame.auto_size,
                                "is_autofit": is_autofit,
                                "word_wrap": text_frame.word_wrap
                            })

        return {
            "file": os.path.basename(presentation_path),
            "total_textboxes": total_textboxes,
            "autofit_textboxes": autofit_textboxes,
            "non_autofit_textboxes": non_autofit_textboxes,
            "autofit_percentage": (autofit_textboxes / total_textboxes * 100) if total_textboxes > 0 else 0,
            "details": textbox_details
        }

    except Exception as e:
        return {"error": f"检查失败: {str(e)}"}

def fix_textbox_autofit_issues(presentation_path: str) -> bool:
    """
    修复PPT中文本框自动调整的问题

    Args:
        presentation_path: PPT文件路径

    Returns:
        修复是否成功
    """
    logger.info(f"开始修复PPT文本框自动调整问题: {os.path.basename(presentation_path)}")

    # 首先检查当前状态
    status = check_textbox_autofit_status(presentation_path)
    if "error" in status:
        logger.error(f"检查失败: {status['error']}")
        return False

    logger.info(f"检查结果: {status['autofit_textboxes']}/{status['total_textboxes']} 个文本框已设置自动调整 ({status['autofit_percentage']:.1f}%)")

    if status['non_autofit_textboxes'] == 0:
        logger.info("所有文本框都已正确设置，无需修复")
        return True

    # 执行修复
    logger.info(f"需要修复 {status['non_autofit_textboxes']} 个文本框...")
    result = ensure_all_textboxes_autofit(presentation_path, verbose=True)

    if result:
        # 再次检查修复结果
        new_status = check_textbox_autofit_status(presentation_path)
        if "error" not in new_status:
            logger.info(f"修复完成: {new_status['autofit_textboxes']}/{new_status['total_textboxes']} 个文本框已设置自动调整 ({new_status['autofit_percentage']:.1f}%)")
            return new_status['non_autofit_textboxes'] == 0

    return result

if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("用法: python ensure_textbox_autofit.py <ppt_file_path> [action]")
        print("action: check (检查状态) | fix (修复问题) | ensure (确保设置)")
        sys.exit(1)

    ppt_path = sys.argv[1]
    action = sys.argv[2] if len(sys.argv) > 2 else "fix"

    # 配置日志
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(levelname)s - %(message)s'
    )

    if action == "check":
        status = check_textbox_autofit_status(ppt_path)
        if "error" in status:
            print(f"错误: {status['error']}")
        else:
            print(f"文件: {status['file']}")
            print(f"总文本框数: {status['total_textboxes']}")
            print(f"已设置自动调整: {status['autofit_textboxes']}")
            print(f"未设置自动调整: {status['non_autofit_textboxes']}")
            print(f"自动调整比例: {status['autofit_percentage']:.1f}%")

    elif action == "fix":
        success = fix_textbox_autofit_issues(ppt_path)
        if success:
            print("✓ 修复成功")
        else:
            print("✗ 修复失败")

    elif action == "ensure":
        success = ensure_all_textboxes_autofit(ppt_path)
        if success:
            print("✓ 确保设置成功")
        else:
            print("✗ 确保设置失败")

    else:
        print(f"未知操作: {action}")
        sys.exit(1)
