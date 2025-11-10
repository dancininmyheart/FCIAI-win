#!/usr/bin/env python3
"""
保持文本框大小不变的文本自适应功能
确保文本框设置为TEXT_TO_FIT_SHAPE时不会改变文本框的原始尺寸
"""
import os
import logging
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt

# 配置日志记录器
logger = logging.getLogger(__name__)

def preserve_textbox_size_with_autofit(presentation_path: str, verbose: bool = True) -> bool:
    """
    设置文本框自适应的同时保持文本框大小不变

    Args:
        presentation_path: PPT文件路径
        verbose: 是否输出详细日志

    Returns:
        调整是否成功
    """
    if not os.path.exists(presentation_path):
        logger.error(f"文件不存在: {presentation_path}")
        return False

    try:
        # 加载演示文稿
        prs = Presentation(presentation_path)

        total_shapes = 0
        total_textboxes = 0
        processed_textboxes = 0
        size_preserved_count = 0

        if verbose:
            logger.info(f"开始处理PPT文件: {os.path.basename(presentation_path)}")
            logger.info("设置文本框自适应并保持原始大小...")

        # 遍历所有幻灯片
        for slide_index, slide in enumerate(prs.slides, 1):
            if verbose:
                logger.debug(f"处理第 {slide_index} 张幻灯片...")

            for shape_index, shape in enumerate(slide.shapes):
                total_shapes += 1

                try:
                    # 处理普通文本框
                    if shape.has_text_frame:
                        total_textboxes += 1

                        # 记录原始尺寸
                        original_width = shape.width
                        original_height = shape.height
                        original_left = shape.left
                        original_top = shape.top

                        if verbose:
                            logger.debug(f"原始尺寸 - 宽度: {original_width}, 高度: {original_height}")

                        # 只设置文本框自适应，不改变其他格式
                        text_frame = shape.text_frame
                        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                        # 检查尺寸是否被改变
                        if (shape.width != original_width or
                            shape.height != original_height or
                            shape.left != original_left or
                            shape.top != original_top):

                            # 恢复原始尺寸和位置
                            shape.width = original_width
                            shape.height = original_height
                            shape.left = original_left
                            shape.top = original_top

                            size_preserved_count += 1
                            if verbose:
                                logger.debug(f"✓ 已恢复文本框原始尺寸: 幻灯片{slide_index}-形状{shape_index+1}")

                        processed_textboxes += 1
                        if verbose:
                            logger.debug(f"✓ 幻灯片{slide_index}-形状{shape_index+1}: 已设置文本框自适应")

                    # 处理表格
                    elif shape.has_table:
                        table = shape.table
                        if verbose:
                            logger.debug(f"处理表格: {table.rows} 行 x {table.columns} 列")

                        # 记录表格原始尺寸
                        table_original_width = shape.width
                        table_original_height = shape.height
                        table_original_left = shape.left
                        table_original_top = shape.top

                        for row_index, row in enumerate(table.rows):
                            for col_index, cell in enumerate(row.cells):
                                total_textboxes += 1

                                # 只设置表格单元格文本框自适应，不改变其他格式
                                text_frame = cell.text_frame
                                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                                processed_textboxes += 1
                                if verbose:
                                    logger.debug(f"✓ 幻灯片{slide_index}-表格单元格({row_index+1},{col_index+1}): 已设置自适应")

                        # 确保表格整体尺寸不变
                        if (shape.width != table_original_width or
                            shape.height != table_original_height or
                            shape.left != table_original_left or
                            shape.top != table_original_top):

                            shape.width = table_original_width
                            shape.height = table_original_height
                            shape.left = table_original_left
                            shape.top = table_original_top

                            size_preserved_count += 1
                            if verbose:
                                logger.debug(f"✓ 已恢复表格原始尺寸: 幻灯片{slide_index}-表格")

                except Exception as shape_error:
                    logger.warning(f"处理幻灯片{slide_index}-形状{shape_index+1}时出错: {shape_error}")

        # 保存演示文稿
        prs.save(presentation_path)

        if verbose:
            logger.info(f"文本框自适应设置完成（保持原始大小）:")
            logger.info(f"  - 总形状数: {total_shapes}")
            logger.info(f"  - 文本框总数: {total_textboxes}")
            logger.info(f"  - 已处理文本框: {processed_textboxes}")
            logger.info(f"  - 尺寸保护次数: {size_preserved_count}")
            logger.info(f"  - 成功率: {(processed_textboxes/total_textboxes*100):.1f}%" if total_textboxes > 0 else "  - 成功率: N/A")

        return True

    except Exception as e:
        logger.error(f"设置文本框自适应失败: {e}")
        import traceback
        logger.error(f"错误详情: {traceback.format_exc()}")
        return False

def check_textbox_size_changes(presentation_path: str) -> dict:
    """
    检查PPT处理前后文本框大小的变化

    Args:
        presentation_path: PPT文件路径

    Returns:
        包含大小变化信息的字典
    """
    if not os.path.exists(presentation_path):
        return {"error": f"文件不存在: {presentation_path}"}

    try:
        # 加载演示文稿
        prs = Presentation(presentation_path)

        textbox_info = []

        # 遍历所有幻灯片
        for slide_index, slide in enumerate(prs.slides, 1):
            for shape_index, shape in enumerate(slide.shapes):
                # 处理普通文本框
                if shape.has_text_frame:
                    textbox_info.append({
                        "slide": slide_index,
                        "shape": shape_index + 1,
                        "type": "textbox",
                        "width": shape.width,
                        "height": shape.height,
                        "left": shape.left,
                        "top": shape.top,
                        "auto_size": shape.text_frame.auto_size,
                        "word_wrap": shape.text_frame.word_wrap
                    })

                # 处理表格
                elif shape.has_table:
                    table = shape.table
                    textbox_info.append({
                        "slide": slide_index,
                        "shape": f"table-{shape_index+1}",
                        "type": "table",
                        "width": shape.width,
                        "height": shape.height,
                        "left": shape.left,
                        "top": shape.top,
                        "rows": table.rows,
                        "columns": table.columns
                    })

                    for row_index, row in enumerate(table.rows):
                        for col_index, cell in enumerate(row.cells):
                            textbox_info.append({
                                "slide": slide_index,
                                "shape": f"table-{shape_index+1}-cell-{row_index+1}-{col_index+1}",
                                "type": "table_cell",
                                "auto_size": cell.text_frame.auto_size,
                                "word_wrap": cell.text_frame.word_wrap
                            })

        return {
            "file": os.path.basename(presentation_path),
            "total_textboxes": len([info for info in textbox_info if info["type"] in ["textbox", "table_cell"]]),
            "textbox_info": textbox_info
        }

    except Exception as e:
        return {"error": f"检查失败: {str(e)}"}

def compare_textbox_sizes(before_info: dict, after_info: dict) -> dict:
    """
    比较处理前后文本框大小的变化

    Args:
        before_info: 处理前的文本框信息
        after_info: 处理后的文本框信息

    Returns:
        包含变化统计的字典
    """
    if "error" in before_info or "error" in after_info:
        return {"error": "输入数据包含错误"}

    changes = []
    unchanged_count = 0
    changed_count = 0

    # 创建查找字典
    before_dict = {f"{info['slide']}-{info['shape']}": info for info in before_info["textbox_info"]}
    after_dict = {f"{info['slide']}-{info['shape']}": info for info in after_info["textbox_info"]}

    # 比较每个文本框
    for key, before in before_dict.items():
        if key in after_dict:
            after = after_dict[key]

            if before["type"] in ["textbox", "table"]:
                # 检查尺寸变化
                width_changed = before.get("width") != after.get("width")
                height_changed = before.get("height") != after.get("height")
                left_changed = before.get("left") != after.get("left")
                top_changed = before.get("top") != after.get("top")

                if width_changed or height_changed or left_changed or top_changed:
                    changes.append({
                        "slide": before["slide"],
                        "shape": before["shape"],
                        "type": before["type"],
                        "width_change": after.get("width", 0) - before.get("width", 0),
                        "height_change": after.get("height", 0) - before.get("height", 0),
                        "left_change": after.get("left", 0) - before.get("left", 0),
                        "top_change": after.get("top", 0) - before.get("top", 0)
                    })
                    changed_count += 1
                else:
                    unchanged_count += 1
            else:
                unchanged_count += 1

    return {
        "total_compared": len(before_dict),
        "unchanged_count": unchanged_count,
        "changed_count": changed_count,
        "changes": changes,
        "size_preservation_rate": (unchanged_count / len(before_dict) * 100) if len(before_dict) > 0 else 0
    }

if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("用法: python preserve_textbox_size.py <ppt_file_path> [action]")
        print("action: preserve (保持大小) | check (检查信息) | compare (比较变化)")
        sys.exit(1)

    ppt_path = sys.argv[1]
    action = sys.argv[2] if len(sys.argv) > 2 else "preserve"

    # 配置日志
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(levelname)s - %(message)s'
    )

    if action == "preserve":
        success = preserve_textbox_size_with_autofit(ppt_path)
        if success:
            print("✓ 文本框自适应设置成功（保持原始大小）")
        else:
            print("✗ 文本框自适应设置失败")

    elif action == "check":
        info = check_textbox_size_changes(ppt_path)
        if "error" in info:
            print(f"错误: {info['error']}")
        else:
            print(f"文件: {info['file']}")
            print(f"总文本框数: {info['total_textboxes']}")
            print("文本框信息:")
            for textbox in info['textbox_info']:
                if textbox['type'] in ['textbox', 'table']:
                    print(f"  幻灯片{textbox['slide']}-{textbox['shape']}: {textbox['width']}x{textbox['height']} @ ({textbox['left']},{textbox['top']})")

    else:
        print(f"未知操作: {action}")
        sys.exit(1)
