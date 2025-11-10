"""
PPT颜色备份和恢复模块
在LibreOffice渲染前备份颜色，渲染后恢复颜色
解决LibreOffice渲染导致的颜色变化问题
"""
import os
import json
import logging
from typing import Dict, List, Any, Optional, Tuple
import tempfile

logger = logging.getLogger(__name__)


class PPTColorBackupRestore:
    """PPT颜色备份和恢复器"""
    
    def __init__(self):
        self.color_backup = {}
    
    def backup_colors_from_ppt(self, ppt_path: str) -> Dict[str, Any]:
        """
        从PPT文件备份所有颜色信息
        
        Args:
            ppt_path: PPT文件路径
            
        Returns:
            Dict: 颜色备份信息
        """
        try:
            from pptx import Presentation
            
            prs = Presentation(ppt_path)
            backup_data = {
                "file_path": ppt_path,
                "slides": []
            }
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_data = {
                    "slide_index": slide_idx,
                    "shapes": []
                }
                
                for shape_idx, shape in enumerate(slide.shapes):
                    shape_data = {
                        "shape_index": shape_idx,
                        "background_fill": None,
                        "text_runs": []
                    }
                    
                    # 备份形状背景颜色
                    try:
                        if hasattr(shape, 'fill') and shape.fill.type is not None:
                            fill_info = self._extract_fill_info(shape.fill)
                            if fill_info:
                                shape_data["background_fill"] = fill_info
                    except Exception as e:
                        logger.debug(f"备份形状背景失败: {e}")
                    
                    # 备份文本颜色
                    if shape.has_text_frame:
                        for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                            for run_idx, run in enumerate(paragraph.runs):
                                if run.text.strip():
                                    run_data = self._extract_run_formatting(run, para_idx, run_idx)
                                    if run_data:
                                        shape_data["text_runs"].append(run_data)
                    
                    # 备份表格颜色
                    elif shape.has_table:
                        table = shape.table
                        for row_idx, row in enumerate(table.rows):
                            for col_idx, cell in enumerate(row.cells):
                                # 备份单元格背景
                                try:
                                    if hasattr(cell, 'fill') and cell.fill.type is not None:
                                        fill_info = self._extract_fill_info(cell.fill)
                                        if fill_info:
                                            cell_data = {
                                                "type": "table_cell",
                                                "row": row_idx,
                                                "col": col_idx,
                                                "background_fill": fill_info
                                            }
                                            shape_data["text_runs"].append(cell_data)
                                except:
                                    pass
                                
                                # 备份单元格文本
                                for para_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                                    for run_idx, run in enumerate(paragraph.runs):
                                        if run.text.strip():
                                            run_data = self._extract_run_formatting(
                                                run, para_idx, run_idx, 
                                                table_pos=(row_idx, col_idx)
                                            )
                                            if run_data:
                                                shape_data["text_runs"].append(run_data)
                    
                    if shape_data["background_fill"] or shape_data["text_runs"]:
                        slide_data["shapes"].append(shape_data)
                
                backup_data["slides"].append(slide_data)
            
            self.color_backup = backup_data
            logger.info(f"颜色备份完成: {len(backup_data['slides'])} 张幻灯片")
            return backup_data
            
        except Exception as e:
            logger.error(f"颜色备份失败: {e}")
            return {}
    
    def _extract_fill_info(self, fill) -> Optional[Dict[str, Any]]:
        """提取填充信息"""
        try:
            fill_info = {"type": str(fill.type)}
            
            if hasattr(fill, 'fore_color') and fill.fore_color:
                color_info = self._extract_color_info(fill.fore_color)
                if color_info:
                    fill_info["fore_color"] = color_info
            
            if hasattr(fill, 'back_color') and fill.back_color:
                color_info = self._extract_color_info(fill.back_color)
                if color_info:
                    fill_info["back_color"] = color_info
            
            return fill_info if len(fill_info) > 1 else None
            
        except Exception as e:
            logger.debug(f"提取填充信息失败: {e}")
            return None
    
    def _extract_color_info(self, color) -> Optional[Dict[str, Any]]:
        """提取颜色信息"""
        try:
            color_info = {}
            
            if hasattr(color, 'type'):
                color_info["type"] = color.type
                
                if color.type == 1:  # RGB颜色
                    if hasattr(color, 'rgb') and color.rgb:
                        rgb = color.rgb
                        color_info["rgb"] = [rgb.r, rgb.g, rgb.b]
                
                elif color.type == 2:  # 主题颜色
                    if hasattr(color, 'theme_color'):
                        color_info["theme_color"] = str(color.theme_color)
                    if hasattr(color, 'brightness'):
                        color_info["brightness"] = color.brightness
            
            return color_info if color_info else None
            
        except Exception as e:
            logger.debug(f"提取颜色信息失败: {e}")
            return None
    
    def _extract_run_formatting(self, run, para_idx: int, run_idx: int, 
                               table_pos: Optional[Tuple[int, int]] = None) -> Optional[Dict[str, Any]]:
        """提取文本运行格式"""
        try:
            run_data = {
                "paragraph_index": para_idx,
                "run_index": run_idx,
                "text": run.text,
                "font": {}
            }
            
            if table_pos:
                run_data["table_position"] = table_pos
            
            # 提取字体信息
            font = run.font
            
            # 字体颜色
            if hasattr(font, 'color') and font.color:
                color_info = self._extract_color_info(font.color)
                if color_info:
                    run_data["font"]["color"] = color_info
            
            # 字体大小
            if hasattr(font, 'size') and font.size:
                run_data["font"]["size"] = font.size.pt
            
            # 字体名称
            if hasattr(font, 'name') and font.name:
                run_data["font"]["name"] = font.name
            
            # 字体样式
            if hasattr(font, 'bold') and font.bold is not None:
                run_data["font"]["bold"] = font.bold
            
            if hasattr(font, 'italic') and font.italic is not None:
                run_data["font"]["italic"] = font.italic
            
            if hasattr(font, 'underline') and font.underline is not None:
                run_data["font"]["underline"] = str(font.underline)
            
            return run_data if run_data["font"] else None
            
        except Exception as e:
            logger.debug(f"提取文本格式失败: {e}")
            return None
    
    def restore_colors_to_ppt(self, ppt_path: str, backup_data: Optional[Dict[str, Any]] = None) -> bool:
        """
        将备份的颜色恢复到PPT文件
        
        Args:
            ppt_path: PPT文件路径
            backup_data: 颜色备份数据（可选，默认使用内部备份）
            
        Returns:
            bool: 是否成功
        """
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.util import Pt
            
            if backup_data is None:
                backup_data = self.color_backup
            
            if not backup_data:
                logger.warning("没有颜色备份数据")
                return False
            
            prs = Presentation(ppt_path)
            restored_count = 0
            
            for slide_data in backup_data.get("slides", []):
                slide_idx = slide_data["slide_index"]
                
                if slide_idx < len(prs.slides):
                    slide = prs.slides[slide_idx]
                    
                    for shape_data in slide_data.get("shapes", []):
                        shape_idx = shape_data["shape_index"]
                        
                        if shape_idx < len(slide.shapes):
                            shape = slide.shapes[shape_idx]
                            
                            # 恢复形状背景
                            bg_fill = shape_data.get("background_fill")
                            if bg_fill and hasattr(shape, 'fill'):
                                if self._restore_fill(shape.fill, bg_fill):
                                    restored_count += 1
                            
                            # 恢复文本格式
                            for run_data in shape_data.get("text_runs", []):
                                if run_data.get("type") == "table_cell":
                                    # 恢复表格单元格背景
                                    if shape.has_table:
                                        row_idx = run_data["row"]
                                        col_idx = run_data["col"]
                                        table = shape.table
                                        if (row_idx < len(table.rows) and 
                                            col_idx < len(table.rows[row_idx].cells)):
                                            cell = table.rows[row_idx].cells[col_idx]
                                            bg_fill = run_data.get("background_fill")
                                            if bg_fill and hasattr(cell, 'fill'):
                                                if self._restore_fill(cell.fill, bg_fill):
                                                    restored_count += 1
                                else:
                                    # 恢复文本运行格式
                                    if self._restore_run_formatting(shape, run_data):
                                        restored_count += 1
            
            # 保存文件
            prs.save(ppt_path)
            logger.info(f"颜色恢复完成: {restored_count} 个元素")
            return True
            
        except Exception as e:
            logger.error(f"颜色恢复失败: {e}")
            return False
    
    def _restore_fill(self, fill, fill_data: Dict[str, Any]) -> bool:
        """恢复填充"""
        try:
            fore_color = fill_data.get("fore_color")
            if fore_color:
                fill.solid()
                return self._restore_color(fill.fore_color, fore_color)
            return False
        except Exception as e:
            logger.debug(f"恢复填充失败: {e}")
            return False
    
    def _restore_color(self, color, color_data: Dict[str, Any]) -> bool:
        """恢复颜色"""
        try:
            from pptx.dml.color import RGBColor
            
            color_type = color_data.get("type")
            
            if color_type == 1:  # RGB颜色
                rgb_data = color_data.get("rgb")
                if rgb_data and len(rgb_data) == 3:
                    color.rgb = RGBColor(rgb_data[0], rgb_data[1], rgb_data[2])
                    return True
            
            elif color_type == 2:  # 主题颜色
                theme_color = color_data.get("theme_color")
                brightness = color_data.get("brightness")
                if theme_color:
                    # 主题颜色恢复比较复杂，这里先记录
                    logger.debug(f"需要恢复主题颜色: {theme_color}")
                    return False  # 暂时不支持主题颜色恢复
            
            return False
            
        except Exception as e:
            logger.debug(f"恢复颜色失败: {e}")
            return False
    
    def _restore_run_formatting(self, shape, run_data: Dict[str, Any]) -> bool:
        """恢复文本运行格式"""
        try:
            from pptx.util import Pt
            
            para_idx = run_data["paragraph_index"]
            run_idx = run_data["run_index"]
            table_pos = run_data.get("table_position")
            
            # 获取目标运行
            target_run = None
            
            if table_pos and shape.has_table:
                # 表格文本
                row_idx, col_idx = table_pos
                table = shape.table
                if (row_idx < len(table.rows) and 
                    col_idx < len(table.rows[row_idx].cells)):
                    cell = table.rows[row_idx].cells[col_idx]
                    if para_idx < len(cell.text_frame.paragraphs):
                        paragraph = cell.text_frame.paragraphs[para_idx]
                        if run_idx < len(paragraph.runs):
                            target_run = paragraph.runs[run_idx]
            
            elif shape.has_text_frame:
                # 普通文本
                if para_idx < len(shape.text_frame.paragraphs):
                    paragraph = shape.text_frame.paragraphs[para_idx]
                    if run_idx < len(paragraph.runs):
                        target_run = paragraph.runs[run_idx]
            
            if not target_run:
                return False
            
            # 恢复字体格式
            font_data = run_data.get("font", {})
            restored = False
            
            # 恢复颜色
            color_data = font_data.get("color")
            if color_data:
                if self._restore_color(target_run.font.color, color_data):
                    restored = True
            
            # 恢复字体大小
            font_size = font_data.get("size")
            if font_size:
                target_run.font.size = Pt(font_size)
                restored = True
            
            # 恢复字体名称
            font_name = font_data.get("name")
            if font_name:
                target_run.font.name = font_name
                restored = True
            
            # 恢复字体样式
            if "bold" in font_data:
                target_run.font.bold = font_data["bold"]
                restored = True
            
            if "italic" in font_data:
                target_run.font.italic = font_data["italic"]
                restored = True
            
            if "underline" in font_data:
                # 下划线恢复比较复杂，先跳过
                pass
            
            return restored
            
        except Exception as e:
            logger.debug(f"恢复文本格式失败: {e}")
            return False
    
    def save_backup_to_file(self, backup_path: str, backup_data: Optional[Dict[str, Any]] = None) -> bool:
        """保存备份到文件"""
        try:
            data = backup_data or self.color_backup
            with open(backup_path, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
            logger.info(f"颜色备份已保存到: {backup_path}")
            return True
        except Exception as e:
            logger.error(f"保存备份失败: {e}")
            return False
    
    def load_backup_from_file(self, backup_path: str) -> Dict[str, Any]:
        """从文件加载备份"""
        try:
            with open(backup_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            self.color_backup = data
            logger.info(f"颜色备份已从文件加载: {backup_path}")
            return data
        except Exception as e:
            logger.error(f"加载备份失败: {e}")
            return {}


def backup_and_restore_workflow(ppt_path: str, process_func, *args, **kwargs) -> bool:
    """
    颜色备份和恢复工作流程
    
    Args:
        ppt_path: PPT文件路径
        process_func: 处理函数（如翻译、渲染等）
        *args, **kwargs: 传递给处理函数的参数
        
    Returns:
        bool: 是否成功
    """
    try:
        # 1. 备份颜色
        logger.info("开始备份PPT颜色...")
        backup_manager = PPTColorBackupRestore()
        backup_data = backup_manager.backup_colors_from_ppt(ppt_path)
        
        if not backup_data:
            logger.warning("颜色备份失败，继续处理...")
        
        # 2. 执行处理函数
        logger.info("执行处理函数...")
        result = process_func(ppt_path, *args, **kwargs)
        
        # 3. 恢复颜色
        if backup_data:
            logger.info("恢复PPT颜色...")
            restore_success = backup_manager.restore_colors_to_ppt(ppt_path, backup_data)
            if restore_success:
                logger.info("颜色恢复成功")
            else:
                logger.warning("⚠颜色恢复失败")
        
        return result
        
    except Exception as e:
        logger.error(f"颜色备份恢复工作流程失败: {e}")
        return False


# 便捷函数
def translate_with_color_protection(ppt_path: str, translation_func, *args, **kwargs) -> bool:
    """带颜色保护的翻译"""
    return backup_and_restore_workflow(ppt_path, translation_func, *args, **kwargs)


def render_with_color_protection(ppt_path: str) -> bool:
    """带颜色保护的LibreOffice渲染"""
    from .libreoffice_autofit import libreoffice_ppt_autofit
    return backup_and_restore_workflow(ppt_path, libreoffice_ppt_autofit)


if __name__ == "__main__":
    print("PPT颜色备份和恢复模块")
    print("用于在LibreOffice渲染前后保护颜色格式")
