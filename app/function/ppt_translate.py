import json
import sys
import time
import os

# from mypy.messages import best_matches
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt, Inches
import re
# English-to-Chinese
# from LLa_translate import translate
# 温馨提示: 使用pipeline推理及在线体验功能的时候，尽量输入单句文本，如果是多句长文本建议人工分句，否则可能出现漏译或未译等情况！！！
# from modelscope.pipelines import pipeline
import logging
# from modelscope.utils.constant import Tasks
import difflib
import re

# 导入工具函数
from ..utils.ppt_utils import (
    get_font_color,
    apply_font_color,
    compare_strings_ignore_spaces,
    find_most_similar,
    remove_invalid_utf8_chars,
    is_valid_reference,
    is_page_number
)
from ..utils.translation_utils import build_map
from .translate_by_qwen import translate_qwen, get_field
# 导入其他翻译模型
# from .translate_by_deepseek import translate_deepseek
# from .translate_by_gpt4o import translate_gpt4o
from colorama import init

from ..utils.task_queue import translation_queue

logging.basicConfig(level=logging.INFO, format='%(message)s', encoding='utf-8')
# 初始化 colorama
init()

# 示例字符串
src_language = 'English'  # 源语言
trg_language = 'Chinese'  # 目标语言


def match(text):
    # 使用正则表达式查找被 {} 包裹的内容
    matches = re.findall(r'\{([^}]+)\}', text)
    # 打印匹配到的内容
    # print(matches)


# 字体颜色处理函数已移动到 utils/ppt_utils.py


# def translate(text):
#     outputs = pipeline_ins(input=text)
#     # print(outputs)
#     return outputs['translation']


# 字符串处理函数已移动到 utils/ppt_utils.py


def calculate_translation_similarity(original_text: str, translated_text: str) -> float:
    """
    计算原文和译文的相似度

    Args:
        original_text: 原文
        translated_text: 译文

    Returns:
        相似度分数 (0.0 - 1.0)
    """
    if not original_text or not translated_text:
        return 0.0

    def normalize_text(text):
        # 转小写
        normalized = text.lower()
        # 去除多余空格
        normalized = ' '.join(normalized.split())
        # 去除常见标点符号
        normalized = re.sub(r'[.,!?;:()\[\]{}"\'`~]', '', normalized)
        return normalized.strip()

    norm_original = normalize_text(original_text)
    norm_translated = normalize_text(translated_text)

    if not norm_original or not norm_translated:
        return 0.0

    # 计算字符级相似度
    char_similarity = difflib.SequenceMatcher(None, norm_original, norm_translated).ratio()

    # 计算词级相似度
    words_original = norm_original.split()
    words_translated = norm_translated.split()
    word_similarity = difflib.SequenceMatcher(None, words_original, words_translated).ratio()

    # 综合相似度 (字符相似度权重0.6，词相似度权重0.4)
    combined_similarity = char_similarity * 0.6 + word_similarity * 0.4

    return combined_similarity


def should_skip_translation_insertion(original_text: str, translated_text: str,
                                    threshold: float = 0.9, debug: bool = False) -> bool:
    """
    判断是否应该跳过翻译插入

    Args:
        original_text: 原文
        translated_text: 译文
        threshold: 相似度阈值，默认0.9 (90%)
        debug: 是否输出调试信息

    Returns:
        True表示应该跳过，False表示应该插入
    """
    if not original_text or not translated_text:
        if debug:
            logging.info(f"跳过翻译：空文本 (原文: '{original_text}', 译文: '{translated_text}')")
        return True

    # 如果原文和译文完全相同，跳过
    if original_text.strip() == translated_text.strip():
        if debug:
            logging.info(f"跳过翻译：文本完全相同 ('{original_text}')")
        return True

    # 计算相似度
    similarity = calculate_translation_similarity(original_text, translated_text)
    should_skip = similarity >= threshold

    if debug:
        logging.info(f"相似度检查: '{original_text[:30]}...' vs '{translated_text[:30]}...'")
        logging.info(f"  相似度: {similarity:.3f}, 阈值: {threshold}, 跳过: {should_skip}")

    return should_skip


def save_shape_geometry(shape):
    """保存形状的几何属性"""
    try:
        return {
            'width': shape.width,
            'height': shape.height,
            'left': shape.left,
            'top': shape.top,
            'rotation': getattr(shape, 'rotation', 0)
        }
    except Exception as e:
        logging.debug(f"保存形状几何属性失败: {e}")
        return {}


def restore_shape_geometry(shape, geometry_info):
    """恢复形状的几何属性"""
    try:
        if not geometry_info:
            return False

        # 恢复尺寸和位置
        if 'width' in geometry_info:
            shape.width = geometry_info['width']

        if 'height' in geometry_info:
            shape.height = geometry_info['height']

        if 'left' in geometry_info:
            shape.left = geometry_info['left']

        if 'top' in geometry_info:
            shape.top = geometry_info['top']

        if 'rotation' in geometry_info:
            try:
                shape.rotation = geometry_info['rotation']
            except:
                pass

        return True

    except Exception as e:
        logging.debug(f"恢复形状几何属性失败: {e}")
        return False


def detect_complex_shape_type(shape):
    """检测复杂形状类型"""
    try:
        shape_info = {
            'type': 'simple',
            'has_fill': False,
            'has_line': False,
            'has_shadow': False,
            'is_group': False,
            'is_custom': False,
            'shape_type': shape.shape_type
        }

        # 检查是否为组合形状
        try:
            if hasattr(shape, 'shapes') and shape.shapes:
                shape_info['type'] = 'group'
                shape_info['is_group'] = True
                return shape_info
        except:
            pass

        # 检查填充属性
        try:
            if hasattr(shape, 'fill'):
                fill = shape.fill
                if hasattr(fill, 'type') and fill.type is not None:
                    shape_info['has_fill'] = True
                    if fill.type != 0:  # 不是无填充
                        shape_info['type'] = 'complex'
        except:
            pass

        # 检查线条属性
        try:
            if hasattr(shape, 'line'):
                line = shape.line
                if hasattr(line, 'color') and line.color is not None:
                    shape_info['has_line'] = True
                    shape_info['type'] = 'complex'
        except:
            pass

        # 检查阴影效果
        try:
            if hasattr(shape, 'shadow'):
                shadow = shape.shadow
                if hasattr(shadow, 'visible') and shadow.visible:
                    shape_info['has_shadow'] = True
                    shape_info['type'] = 'complex'
        except:
            pass

        # 检查是否为自定义形状
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                shape_info['is_custom'] = True
                shape_info['type'] = 'custom'
        except:
            pass

        return shape_info

    except Exception as e:
        logging.debug(f"检测形状类型失败: {e}")
        return {'type': 'unknown', 'error': str(e)}


def save_complex_shape_properties(shape):
    """保存复杂形状的完整属性（增强版，更全面的保存）"""
    try:
        properties = {
            'basic_geometry': {
                'width': shape.width,
                'height': shape.height,
                'left': shape.left,
                'top': shape.top,
                'rotation': getattr(shape, 'rotation', 0)
            },
            'shape_properties': {
                'shape_type': shape.shape_type,
                'name': getattr(shape, 'name', ''),
                'shape_id': getattr(shape, 'shape_id', None)
            },
            'fill_properties': {},
            'line_properties': {},
            'shadow_properties': {},
            'text_frame_properties': {},
            'advanced_properties': {}
        }

        # 保存填充属性（增强版）
        try:
            if hasattr(shape, 'fill'):
                fill = shape.fill
                properties['fill_properties'] = {
                    'type': getattr(fill, 'type', None),
                    'transparency': getattr(fill, 'transparency', None),
                    'fore_color_rgb': None,
                    'back_color_rgb': None
                }

                # 保存颜色信息
                try:
                    if hasattr(fill, 'fore_color') and hasattr(fill.fore_color, 'rgb'):
                        properties['fill_properties']['fore_color_rgb'] = fill.fore_color.rgb
                except:
                    pass

                try:
                    if hasattr(fill, 'back_color') and hasattr(fill.back_color, 'rgb'):
                        properties['fill_properties']['back_color_rgb'] = fill.back_color.rgb
                except:
                    pass
        except Exception as e:
            properties['fill_properties']['error'] = str(e)

        # 保存线条属性（完整版，修复边框颜色问题）
        try:
            if hasattr(shape, 'line'):
                line = shape.line
                properties['line_properties'] = {
                    'width': getattr(line, 'width', None),
                    'dash_style': getattr(line, 'dash_style', None),
                    'fill_type': None,
                    'transparency': None,
                    'color_info': None
                }

                # 保存线条填充类型（决定是否有边框）
                try:
                    if hasattr(line, 'fill'):
                        properties['line_properties']['fill_type'] = getattr(line.fill, 'type', None)
                        properties['line_properties']['transparency'] = getattr(line.fill, 'transparency', None)
                except:
                    pass

                # 保存线条颜色（完整版）
                try:
                    if hasattr(line, 'color'):
                        color = line.color
                        color_info = {}

                        # 检查颜色类型
                        try:
                            from pptx.enum.dml import MSO_COLOR_TYPE
                            color_type = getattr(color, 'type', None)

                            if color_type == MSO_COLOR_TYPE.RGB:
                                color_info = {
                                    'color_type': 'rgb',
                                    'rgb_value': getattr(color, 'rgb', None)
                                }
                            elif color_type == MSO_COLOR_TYPE.THEME:
                                color_info = {
                                    'color_type': 'theme',
                                    'theme_color': getattr(color, 'theme_color', None)
                                }
                            elif color_type == MSO_COLOR_TYPE.AUTO:
                                color_info = {'color_type': 'auto'}
                            else:
                                # 尝试获取RGB作为后备
                                color_info = {
                                    'color_type': 'rgb',
                                    'rgb_value': getattr(color, 'rgb', None)
                                }
                        except:
                            # 如果无法获取类型，直接尝试RGB
                            try:
                                color_info = {
                                    'color_type': 'rgb',
                                    'rgb_value': getattr(color, 'rgb', None)
                                }
                            except:
                                color_info = {'color_type': 'error'}

                        properties['line_properties']['color_info'] = color_info

                except Exception as e:
                    properties['line_properties']['color_info'] = {'color_type': 'error', 'error': str(e)}
        except Exception as e:
            properties['line_properties']['error'] = str(e)

        # 保存阴影属性
        try:
            if hasattr(shape, 'shadow'):
                shadow = shape.shadow
                properties['shadow_properties'] = {
                    'visible': getattr(shadow, 'visible', None),
                    'style': getattr(shadow, 'style', None),
                    'blur_radius': getattr(shadow, 'blur_radius', None),
                    'distance': getattr(shadow, 'distance', None),
                    'direction': getattr(shadow, 'direction', None)
                }
        except Exception as e:
            properties['shadow_properties']['error'] = str(e)

        # 保存文本框属性（增强版）
        try:
            if hasattr(shape, 'text_frame'):
                text_frame = shape.text_frame
                properties['text_frame_properties'] = {
                    'auto_size': getattr(text_frame, 'auto_size', None),
                    'word_wrap': getattr(text_frame, 'word_wrap', None),
                    'margin_left': getattr(text_frame, 'margin_left', None),
                    'margin_right': getattr(text_frame, 'margin_right', None),
                    'margin_top': getattr(text_frame, 'margin_top', None),
                    'margin_bottom': getattr(text_frame, 'margin_bottom', None),
                    'vertical_anchor': getattr(text_frame, 'vertical_anchor', None)
                }
        except Exception as e:
            properties['text_frame_properties']['error'] = str(e)

        # 保存高级属性
        try:
            properties['advanced_properties'] = {
                'has_text_frame': getattr(shape, 'has_text_frame', False),
                'has_table': getattr(shape, 'has_table', False),
                'has_chart': getattr(shape, 'has_chart', False),
                'auto_shape_type': getattr(shape, 'auto_shape_type', None)
            }
        except Exception as e:
            properties['advanced_properties']['error'] = str(e)

        return properties

    except Exception as e:
        logging.debug(f"保存复杂形状属性失败: {e}")
        return {}


def restore_complex_shape_properties(shape, properties):
    """恢复复杂形状的完整属性（增强版，更全面的恢复）"""
    try:
        if not properties:
            return False

        success_operations = 0
        total_operations = 0

        # 恢复基本几何属性
        try:
            basic = properties.get('basic_geometry', {})
            if basic:
                total_operations += 4

                if 'width' in basic:
                    shape.width = basic['width']
                    success_operations += 1

                if 'height' in basic:
                    shape.height = basic['height']
                    success_operations += 1

                if 'left' in basic:
                    shape.left = basic['left']
                    success_operations += 1

                if 'top' in basic:
                    shape.top = basic['top']
                    success_operations += 1

                # 旋转属性需要特殊处理
                if 'rotation' in basic:
                    try:
                        shape.rotation = basic['rotation']
                    except:
                        pass
        except Exception as e:
            logging.debug(f"恢复基本几何属性失败: {e}")

        # 恢复填充属性
        try:
            fill_props = properties.get('fill_properties', {})
            if fill_props and hasattr(shape, 'fill') and 'error' not in fill_props:
                fill = shape.fill
                total_operations += 1

                if 'type' in fill_props and fill_props['type'] is not None:
                    try:
                        fill.type = fill_props['type']
                        success_operations += 1
                    except:
                        pass

                if 'transparency' in fill_props and fill_props['transparency'] is not None:
                    try:
                        fill.transparency = fill_props['transparency']
                    except:
                        pass
        except Exception as e:
            logging.debug(f"恢复填充属性失败: {e}")

        # 恢复线条属性（完整版，修复边框颜色问题）
        try:
            line_props = properties.get('line_properties', {})
            if line_props and hasattr(shape, 'line') and 'error' not in line_props:
                line = shape.line
                total_operations += 1

                # 恢复线条宽度
                if 'width' in line_props and line_props['width'] is not None:
                    try:
                        line.width = line_props['width']
                        success_operations += 1
                        logging.debug(f"恢复线条宽度: {line_props['width']}")
                    except Exception as e:
                        logging.debug(f"恢复线条宽度失败: {e}")

                # 恢复线条样式
                if 'dash_style' in line_props and line_props['dash_style'] is not None:
                    try:
                        line.dash_style = line_props['dash_style']
                        logging.debug(f"恢复线条样式: {line_props['dash_style']}")
                    except Exception as e:
                        logging.debug(f"恢复线条样式失败: {e}")

                # 恢复线条填充类型（决定是否有边框）
                if 'fill_type' in line_props and line_props['fill_type'] is not None:
                    try:
                        if hasattr(line, 'fill'):
                            line.fill.type = line_props['fill_type']
                            logging.debug(f"恢复线条填充类型: {line_props['fill_type']}")
                    except Exception as e:
                        logging.debug(f"恢复线条填充类型失败: {e}")

                # 恢复线条透明度
                if 'transparency' in line_props and line_props['transparency'] is not None:
                    try:
                        if hasattr(line, 'fill'):
                            line.fill.transparency = line_props['transparency']
                            logging.debug(f"恢复线条透明度: {line_props['transparency']}")
                    except Exception as e:
                        logging.debug(f"恢复线条透明度失败: {e}")

                # 恢复线条颜色（完整版）
                color_info = line_props.get('color_info')
                if color_info and color_info.get('color_type') != 'error':
                    try:
                        color = line.color
                        color_type = color_info.get('color_type')

                        if color_type == 'rgb' and 'rgb_value' in color_info:
                            rgb_value = color_info['rgb_value']
                            if rgb_value is not None:
                                color.rgb = rgb_value
                                success_operations += 1
                                logging.debug(f"恢复线条RGB颜色: {rgb_value}")

                        elif color_type == 'theme' and 'theme_color' in color_info:
                            theme_color = color_info['theme_color']
                            if theme_color is not None:
                                color.theme_color = theme_color
                                success_operations += 1
                                logging.debug(f"恢复线条主题颜色: {theme_color}")

                        elif color_type == 'auto':
                            from pptx.enum.dml import MSO_COLOR_TYPE
                            color.type = MSO_COLOR_TYPE.AUTO
                            success_operations += 1
                            logging.debug("恢复线条自动颜色")

                        else:
                            # 尝试RGB作为后备
                            if 'rgb_value' in color_info and color_info['rgb_value'] is not None:
                                color.rgb = color_info['rgb_value']
                                success_operations += 1
                                logging.debug(f"使用RGB后备恢复线条颜色: {color_info['rgb_value']}")

                    except Exception as e:
                        logging.debug(f"恢复线条颜色失败: {e}")

        except Exception as e:
            logging.debug(f"恢复线条属性失败: {e}")

        # 恢复文本框属性
        try:
            tf_props = properties.get('text_frame_properties', {})
            if tf_props and hasattr(shape, 'text_frame') and 'error' not in tf_props:
                text_frame = shape.text_frame
                total_operations += 1

                # 恢复边距
                for margin in ['margin_left', 'margin_right', 'margin_top', 'margin_bottom']:
                    if margin in tf_props and tf_props[margin] is not None:
                        try:
                            setattr(text_frame, margin, tf_props[margin])
                        except:
                            pass

                # 恢复其他属性
                for prop in ['word_wrap', 'vertical_anchor']:
                    if prop in tf_props and tf_props[prop] is not None:
                        try:
                            setattr(text_frame, prop, tf_props[prop])
                        except:
                            pass

                success_operations += 1
        except Exception as e:
            logging.debug(f"恢复文本框属性失败: {e}")

        # 计算成功率
        success_rate = (success_operations / max(total_operations, 1)) * 100
        logging.debug(f"形状属性恢复成功率: {success_rate:.1f}% ({success_operations}/{total_operations})")

        return success_operations > 0

    except Exception as e:
        logging.debug(f"恢复复杂形状属性失败: {e}")
        return False


def has_shape_deformed(shape, original_properties):
    """检查形状是否发生变形（增强版，更严格的检测）"""
    try:
        if not original_properties:
            return False

        basic = original_properties.get('basic_geometry', {})
        if not basic:
            return False

        # 使用更严格的容差 - 从2降低到0.5
        tolerance = 0.5  # 0.5个单位的容差，更敏感地检测变形

        # 检查宽度变化
        if 'width' in basic:
            width_diff = abs(shape.width - basic['width'])
            if width_diff > tolerance:
                logging.warning(f"检测到宽度变化: {width_diff:.2f} > {tolerance}")
                return True

        # 检查高度变化
        if 'height' in basic:
            height_diff = abs(shape.height - basic['height'])
            if height_diff > tolerance:
                logging.warning(f"检测到高度变化: {height_diff:.2f} > {tolerance}")
                return True

        # 检查左边距变化
        if 'left' in basic:
            left_diff = abs(shape.left - basic['left'])
            if left_diff > tolerance:
                logging.warning(f"检测到左边距变化: {left_diff:.2f} > {tolerance}")
                return True

        # 检查上边距变化
        if 'top' in basic:
            top_diff = abs(shape.top - basic['top'])
            if top_diff > tolerance:
                logging.warning(f"检测到上边距变化: {top_diff:.2f} > {tolerance}")
                return True

        return False

    except Exception as e:
        logging.debug(f"检查形状变形失败: {e}")
        return True  # 出错时假设已变形，采用保守策略


def safe_set_autofit_with_size_preservation(text_frame, shape):
    """安全地设置自适应并保护形状大小（增强版，支持复杂形状）"""
    try:
        # 1. 检测形状类型
        shape_info = detect_complex_shape_type(shape)
        logging.debug(f"检测到形状类型: {shape_info['type']}")

        # 2. 保存完整属性
        original_properties = save_complex_shape_properties(shape)

        # 3. 根据形状类型选择策略
        if shape_info['type'] == 'group':
            # 组合形状：不设置自适应，避免破坏组合结构
            logging.info("跳过组合形状的自适应设置")
            return True

        elif shape_info['type'] == 'complex':
            # 复杂形状：谨慎设置自适应
            logging.debug("为复杂形状设置自适应")

            # 保存当前状态
            current_auto_size = text_frame.auto_size

            try:
                # 设置自适应
                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                # 立即检查是否有变形
                if has_shape_deformed(shape, original_properties):
                    logging.warning("检测到复杂形状变形，恢复原始状态")
                    # 恢复原始自适应设置
                    text_frame.auto_size = current_auto_size
                    # 恢复所有属性
                    restore_complex_shape_properties(shape, original_properties)
                    return False
                else:
                    logging.debug("复杂形状未变形，保持自适应设置")
                    # 恢复其他属性，保持自适应
                    restore_complex_shape_properties(shape, original_properties)
                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    return True

            except Exception as e:
                logging.warning(f"复杂形状设置自适应失败: {e}")
                # 恢复原始状态
                text_frame.auto_size = current_auto_size
                restore_complex_shape_properties(shape, original_properties)
                return False

        else:
            # 简单形状或自定义形状：使用标准方法
            logging.debug(f"为{shape_info['type']}形状设置自适应")
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            restore_complex_shape_properties(shape, original_properties)
            return True

    except Exception as e:
        logging.error(f"设置自适应失败: {e}")
        return False


def process_presentation_add_annotations(path_to_presentation, annotations, stop_words_list, custom_translations,
                                         source_language, target_language, bilingual_translation):
    prs = Presentation(path_to_presentation)
    all_text = ""
    for slide in prs.slides:
        # 遍历每个形状（包含文本的元素，如文本框）
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:  # 忽略空文本
                        # print(text)
                        # text=text+"\n"
                        all_text += text + "\n"
    field = str(get_field(all_text))
    tage_text = ""
    annotations = annotations["annotations"]
    for item in annotations:
        text = item["ocrResult"].replace("\n", " ")
        tage_text += text + "\n"
    stop_words = list()
    custom_words = dict()
    for i in stop_words_list:
        if i in tage_text:
            stop_words.append(i)
    for k, v in custom_translations.items():
        if k in tage_text:
            custom_words[k] = v
    data = translate_qwen(tage_text, field, stop_words, custom_words, source_language, target_language)
    for item in annotations:
        page = item["page"]
        original_text = item["ocrResult"]
        slide = prs.slides[page - 1]
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        left = slide_width - Inches(2)  # 文本框的左边距
        top = 0  # 文本框的上边距
        new_text = find_most_similar(original_text, list(data.keys()))
        # 设置文本框的宽度和高度
        width = Inches(2)
        height = Inches(1)

        # 添加文本框
        shape = slide.shapes.add_textbox(left, top, width, height)
        text_frame = shape.text_frame
        original_text = re.sub(r'_x000B_|\u000B', '', original_text)
        translated_text=data[new_text]
        translated_text = re.sub(r'_x000B_|\u000B', '', translated_text)
        if str(bilingual_translation) == "1":
            text_frame.text = original_text + "\n" + translated_text
        else:
            text_frame.text = data[new_text]

        # 设置文本框中文字的字体和颜色（保持注释功能的红色，但可配置）
        for p in text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(14)  # 设置字体大小
                # 注释功能使用红色字体以便区分，这是预期行为
                run.font.color.rgb = RGBColor(255, 0, 0)  # 设置字体颜色为红色
    prs.save(path_to_presentation)


def is_valid_reference(text):
    pattern = r'\d+\s*[A-Za-z&\s\.\-]+,\s*\d{4}'
    return bool(re.match(pattern, text))
def is_page_number(text):
    text = text.strip()

    # 常见纯数字页码
    if re.fullmatch(r'\d{1,3}', text):
        return True
    return False

def process_presentation(path_to_presentation, stop_words_list, custom_translations, select_page, source_language,
                         target_language, bilingual_translation, enable_uno_conversion=True):
    logging.info(f"开始处理演示文稿: {os.path.basename(path_to_presentation)}")
    logging.info(f"源语言: {source_language}, 目标语言: {target_language}, 双语翻译: {bilingual_translation}")
    logging.info(f"选中页面: {select_page}")

    try:
        # 加载演示文稿
        logging.info("正在加载演示文稿...")
        prs = Presentation(path_to_presentation)
        total_slides = len(prs.slides)
        logging.info(f"演示文稿加载成功，共 {total_slides} 张幻灯片")

        # 收集所有文本
        logging.info("正在收集所有文本...")
        all_text = ""
        for slide in prs.slides:
            # 遍历每个形状（包含文本的元素，如文本框）
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:  # 忽略空文本
                            all_text += text + "\n"

        logging.info(f"文本收集完成，共收集 {len(all_text)} 个字符")

        # 获取领域
        logging.info("正在分析文本领域...")
        field = str(get_field(all_text))
        logging.info(f"文本领域分析结果: {field}")

        # 处理每张幻灯片
        processed_slides = 0
        skipped_slides = 0

        for current_slide_index, slide in enumerate(prs.slides, 1):
            # 更新翻译进度
            # 由于缺少task_id参数，这里暂时使用用户ID进行更新
            # 这是一个临时解决方案，理想情况下，应该从调用方传入task_id
            # 在实际使用时，应该从任务上下文中获取用户ID，而不是使用硬编码的值
            translation_queue.update_progress_by_user(1, current_slide_index, total_slides)  # 替换为按用户ID更新进度

            # 检查是否需要处理当前幻灯片
            if current_slide_index not in select_page:
                logging.info(f"跳过第 {current_slide_index} 张幻灯片 (不在选中页面列表中)")
                skipped_slides += 1
                continue

            logging.info(f"开始处理第 {current_slide_index} 张幻灯片 ({current_slide_index}/{total_slides})")

            # 收集当前幻灯片的文本
            slide_text = ""
            # 遍历每个形状（包含文本的元素，如文本框）
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:  # 忽略空文本
                            slide_text += text + "\n"

            # 处理表格
            table_count = 0
            for shape in slide.shapes:
                if shape.has_table:  # 检查该形状是否为表格
                    table_count += 1
                    logging.info(f"处理第 {current_slide_index} 张幻灯片中的表格 #{table_count}")
                    # 获取表格
                    table = shape.table
                    # 遍历表格中的每一行
                    for row in table.rows:
                        # 遍历每一列
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    # 获取单元格的文本
                                    text = run.text.strip()
                                    slide_text += "【"+text+"】" + "\n"

            logging.info(f"第 {current_slide_index} 张幻灯片文本收集完成，共 {len(slide_text)} 个字符")

            # 筛选停止词和自定义翻译
            logging.info("正在筛选适用的停止词和自定义翻译...")
            stop_words = list()
            custom_words = dict()
            for i in stop_words_list:
                if i in slide_text:
                    stop_words.append(i)

            for k, v in custom_translations.items():
                if k in slide_text:
                    custom_words[k] = v

            logging.info(f"应用 {len(stop_words)} 个停止词和 {len(custom_words)} 个自定义翻译")

            # 翻译文本
            logging.info(f"开始翻译第 {current_slide_index} 张幻灯片...")
            data = translate_local_qwen(slide_text, field, stop_words, custom_words, source_language, target_language)
            logging.info(f"翻译完成，获得 {len(list(data.keys()))} 个翻译结果")

            # 应用翻译结果
            logging.info(f"开始应用翻译结果到第 {current_slide_index} 张幻灯片...")
            text_blocks_updated = 0

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                    for paragraph in text_frame.paragraphs:
                        original_text = paragraph.text.strip()
                        original_text = remove_invalid_utf8_chars(original_text)
                        if original_text:  # Ensure text is not empty
                            if is_valid_reference(original_text):
                                continue

                            # 保存原始颜色
                            if paragraph.runs:
                                original_color = get_font_color(paragraph.runs[0])
                            else:
                                original_color = None

                            # 查找翻译
                            translated_text = ""
                            new_text = find_most_similar(original_text, list(data.keys()))
                            if new_text in data:
                                clean_text1 = re.sub(r'[^\w]', '', original_text)
                                clean_text2 = re.sub(r'[^\w]', '', data[new_text])
                                if clean_text1 != clean_text2:
                                    translated_text = data[new_text]

                            # 应用翻译
                            original_text=re.sub(r'_x000B_|\u000B', '', original_text)
                            translated_text=re.sub(r'_x000B_|\u000B', '', translated_text)
                            if not is_page_number(original_text):
                                if translated_text != original_text and translated_text:
                                    # 检查相似度，如果相似度过高则跳过翻译
                                    if should_skip_translation_insertion(original_text, translated_text, threshold=0.9, debug=True):
                                        logging.info(f"跳过高相似度翻译: '{original_text[:30]}...' -> '{translated_text[:30]}...'")
                                        continue

                                    text_blocks_updated += 1
                                    paragraph.clear()
                                    run = paragraph.add_run()
                                    if str(bilingual_translation) == "1":
                                        run.text = original_text + "\n" + translated_text
                                    else:
                                        run.text = translated_text
                                    run.font.size = Pt(24)  # 例如字体大小，可根据需要调整

                            # 恢复颜色
                            if original_color:
                                apply_font_color(run, original_color)

                            # 安全地设置自适应并保护形状大小
                            safe_set_autofit_with_size_preservation(text_frame, shape)
                elif shape.has_table:  # 检查该形状是否为表格
                    # 处理表格翻译
                    table = shape.table
                    cells_updated = 0
                    # 遍历表格中的每一行
                    for row in table.rows:
                        # 遍历每一列
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.size = Pt(10)
                                    # 获取单元格的文本
                                    new_text = find_most_similar(run.text, list(data.keys()))
                                    if new_text is None or new_text not in data:
                                        clean_text2 = ''
                                    else:
                                        clean_text1 = re.sub(r'[^\w]', '', run.text)
                                        clean_text2 = re.sub(r'[^\w]', '', data[new_text])

                                        if clean_text1 != clean_text2:
                                            # 检查相似度，如果相似度过高则跳过翻译
                                            if should_skip_translation_insertion(run.text, data[new_text], threshold=0.9, debug=True):
                                                logging.info(f"跳过表格高相似度翻译: '{run.text[:30]}...' -> '{data[new_text][:30]}...'")
                                                continue

                                            cells_updated += 1
                                            if str(bilingual_translation) == "1":
                                                run.text = run.text + "\n" + data[new_text] + "\n"
                                            else:
                                                run.text = data[new_text] + "\n"

            logging.info(f"第 {current_slide_index} 张幻灯片处理完成，更新了 {text_blocks_updated} 个文本块")
            processed_slides += 1

        # 保存演示文稿
        logging.info("正在保存演示文稿...")
        prs.save(path_to_presentation)
        logging.info(f"演示文稿处理完成: 处理了 {processed_slides} 张幻灯片，跳过了 {skipped_slides} 张幻灯片")
        return True
    except Exception as e:
        logging.error(f"处理演示文稿时出错: {str(e)}")
        import traceback
        logging.error(traceback.format_exc())
        return False

def has_meaningful_text_content(text_frame):
    """
    检测文本框是否包含有意义的文字内容

    Args:
        text_frame: PPT文本框对象

    Returns:
        bool: True表示包含有意义的文字，False表示只有形状或空白
    """
    try:
        if not text_frame or not hasattr(text_frame, 'paragraphs'):
            return False

        # 检查所有段落
        total_text = ""
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text:
                    total_text += run.text

        # 去除空白字符
        total_text = total_text.strip()

        # 如果没有文本，返回False
        if not total_text:
            return False

        # 检查是否只是空白字符、换行符等
        if re.match(r'^[\s\n\r\t]*$', total_text):
            return False

        # 检查是否只是纯数字（页码等）
        if re.match(r'^[\d\s\.,\-%]+$', total_text):
            return False

        # 检查是否只是纯标点符号
        if re.match(r'^[^\w\s]+$', total_text):
            return False

        # 检查是否只是单个字符
        if len(total_text) <= 1:
            return False

        # 检查是否只是特殊字符
        if re.match(r'^[\s\-_=+\*#@$%^&()]+$', total_text):
            return False

        # 如果通过了所有检查，认为包含有意义的文字
        return True

    except Exception as e:
        # 出错时保守处理，认为包含文字（避免跳过需要处理的文本框）
        logging.debug(f"检测文本内容时出错: {e}")
        return True


def should_adjust_textbox_layout(shape):
    """
    判断是否应该调整文本框布局

    Args:
        shape: PPT形状对象

    Returns:
        bool: True表示应该调整，False表示跳过
    """
    try:
        # 检查是否有文本框
        if not shape.has_text_frame:
            return False

        text_frame = shape.text_frame

        # 检查是否包含有意义的文字内容
        if not has_meaningful_text_content(text_frame):
            return False

        return True

    except Exception as e:
        # 出错时保守处理，进行调整（避免跳过需要处理的文本框）
        logging.debug(f"判断是否调整文本框时出错: {e}")
        return True


def get_textbox_content_summary(text_frame):
    """
    获取文本框内容摘要（用于调试）

    Args:
        text_frame: PPT文本框对象

    Returns:
        str: 内容摘要
    """
    try:
        if not text_frame or not hasattr(text_frame, 'paragraphs'):
            return "无文本框"

        total_text = ""
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text:
                    total_text += run.text

        total_text = total_text.strip()

        if not total_text:
            return "空白"
        elif len(total_text) > 30:
            return f"'{total_text[:30]}...'"
        else:
            return f"'{total_text}'"

    except Exception as e:
        return f"检测失败: {e}"


def safe_set_autofit_with_content_check(text_frame, shape, debug=False):
    """
    安全地设置自适应，只对包含文字内容的文本框进行调整

    Args:
        text_frame: 文本框对象
        shape: 形状对象
        debug: 是否输出调试信息

    Returns:
        dict: 处理结果信息
    """
    try:
        # 检查是否应该调整
        if not should_adjust_textbox_layout(shape):
            if debug:
                content_summary = get_textbox_content_summary(text_frame)
                logging.info(f"跳过文本框调整: {content_summary}")

            return {
                'adjusted': False,
                'reason': 'no_meaningful_content',
                'content': get_textbox_content_summary(text_frame)
            }

        # 包含有意义的文字，进行调整
        if debug:
            content_summary = get_textbox_content_summary(text_frame)
            logging.info(f"调整文本框: {content_summary}")

        # 使用现有的复杂形状处理逻辑
        success = safe_set_autofit_with_size_preservation(text_frame, shape)

        return {
            'adjusted': True,
            'success': success,
            'reason': 'has_meaningful_content',
            'content': get_textbox_content_summary(text_frame)
        }

    except Exception as e:
        if debug:
            logging.error(f"文本框调整出错: {e}")

        return {
            'adjusted': False,
            'success': False,
            'reason': 'error',
            'error': str(e)
        }


# start_time = time.time()
# process_presentation("./test.pptx")
# end_time = time.time()
# print(f"Training took {end_time - start_time:.2f} seconds")
# 读取现有的PPT文件
# print(pipeline_ins.model)
# print(pipeline_ins)
# outs=translate("These models have an advantage over naïve OLS models in that they predict the effect of change in neighbourhood characteristics on change in educational commitment, and therefore offer a more dynamic approach to modelling neighbourhood effects.")
# print(outs)
