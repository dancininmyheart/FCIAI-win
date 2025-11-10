"""
PPT颜色保护模块
确保在翻译和文本框自适应处理过程中保持原始颜色
"""
import logging
from typing import Dict, List, Tuple, Optional, Any

try:
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.error("python-pptx不可用")

logger = logging.getLogger(__name__)


class ColorProtector:
    """PPT颜色保护器"""
    
    def __init__(self):
        self.saved_colors = {}
    
    def save_text_colors(self, text_frame, identifier: str = None) -> Dict[str, Any]:
        """
        保存文本框中所有文本的颜色信息
        
        Args:
            text_frame: PPT文本框对象
            identifier: 标识符，用于区分不同的文本框
            
        Returns:
            Dict: 保存的颜色信息
        """
        if not PPTX_AVAILABLE:
            return {}
        
        color_info = {
            'identifier': identifier,
            'paragraphs': []
        }
        
        try:
            for para_idx, paragraph in enumerate(text_frame.paragraphs):
                para_colors = {
                    'paragraph_index': para_idx,
                    'runs': []
                }
                
                for run_idx, run in enumerate(paragraph.runs):
                    if run.text.strip():  # 只处理有内容的run
                        run_color = self._extract_run_color(run)
                        run_highlight = self._extract_run_highlight(run)
                        run_info = {
                            'run_index': run_idx,
                            'text': run.text,
                            'color': run_color,
                            'highlight': run_highlight,
                            'font_name': run.font.name,
                            'font_size': run.font.size,
                            'font_bold': run.font.bold,
                            'font_italic': run.font.italic,
                            'font_underline': run.font.underline
                        }
                        para_colors['runs'].append(run_info)
                
                color_info['paragraphs'].append(para_colors)
            
            # 保存到实例变量中
            if identifier:
                self.saved_colors[identifier] = color_info
            
            return color_info
            
        except Exception as e:
            logger.error(f"保存文本颜色失败: {e}")
            return {}
    
    def restore_text_colors(self, text_frame, color_info: Dict[str, Any]) -> bool:
        """
        恢复文本框中所有文本的颜色信息
        
        Args:
            text_frame: PPT文本框对象
            color_info: 之前保存的颜色信息
            
        Returns:
            bool: 恢复是否成功
        """
        if not PPTX_AVAILABLE or not color_info:
            return False
        
        try:
            success_count = 0
            total_count = 0
            
            for para_info in color_info.get('paragraphs', []):
                para_idx = para_info.get('paragraph_index', 0)
                
                # 确保段落存在
                if para_idx < len(text_frame.paragraphs):
                    paragraph = text_frame.paragraphs[para_idx]
                    
                    for run_info in para_info.get('runs', []):
                        run_idx = run_info.get('run_index', 0)
                        total_count += 1
                        
                        # 确保run存在
                        if run_idx < len(paragraph.runs):
                            run = paragraph.runs[run_idx]
                            
                            # 恢复颜色和高亮
                            color_restored = self._apply_run_color(run, run_info.get('color'))
                            highlight_restored = self._apply_run_highlight(run, run_info.get('highlight'))
                            if color_restored or highlight_restored:
                                success_count += 1
                            
                            # 恢复其他字体属性（除了大小，因为可能需要自适应）
                            try:
                                if run_info.get('font_name'):
                                    run.font.name = run_info['font_name']
                                if run_info.get('font_bold') is not None:
                                    run.font.bold = run_info['font_bold']
                                if run_info.get('font_italic') is not None:
                                    run.font.italic = run_info['font_italic']
                                if run_info.get('font_underline') is not None:
                                    run.font.underline = run_info['font_underline']
                            except Exception as e:
                                logger.debug(f"恢复字体属性失败: {e}")
            
            logger.debug(f"颜色恢复成功率: {success_count}/{total_count}")
            return success_count > 0
            
        except Exception as e:
            logger.error(f"恢复文本颜色失败: {e}")
            return False
    
    def _extract_run_color(self, run) -> Optional[Dict[str, Any]]:
        """提取run的颜色信息"""
        try:
            color = run.font.color
            color_info = {
                'type': None,
                'value': None
            }

            if color.type == MSO_COLOR_TYPE.RGB:
                # RGB颜色
                rgb = color.rgb
                color_info['type'] = 'rgb'
                color_info['value'] = (rgb.r, rgb.g, rgb.b)

            elif color.type == MSO_COLOR_TYPE.THEME:
                # 主题颜色
                color_info['type'] = 'theme'
                color_info['value'] = color.theme_color

            elif color.type == MSO_COLOR_TYPE.AUTO:
                # 自动颜色
                color_info['type'] = 'auto'
                color_info['value'] = None

            else:
                # 其他类型，尝试获取RGB值
                try:
                    rgb = color.rgb
                    color_info['type'] = 'rgb'
                    color_info['value'] = (rgb.r, rgb.g, rgb.b)
                except:
                    color_info['type'] = 'unknown'
                    color_info['value'] = None

            return color_info

        except Exception as e:
            logger.debug(f"提取颜色信息失败: {e}")
            return None

    def _extract_run_highlight(self, run) -> Optional[Dict[str, Any]]:
        """提取run的高亮色信息"""
        try:
            # 检查是否有高亮色
            highlight_info = {
                'type': None,
                'value': None
            }

            # 尝试获取高亮色（不同版本的python-pptx可能有不同的属性名）
            highlight_attrs = ['highlight_color', 'highlight', 'background_color']

            for attr in highlight_attrs:
                if hasattr(run.font, attr):
                    highlight = getattr(run.font, attr)
                    if highlight:
                        try:
                            if hasattr(highlight, 'type'):
                                if highlight.type == MSO_COLOR_TYPE.RGB:
                                    rgb = highlight.rgb
                                    highlight_info['type'] = 'rgb'
                                    highlight_info['value'] = (rgb.r, rgb.g, rgb.b)
                                elif highlight.type == MSO_COLOR_TYPE.THEME:
                                    highlight_info['type'] = 'theme'
                                    highlight_info['value'] = highlight.theme_color
                                elif highlight.type == MSO_COLOR_TYPE.AUTO:
                                    highlight_info['type'] = 'auto'
                                    highlight_info['value'] = None
                                break
                        except:
                            continue

            return highlight_info if highlight_info['type'] else None

        except Exception as e:
            logger.debug(f"提取高亮信息失败: {e}")
            return None
    
    def _apply_run_color(self, run, color_info: Optional[Dict[str, Any]]) -> bool:
        """应用颜色信息到run"""
        if not color_info:
            return False

        try:
            color_type = color_info.get('type')
            color_value = color_info.get('value')

            if color_type == 'rgb' and color_value:
                # 应用RGB颜色
                r, g, b = color_value
                run.font.color.rgb = RGBColor(r, g, b)
                return True

            elif color_type == 'theme' and color_value:
                # 应用主题颜色
                run.font.color.theme_color = color_value
                return True

            elif color_type == 'auto':
                # 应用自动颜色
                run.font.color.type = MSO_COLOR_TYPE.AUTO
                return True

            return False

        except Exception as e:
            logger.debug(f"应用颜色失败: {e}")
            return False

    def _apply_run_highlight(self, run, highlight_info: Optional[Dict[str, Any]]) -> bool:
        """应用高亮色信息到run"""
        if not highlight_info:
            return False

        try:
            highlight_type = highlight_info.get('type')
            highlight_value = highlight_info.get('value')

            # 尝试不同的高亮属性名
            highlight_attrs = ['highlight_color', 'highlight', 'background_color']

            for attr in highlight_attrs:
                if hasattr(run.font, attr):
                    highlight = getattr(run.font, attr)
                    if highlight:
                        try:
                            if highlight_type == 'rgb' and highlight_value:
                                r, g, b = highlight_value
                                highlight.rgb = RGBColor(r, g, b)
                                return True
                            elif highlight_type == 'theme' and highlight_value:
                                highlight.theme_color = highlight_value
                                return True
                            elif highlight_type == 'auto':
                                highlight.type = MSO_COLOR_TYPE.AUTO
                                return True
                        except:
                            continue

            return False

        except Exception as e:
            logger.debug(f"应用高亮失败: {e}")
            return False
    
    def get_saved_colors(self, identifier: str) -> Optional[Dict[str, Any]]:
        """获取保存的颜色信息"""
        return self.saved_colors.get(identifier)
    
    def clear_saved_colors(self, identifier: str = None):
        """清除保存的颜色信息"""
        if identifier:
            self.saved_colors.pop(identifier, None)
        else:
            self.saved_colors.clear()


# 全局颜色保护器实例
_global_color_protector = ColorProtector()


def save_textframe_colors(text_frame, identifier: str = None) -> Dict[str, Any]:
    """
    保存文本框颜色（便捷函数）
    
    Args:
        text_frame: PPT文本框对象
        identifier: 标识符
        
    Returns:
        Dict: 颜色信息
    """
    return _global_color_protector.save_text_colors(text_frame, identifier)


def restore_textframe_colors(text_frame, color_info: Dict[str, Any]) -> bool:
    """
    恢复文本框颜色（便捷函数）
    
    Args:
        text_frame: PPT文本框对象
        color_info: 颜色信息
        
    Returns:
        bool: 是否成功
    """
    return _global_color_protector.restore_text_colors(text_frame, color_info)


def protect_colors_during_processing(text_frame, processing_func, *args, **kwargs):
    """
    在处理过程中保护颜色

    Args:
        text_frame: 文本框对象
        processing_func: 处理函数
        *args, **kwargs: 传递给处理函数的参数

    Returns:
        处理函数的返回值
    """
    # 保存颜色
    color_info = save_textframe_colors(text_frame)

    try:
        # 执行处理函数
        result = processing_func(*args, **kwargs)

        # 恢复颜色
        restore_textframe_colors(text_frame, color_info)

        return result

    except Exception as e:
        # 即使处理失败，也尝试恢复颜色
        restore_textframe_colors(text_frame, color_info)
        raise e


def safe_replace_paragraph_text(paragraph, new_text, preserve_formatting=True) -> bool:
    """
    安全地替换段落文本，保持格式

    Args:
        paragraph: PPT段落对象
        new_text: 要写入的完整新文本（已在上层处理双语/单语逻辑）
        preserve_formatting: 是否尽量保持原有格式

    Returns:
        bool: 是否成功
    """
    if not PPTX_AVAILABLE:
        return False

    try:
        # 保存原始格式
        original_formats = []
        if preserve_formatting and paragraph.runs:
            for run in paragraph.runs:
                if run.text.strip():
                    run_format = {
                        'font_name': run.font.name,
                        'font_size': run.font.size,
                        'font_bold': run.font.bold,
                        'font_italic': run.font.italic,
                        'font_underline': run.font.underline,
                        'color': None,
                        'highlight': None
                    }

                    # 保存颜色
                    try:
                        color = run.font.color
                        if color.type == MSO_COLOR_TYPE.RGB:
                            rgb = color.rgb
                            run_format['color'] = ('rgb', (rgb.r, rgb.g, rgb.b))
                        elif color.type == MSO_COLOR_TYPE.THEME:
                            run_format['color'] = ('theme', color.theme_color)
                        elif color.type == MSO_COLOR_TYPE.AUTO:
                            run_format['color'] = ('auto', None)
                    except:
                        pass

                    # 保存高亮色
                    # try:
                    #     highlight_attrs = ['highlight_color', 'highlight', 'background_color']
                    #     for attr in highlight_attrs:
                    #         if hasattr(run.font, attr):
                    #             highlight = getattr(run.font, attr)
                    #             if highlight and hasattr(highlight, 'type'):
                    #                 if highlight.type == MSO_COLOR_TYPE.RGB:
                    #                     rgb = highlight.rgb
                    #                     run_format['highlight'] = ('rgb', (rgb.r, rgb.g, rgb.b))
                    #                     break
                    #                 elif highlight.type == MSO_COLOR_TYPE.THEME:
                    #                     run_format['highlight'] = ('theme', highlight.theme_color)
                    #                     break
                    # except:
                    #     pass

                    original_formats.append(run_format)
                    break  # 只取第一个有内容的run的格式

        # 直接写入新文本（上层已决定双语/单语形式）
        paragraph.text = new_text
        # 恢复格式
        if preserve_formatting and original_formats:
            format_info = original_formats[0]  # 使用第一个格式

            try:
                # 尝试恢复第一个 run 的格式
                if paragraph.runs:
                    run = paragraph.runs[0]
                    if format_info.get('font_name'):
                        run.font.name = format_info['font_name']
                    if format_info.get('font_size'):
                        run.font.size = format_info['font_size']
                    if format_info.get('font_bold') is not None:
                        run.font.bold = format_info['font_bold']
                    if format_info.get('font_italic') is not None:
                        run.font.italic = format_info['font_italic']
                    if format_info.get('font_underline') is not None:
                        run.font.underline = format_info['font_underline']

                # 恢复颜色
                if paragraph.runs and format_info.get('color'):
                    run = paragraph.runs[0]
                    color_type, color_value = format_info['color']
                    if color_type == 'rgb' and color_value:
                        r, g, b = color_value
                        run.font.color.rgb = RGBColor(r, g, b)
                    elif color_type == 'theme' and color_value:
                        run.font.color.theme_color = color_value
                    elif color_type == 'auto':
                        run.font.color.type = MSO_COLOR_TYPE.AUTO



            except Exception as e:
                logger.debug(f"恢复段落格式失败: {e}")

        return True

    except Exception as e:
        logger.error(f"安全替换段落文本失败: {e}")
        return False


def ensure_color_consistency(ppt_path: str) -> bool:
    """
    确保PPT文件的颜色一致性
    
    Args:
        ppt_path: PPT文件路径
        
    Returns:
        bool: 处理是否成功
    """
    if not PPTX_AVAILABLE:
        return False
    
    try:
        from pptx import Presentation
        
        prs = Presentation(ppt_path)
        processed_count = 0
        
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    
                    # 检查并修复无颜色的文本
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip():
                                try:
                                    # 检查是否有颜色设置
                                    color_type = run.font.color.type
                                    if color_type is None or color_type == 0:
                                        # 设置默认黑色
                                        run.font.color.rgb = RGBColor(0, 0, 0)
                                        processed_count += 1
                                except:
                                    # 如果获取颜色失败，设置为黑色
                                    run.font.color.rgb = RGBColor(0, 0, 0)
                                    processed_count += 1
        
        if processed_count > 0:
            prs.save(ppt_path)
            logger.info(f"修复了 {processed_count} 个无颜色文本")
        
        return True
        
    except Exception as e:
        logger.error(f"确保颜色一致性失败: {e}")
        return False
