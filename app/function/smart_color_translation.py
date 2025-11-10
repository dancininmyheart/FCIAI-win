"""
智能颜色保护翻译模块
自动选择最佳的颜色保护方法，解决UNO接口Python版本冲突问题
"""
import os
import logging
from typing import Dict, List, Any, Optional, Tuple

logger = logging.getLogger(__name__)


class SmartColorTranslator:
    """智能颜色保护翻译器"""

    def __init__(self):
        self.available_methods = []
        self.preferred_method = None
        self._detect_available_methods()

    def _detect_available_methods(self):
        """检测可用的翻译方法"""
        methods = []

        # 方法1: LibreOffice外部处理器（推荐，避免Python版本冲突）
        try:
            from .libreoffice_uno_alternative import LibreOfficeExternalProcessor
            processor = LibreOfficeExternalProcessor()
            if processor.find_libreoffice_executable():
                methods.append(('uno_external', 'LibreOffice外部处理器', 95))
                logger.info("✅ LibreOffice外部处理器可用（推荐）")
        except ImportError:
            logger.debug("LibreOffice外部处理器不可用")
        
        # 方法2: LibreOffice外部处理器
        try:
            from .libreoffice_uno_alternative import LibreOfficeExternalProcessor
            processor = LibreOfficeExternalProcessor()
            if processor.find_libreoffice_executable():
                methods.append(('uno_external', 'UNO外部处理器', 90))
                logger.info("✅ UNO外部处理器可用")
        except ImportError:
            logger.debug("UNO外部处理器不可用")
        
        # 方法3: 增强的python-pptx颜色保护
        try:
            from .color_protection import ColorProtector
            methods.append(('enhanced_pptx', '增强python-pptx', 80))
            logger.info("✅ 增强python-pptx颜色保护可用")
        except ImportError:
            logger.debug("增强颜色保护不可用")
        
        # 方法4: 页面翻译（已有颜色保护）
        try:
            from .page_based_translation import translate_slide_by_page
            methods.append(('page_based', '页面翻译', 70))
            logger.info("✅ 页面翻译可用")
        except ImportError:
            logger.debug("页面翻译不可用")
        
        # 方法5: 基础python-pptx
        try:
            from pptx import Presentation
            methods.append(('basic_pptx', '基础python-pptx', 60))
            logger.info("✅ 基础python-pptx可用")
        except ImportError:
            logger.debug("python-pptx不可用")
        
        # 按优先级排序
        methods.sort(key=lambda x: x[2], reverse=True)
        self.available_methods = methods
        
        if methods:
            self.preferred_method = methods[0][0]
            logger.info(f"首选翻译方法: {methods[0][1]}")
        else:
            logger.error("没有可用的翻译方法")
    
    def get_available_methods(self) -> List[Tuple[str, str, int]]:
        """获取可用方法列表"""
        return self.available_methods.copy()
    
    def translate_ppt_with_best_color_protection(
        self,
        ppt_path: str,
        translation_data: Dict[str, str],
        output_path: str = None,
        bilingual_mode: bool = False,
        preferred_method: str = None
    ) -> Tuple[bool, str]:
        """
        使用最佳颜色保护方法翻译PPT
        
        Args:
            ppt_path: PPT文件路径
            translation_data: 翻译数据
            output_path: 输出路径
            bilingual_mode: 双语模式
            preferred_method: 首选方法
            
        Returns:
            Tuple[bool, str]: (是否成功, 使用的方法)
        """
        if not os.path.exists(ppt_path):
            logger.error(f"PPT文件不存在: {ppt_path}")
            return False, "file_not_found"
        
        if not self.available_methods:
            logger.error("没有可用的翻译方法")
            return False, "no_methods"
        
        # 确定尝试顺序
        methods_to_try = []
        if preferred_method:
            # 如果指定了首选方法，先尝试它
            for method_id, method_name, priority in self.available_methods:
                if method_id == preferred_method:
                    methods_to_try.append((method_id, method_name, priority))
                    break
        
        # 添加其他方法
        for method_id, method_name, priority in self.available_methods:
            if not preferred_method or method_id != preferred_method:
                methods_to_try.append((method_id, method_name, priority))
        
        # 逐个尝试方法
        for method_id, method_name, priority in methods_to_try:
            logger.info(f"尝试使用: {method_name}")
            
            try:
                success = self._translate_with_method(
                    method_id, ppt_path, translation_data, output_path, bilingual_mode
                )
                
                if success:
                    logger.info(f"✅ 翻译成功，使用方法: {method_name}")
                    return True, method_id
                else:
                    logger.warning(f"⚠️ 方法失败: {method_name}")
                    
            except Exception as e:
                logger.warning(f"⚠️ 方法异常: {method_name} - {e}")
                continue
        
        logger.error("❌ 所有翻译方法都失败了")
        return False, "all_failed"
    
    def _translate_with_method(
        self,
        method_id: str,
        ppt_path: str,
        translation_data: Dict[str, str],
        output_path: str,
        bilingual_mode: bool
    ) -> bool:
        """使用指定方法翻译"""
        
        if method_id == 'uno_direct':
            return self._translate_uno_direct(ppt_path, translation_data, output_path, bilingual_mode)
        
        elif method_id == 'uno_external':
            return self._translate_uno_external(ppt_path, translation_data, output_path, bilingual_mode)
        
        elif method_id == 'enhanced_pptx':
            return self._translate_enhanced_pptx(ppt_path, translation_data, output_path, bilingual_mode)
        
        elif method_id == 'page_based':
            return self._translate_page_based(ppt_path, translation_data, output_path, bilingual_mode)
        
        elif method_id == 'basic_pptx':
            return self._translate_basic_pptx(ppt_path, translation_data, output_path, bilingual_mode)
        
        else:
            logger.error(f"未知的翻译方法: {method_id}")
            return False
    
    def _translate_uno_direct(self, ppt_path: str, translation_data: Dict[str, str], output_path: str, bilingual_mode: bool) -> bool:
        """直接UNO接口翻译"""
        try:
            from .ppt_translate_uno import translate_ppt_with_uno_color_protection
            return translate_ppt_with_uno_color_protection(
                ppt_path, translation_data, output_path, bilingual_mode
            )
        except Exception as e:
            logger.debug(f"直接UNO翻译失败: {e}")
            return False
    
    def _translate_uno_external(self, ppt_path: str, translation_data: Dict[str, str], output_path: str, bilingual_mode: bool) -> bool:
        """外部UNO处理器翻译"""
        try:
            from .libreoffice_uno_alternative import translate_ppt_external_uno
            
            # 外部处理器目前不支持双语模式，需要预处理
            if bilingual_mode:
                processed_data = {}
                for original, translated in translation_data.items():
                    if original != translated:
                        processed_data[original] = f"{original}\n{translated}"
                    else:
                        processed_data[original] = original
                translation_data = processed_data
            
            return translate_ppt_external_uno(ppt_path, translation_data, output_path)
        except Exception as e:
            logger.debug(f"外部UNO翻译失败: {e}")
            return False
    
    def _translate_enhanced_pptx(self, ppt_path: str, translation_data: Dict[str, str], output_path: str, bilingual_mode: bool) -> bool:
        """增强python-pptx翻译"""
        try:
            from .color_protection import ColorProtector
            from pptx import Presentation
            
            prs = Presentation(ppt_path)
            protector = ColorProtector()
            
            # 处理每个幻灯片
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_frame = shape.text_frame
                        
                        # 保存颜色
                        color_info = protector.save_text_colors(text_frame)
                        
                        # 应用翻译
                        for paragraph in text_frame.paragraphs:
                            original_text = paragraph.text.strip()
                            if original_text in translation_data:
                                translated = translation_data[original_text]
                                if bilingual_mode:
                                    paragraph.text = f"{original_text}\n{translated}"
                                else:
                                    paragraph.text = translated
                        
                        # 恢复颜色
                        protector.restore_text_colors(text_frame, color_info)
            
            # 保存
            save_path = output_path or ppt_path
            prs.save(save_path)
            return True
            
        except Exception as e:
            logger.debug(f"增强python-pptx翻译失败: {e}")
            return False
    
    def _translate_page_based(self, ppt_path: str, translation_data: Dict[str, str], output_path: str, bilingual_mode: bool) -> bool:
        """页面翻译"""
        try:
            from .page_based_translation import translate_ppt_by_pages
            
            # 转换翻译数据格式
            translations_list = list(translation_data.values())
            
            return translate_ppt_by_pages(
                ppt_path, translations_list, str(int(bilingual_mode)), output_path
            )
        except Exception as e:
            logger.debug(f"页面翻译失败: {e}")
            return False
    
    def _translate_basic_pptx(self, ppt_path: str, translation_data: Dict[str, str], output_path: str, bilingual_mode: bool) -> bool:
        """基础python-pptx翻译"""
        try:
            from pptx import Presentation
            
            prs = Presentation(ppt_path)
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            original_text = paragraph.text.strip()
                            if original_text in translation_data:
                                translated = translation_data[original_text]
                                if bilingual_mode:
                                    paragraph.text = f"{original_text}\n{translated}"
                                else:
                                    paragraph.text = translated
            
            save_path = output_path or ppt_path
            prs.save(save_path)
            return True
            
        except Exception as e:
            logger.debug(f"基础python-pptx翻译失败: {e}")
            return False


# 全局智能翻译器实例
_smart_translator = SmartColorTranslator()


def smart_translate_ppt(
    ppt_path: str,
    translation_data: Dict[str, str],
    output_path: str = None,
    bilingual_mode: bool = False,
    preferred_method: str = None
) -> Tuple[bool, str]:
    """
    智能PPT颜色保护翻译（便捷函数）
    
    Args:
        ppt_path: PPT文件路径
        translation_data: 翻译数据
        output_path: 输出路径
        bilingual_mode: 双语模式
        preferred_method: 首选方法
        
    Returns:
        Tuple[bool, str]: (是否成功, 使用的方法)
    """
    return _smart_translator.translate_ppt_with_best_color_protection(
        ppt_path, translation_data, output_path, bilingual_mode, preferred_method
    )


def get_translation_capabilities() -> Dict[str, Any]:
    """获取翻译能力信息"""
    methods = _smart_translator.get_available_methods()
    
    return {
        'available_methods': [
            {'id': method_id, 'name': method_name, 'priority': priority}
            for method_id, method_name, priority in methods
        ],
        'preferred_method': _smart_translator.preferred_method,
        'total_methods': len(methods),
        'best_available': methods[0] if methods else None
    }


if __name__ == "__main__":
    print("智能颜色保护翻译模块")
    print("=" * 50)
    
    capabilities = get_translation_capabilities()
    print(f"可用方法数量: {capabilities['total_methods']}")
    print(f"首选方法: {capabilities['preferred_method']}")
    
    print("\n可用的翻译方法:")
    for method in capabilities['available_methods']:
        print(f"  • {method['name']} (优先级: {method['priority']})")
    
    if capabilities['best_available']:
        print(f"\n✅ 推荐使用: {capabilities['best_available'][1]}")
    else:
        print("\n❌ 没有可用的翻译方法")
