import os
import json
import string
import tempfile
import shutil
from typing import List, Dict, Optional, Tuple, Any
from pathlib import Path
from pptx import Presentation
from pptx.shapes.picture import Picture
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys
import os
import re
sys.path.insert(0, os.path.dirname(__file__))

# æ·»åŠ QwenTranslatorå¯¼å…¥
from .translator import QwenTranslator

def perform_ocr_on_image(image_path: str, api_key: str) -> Optional[Dict]:
    """
    å¯¹å•ä¸ªå›¾ç‰‡æ‰§è¡ŒOCRè¯†åˆ«
    
    Args:
        image_path: å›¾ç‰‡è·¯å¾„
        api_key: OCR APIå¯†é’¥
        
    Returns:
        OCRç»“æœå­—å…¸æˆ–None
    """
    logger.info(f"[OCRè¯†åˆ«] å¼€å§‹å¯¹å›¾ç‰‡æ‰§è¡ŒOCRè¯†åˆ«: {image_path}")
    logger.info(f"[OCRè¯†åˆ«] å›¾ç‰‡æ–‡ä»¶å¤§å°: {os.path.getsize(image_path) if os.path.exists(image_path) else 'æ–‡ä»¶ä¸å­˜åœ¨'} å­—èŠ‚")
    
    try:
        # æ£€æŸ¥æ–‡ä»¶æ˜¯å¦å­˜åœ¨
        if not os.path.exists(image_path):
            logger.error(f"[OCRè¯†åˆ«] å›¾ç‰‡æ–‡ä»¶ä¸å­˜åœ¨: {image_path}")
            return None
            
        # æ£€æŸ¥APIå¯†é’¥
        if not api_key:
            logger.error("[OCRè¯†åˆ«] APIå¯†é’¥æœªæä¾›")
            return None
            
        # è¯»å–å›¾ç‰‡æ–‡ä»¶
        with open(image_path, 'rb') as f:
            image_data = f.read()
            
        logger.info(f"[OCRè¯†åˆ«] æˆåŠŸè¯»å–å›¾ç‰‡æ•°æ®ï¼Œå¤§å°: {len(image_data)} å­—èŠ‚")
        
        # è°ƒç”¨OCR API
        headers = {
            'Authorization': f'Bearer {api_key}',
            'Content-Type': 'application/octet-stream'
        }
        
        # æ ¹æ®æ–‡ä»¶æ‰©å±•åç¡®å®šOCRç±»å‹
        _, ext = os.path.splitext(image_path)
        ext = ext.lower()
        
        if ext in ['.pdf']:
            ocr_type = 'pdf'
        elif ext in ['.png', '.jpg', '.jpeg', '.bmp', '.gif']:
            ocr_type = 'image'
        else:
            ocr_type = 'image'  # é»˜è®¤ä¸ºå›¾ç‰‡
            
        logger.info(f"[OCRè¯†åˆ«] å›¾ç‰‡ç±»å‹: {ext}, OCRç±»å‹: {ocr_type}")
        
        # æ„é€ API URL
        api_url = "https://dashscope.aliyuncs.com/api/v1/services/ocr/general"
        params = {'ocr_type': ocr_type}
        
        logger.info(f"[OCRè¯†åˆ«] è°ƒç”¨OCR API: {api_url}")
        logger.info(f"[OCRè¯†åˆ«] è¯·æ±‚å‚æ•°: {params}")
        
        # å‘é€è¯·æ±‚
        response = requests.post(
            api_url,
            headers=headers,
            params=params,
            data=image_data,
            timeout=300  # 5åˆ†é’Ÿè¶…æ—¶
        )
        
        logger.info(f"[OCRè¯†åˆ«] OCR APIå“åº”çŠ¶æ€ç : {response.status_code}")
        
        # æ£€æŸ¥å“åº”
        if response.status_code == 200:
            result = response.json()
            logger.info(f"[OCRè¯†åˆ«] OCR APIå“åº”æˆåŠŸ")
            logger.info(f"[OCRè¯†åˆ«] OCRç»“æœç»“æ„: {list(result.keys()) if isinstance(result, dict) else 'éå­—å…¸ç»“æ„'}")
            
            # æå–æ–‡æœ¬ä¿¡æ¯
            if 'output' in result and 'recognized_texts' in result['output']:
                texts = result['output']['recognized_texts']
                logger.info(f"[OCRè¯†åˆ«] è¯†åˆ«åˆ° {len(texts)} ä¸ªæ–‡æœ¬å—")
                
                # æ„é€ è¿”å›ç»“æœ
                extracted_texts = {}
                for i, text_info in enumerate(texts):
                    if isinstance(text_info, dict) and 'text' in text_info:
                        text_key = f"text_{i+1:03d}"
                        extracted_texts[text_key] = text_info['text']
                        logger.info(f"[OCRè¯†åˆ«] æ–‡æœ¬ {text_key}: {text_info['text'][:50]}...")
                
                return {
                    'texts': extracted_texts,
                    'raw_result': result
                }
            else:
                logger.warning(f"[OCRè¯†åˆ«] OCRç»“æœæ ¼å¼ä¸ç¬¦åˆé¢„æœŸ: {list(result.keys()) if isinstance(result, dict) else 'éå­—å…¸ç»“æ„'}")
                return {
                    'texts': {},
                    'raw_result': result
                }
        else:
            logger.error(f"[OCRè¯†åˆ«] OCR APIè°ƒç”¨å¤±è´¥ï¼ŒçŠ¶æ€ç : {response.status_code}")
            logger.error(f"[OCRè¯†åˆ«] é”™è¯¯å“åº”: {response.text}")
            return None
            
    except Exception as e:
        logger.error(f"[OCRè¯†åˆ«] OCRè¯†åˆ«è¿‡ç¨‹ä¸­å‡ºé”™: {str(e)}")
        logger.exception("[OCRè¯†åˆ«] è¯¦ç»†é”™è¯¯ä¿¡æ¯")
        return None
# å¯¼å…¥æ—¥å¿—ç³»ç»Ÿ
from logger_config_ocr import get_logger

# å¯¼å…¥OCR APIå¤„ç†å™¨
from ocr_api import OCRProcessor

# å¯¼å…¥OCR QWEN APIå¤„ç†ç¨‹åº
from qwen_ocr_api import process_folder_with_mapping

# å¯¼å…¥ç¿»è¯‘æ¨¡å—
from translator import TranslationManager

# è·å–æ—¥å¿—è®°å½•å™¨
logger = get_logger("ocr_controller")


def _extract_image_paths_from_markdown(markdown_content: str, markdown_dir: str) -> List[str]:
    """
    ä»Markdownå†…å®¹ä¸­æå–å›¾ç‰‡è·¯å¾„
    
    Args:
        markdown_content: Markdownå†…å®¹
        markdown_dir: Markdownæ–‡ä»¶æ‰€åœ¨ç›®å½•
        
    Returns:
        å›¾ç‰‡è·¯å¾„åˆ—è¡¨
    """
    logger.info("[å›¾ç‰‡æå–] å¼€å§‹ä»Markdownå†…å®¹ä¸­æå–å›¾ç‰‡è·¯å¾„")
    logger.info(f"[å›¾ç‰‡æå–] Markdownç›®å½•: {markdown_dir}")
    logger.info(f"[å›¾ç‰‡æå–] Markdownå†…å®¹é•¿åº¦: {len(markdown_content)}")
    
    # åŒ¹é…Markdownå›¾ç‰‡è¯­æ³•: ![alt](path)
    image_pattern = r'!\[.*?\]\((.*?)\)'
    matches = re.findall(image_pattern, markdown_content)
    
    logger.info(f"[å›¾ç‰‡æå–] æ‰¾åˆ° {len(matches)} ä¸ªå›¾ç‰‡å¼•ç”¨")
    for i, match in enumerate(matches):
        logger.info(f"[å›¾ç‰‡æå–] å›¾ç‰‡å¼•ç”¨ {i+1}: {match}")
    
    image_paths = []
    for match in matches:
        # å¤„ç†ç›¸å¯¹è·¯å¾„
        if match.startswith('images/'):
            # ç›¸å¯¹äºmarkdownæ–‡ä»¶çš„imagesç›®å½•
            image_path = os.path.join(markdown_dir, match)
        elif match.startswith('./images/'):
            # ç›¸å¯¹äºmarkdownæ–‡ä»¶çš„imagesç›®å½•
            image_path = os.path.join(markdown_dir, match[2:])  # å»æ‰ ./
        elif not os.path.isabs(match):
            # å…¶ä»–ç›¸å¯¹è·¯å¾„
            image_path = os.path.join(markdown_dir, match)
        else:
            # ç»å¯¹è·¯å¾„
            image_path = match
            
        logger.info(f"[å›¾ç‰‡æå–] å¤„ç†å›¾ç‰‡è·¯å¾„: {match} -> {image_path}")
        
        # æ£€æŸ¥æ–‡ä»¶æ˜¯å¦å­˜åœ¨
        if os.path.exists(image_path):
            image_paths.append(image_path)
            logger.info(f"[å›¾ç‰‡æå–] å›¾ç‰‡æ–‡ä»¶å­˜åœ¨: {image_path}")
        else:
            logger.warning(f"[å›¾ç‰‡æå–] å›¾ç‰‡æ–‡ä»¶ä¸å­˜åœ¨: {image_path}")
    
    logger.info(f"[å›¾ç‰‡æå–] æœ€ç»ˆæå–åˆ° {len(image_paths)} ä¸ªæœ‰æ•ˆå›¾ç‰‡è·¯å¾„")
    for i, path in enumerate(image_paths):
        logger.info(f"[å›¾ç‰‡æå–] æœ‰æ•ˆå›¾ç‰‡ {i+1}: {path}")
        
    return image_paths


def _create_image_mapping(image_paths: List[str]) -> Dict:
    """
    åˆ›å»ºå›¾ç‰‡æ˜ å°„æ•°æ®ç»“æ„
    
    Args:
        image_paths: å›¾ç‰‡è·¯å¾„åˆ—è¡¨
        
    Returns:
        å›¾ç‰‡æ˜ å°„å­—å…¸
    """
    mapping = {"slide_1": {"slide_number": 1, "images": []}}
    
    for i, image_path in enumerate(image_paths):
        filename = os.path.basename(image_path)
        image_info = {
            "filename": filename,
            "filepath": image_path,
            "shape_id": i,
            "left": 0,
            "top": 0,
            "width": 0,
            "height": 0,
            "content_type": "",
            "original_format": "",
            "file_size": os.path.getsize(image_path) if os.path.exists(image_path) else 0
        }
        mapping["slide_1"]["images"].append(image_info)
    
    return mapping


def _translate_text(text: str, target_language: str, source_language: str) -> str:
    """
    ç¿»è¯‘æ–‡æœ¬
    
    Args:
        text: å¾…ç¿»è¯‘æ–‡æœ¬
        target_language: ç›®æ ‡è¯­è¨€
        source_language: æºè¯­è¨€
        
    Returns:
        ç¿»è¯‘åçš„æ–‡æœ¬ï¼Œå¦‚æœç¿»è¯‘å¤±è´¥è¿”å›None
    """
    try:
        # ä½¿ç”¨QwenTranslatorè¿›è¡Œç¿»è¯‘
        translator = QwenTranslator(target_language=target_language)
        translated = translator.translate_text(text, source_language=source_language)
        # åªæœ‰å½“ç¿»è¯‘ç»“æœä¸åŸæ–‡ä¸åŒæ—¶æ‰è¿”å›ç¿»è¯‘ç»“æœï¼Œå¦åˆ™è¿”å›Noneè¡¨ç¤ºç¿»è¯‘å¤±è´¥
        if translated and translated.strip() and translated.strip() != text.strip():
            return translated
        else:
            return None
    except Exception as e:
        logger.error(f"ç¿»è¯‘æ–‡æœ¬æ—¶å‡ºé”™: {str(e)}")
        logger.exception("ç¿»è¯‘é”™è¯¯è¯¦æƒ…")
        return None


def process_markdown_images_ocr_and_translate(
    markdown_content: str,
    markdown_dir: str,
    target_language: str = "zh",
    source_language: str = "en"
) -> List[Dict]:
    """
    å¤„ç†Markdownä¸­çš„å›¾ç‰‡OCRè¯†åˆ«å’Œç¿»è¯‘
    
    Args:
        markdown_content: Markdownå†…å®¹
        markdown_dir: Markdownæ–‡ä»¶æ‰€åœ¨ç›®å½•
        target_language: ç›®æ ‡è¯­è¨€
        source_language: æºè¯­è¨€
        
    Returns:
        List of OCRç»“æœå­—å…¸ï¼Œæ¯ä¸ªå­—å…¸åŒ…å«:
        {
            "success": bool,
            "image_path": str,
            "ocr_text_combined": str,
            "translation_text_combined": str
        }
    """
    logger.info("[Markdown OCR] å¼€å§‹å¤„ç†Markdownä¸­çš„å›¾ç‰‡OCRè¯†åˆ«å’Œç¿»è¯‘")
    logger.info(f"[Markdown OCR] Markdownç›®å½•: {markdown_dir}")
    logger.info(f"[Markdown OCR] ç›®æ ‡è¯­è¨€: {target_language}, æºæ ‡è¯­è¨€: {source_language}")
    logger.info(f"[Markdown OCR] Markdownå†…å®¹é•¿åº¦: {len(markdown_content)} å­—ç¬¦")
    
    results = []
    
    try:
        # åˆ›å»ºä¸´æ—¶ç›®å½•ç”¨äºå¤„ç†
        with tempfile.TemporaryDirectory(prefix="md_ocr_") as temp_dir:
            logger.info(f"[Markdown OCR] åˆ›å»ºä¸´æ—¶ç›®å½•: {temp_dir}")
            
            # æå–Markdownä¸­çš„å›¾ç‰‡é“¾æ¥
            image_paths = _extract_image_paths_from_markdown(markdown_content, markdown_dir)
            logger.info(f"[Markdown OCR] æ‰¾åˆ° {len(image_paths)} ä¸ªå›¾ç‰‡è·¯å¾„")
            
            if not image_paths:
                logger.info("[Markdown OCR] æœªæ‰¾åˆ°ä»»ä½•å›¾ç‰‡ï¼Œè·³è¿‡OCRå¤„ç†")
                return results
            
            # å¤åˆ¶å›¾ç‰‡åˆ°ä¸´æ—¶ç›®å½•
            temp_image_paths = []
            for i, image_path in enumerate(image_paths):
                if os.path.exists(image_path):
                    # ç”Ÿæˆæ–°çš„æ–‡ä»¶å
                    ext = os.path.splitext(image_path)[1]
                    new_filename = f"image_{i+1:04d}{ext}"
                    temp_image_path = os.path.join(temp_dir, new_filename)
                    
                    # å¤åˆ¶æ–‡ä»¶
                    shutil.copy2(image_path, temp_image_path)
                    temp_image_paths.append(temp_image_path)
                    logger.info(f"[Markdown OCR] å·²å¤åˆ¶å›¾ç‰‡: {os.path.basename(image_path)} -> {new_filename}")
                else:
                    logger.warning(f"[Markdown OCR] å›¾ç‰‡æ–‡ä»¶ä¸å­˜åœ¨: {image_path}")
            
            if not temp_image_paths:
                logger.warning("[Markdown OCR] æ²¡æœ‰æœ‰æ•ˆçš„å›¾ç‰‡æ–‡ä»¶ï¼Œè·³è¿‡OCRå¤„ç†")
                return results
            
            logger.info(f"[Markdown OCR] æˆåŠŸå¤åˆ¶ {len(temp_image_paths)} ä¸ªå›¾ç‰‡æ–‡ä»¶åˆ°ä¸´æ—¶ç›®å½•")
            
            # åˆ›å»ºå›¾ç‰‡æ˜ å°„JSONæ–‡ä»¶
            image_mapping = _create_image_mapping(temp_image_paths)
            json_path = os.path.join(temp_dir, "image_mapping.json")
            with open(json_path, 'w', encoding='utf-8') as f:
                json.dump(image_mapping, f, ensure_ascii=False, indent=2)
            logger.info(f"[Markdown OCR] å·²åˆ›å»ºå›¾ç‰‡æ˜ å°„æ–‡ä»¶: {json_path}")
            logger.info(f"[Markdown OCR] å›¾ç‰‡æ˜ å°„å†…å®¹: {json.dumps(image_mapping, ensure_ascii=False, indent=2)}")
            
            # æ‰§è¡ŒOCRè¯†åˆ«
            API_KEY = os.getenv("QWEN_API_KEY")
            if not API_KEY:
                logger.error("[Markdown OCR] æœªè®¾ç½®QWEN_API_KEYç¯å¢ƒå˜é‡ï¼Œæ— æ³•æ‰§è¡ŒOCRè¯†åˆ«")
                return results
            
            logger.info("[Markdown OCR] å¼€å§‹æ‰§è¡ŒOCRè¯†åˆ«")
            process_folder_with_mapping(temp_dir, json_path, API_KEY)
            logger.info("[Markdown OCR] OCRè¯†åˆ«å®Œæˆ")
            
            # è¯»å–OCRç»“æœ
            if os.path.exists(json_path):
                with open(json_path, 'r', encoding='utf-8') as f:
                    updated_mapping = json.load(f)
                
                logger.info(f"[Markdown OCR] OCRç»“æœæ–‡ä»¶åŠ è½½æˆåŠŸ")
                logger.info(f"[Markdown OCR] OCRç»“æœå†…å®¹: {json.dumps(updated_mapping, ensure_ascii=False, indent=2)}")
                
                # å¤„ç†æ¯ä¸ªå›¾ç‰‡çš„OCRç»“æœ
                for slide_key, slide_data in updated_mapping.items():
                    if 'images' not in slide_data:
                        logger.warning(f"[Markdown OCR] Slideæ•°æ®ä¸­æ²¡æœ‰imageså­—æ®µ: {slide_key}")
                        continue
                    
                    for image_info in slide_data['images']:
                        result = {
                            "image_path": image_info.get("filepath", ""),
                            "success": False,
                            "ocr_text_combined": "",
                            "translation_text_combined": ""
                        }
                        
                        # è·å–OCRæ–‡æœ¬
                        all_text = image_info.get('all_text', {})
                        logger.info(f"[Markdown OCR] å›¾ç‰‡ {os.path.basename(result['image_path'])} çš„OCRæ–‡æœ¬: {all_text}")
                        
                        if all_text:
                            # åˆå¹¶æ‰€æœ‰OCRæ–‡æœ¬
                            ocr_texts = [text.strip() for text in all_text.values() if text.strip()]
                            result["ocr_text_combined"] = '\n'.join(ocr_texts)
                            result["success"] = True
                            logger.info(f"[Markdown OCR] å›¾ç‰‡OCRè¯†åˆ«æˆåŠŸ: {os.path.basename(result['image_path'])}, æ–‡æœ¬é•¿åº¦: {len(result['ocr_text_combined'])}")
                        else:
                            logger.info(f"[Markdown OCR] å›¾ç‰‡æœªè¯†åˆ«åˆ°æ–‡æœ¬: {os.path.basename(result['image_path'])}")
                        
                        # ç¿»è¯‘OCRæ–‡æœ¬
                        if all_text:
                            try:
                                # ç¿»è¯‘æ‰€æœ‰æ–‡æœ¬
                                translated_texts = {}
                                for text_key, text_value in all_text.items():
                                    if text_value and text_value.strip():
                                        logger.info(f"[Markdown OCR] ç¿»è¯‘æ–‡æœ¬: {text_key} = {text_value.strip()}")
                                        translated = _translate_text(text_value.strip(), target_language, source_language)
                                        # æ£€æŸ¥ç¿»è¯‘æ˜¯å¦æˆåŠŸ
                                        if translated is not None:
                                            translated_texts[text_key] = translated
                                            logger.info(f"[Markdown OCR] ç¿»è¯‘ç»“æœ: {text_key} = {translated}")
                                        else:
                                            logger.warning(f"[Markdown OCR] ç¿»è¯‘å¤±è´¥æˆ–ç»“æœä¸åŸæ–‡ç›¸åŒï¼Œè·³è¿‡è¯¥æ–‡æœ¬æ®µ: {text_key}")
                                
                                # åˆå¹¶æ‰€æœ‰ç¿»è¯‘æ–‡æœ¬ï¼ˆåŒ…æ‹¬æœªç¿»è¯‘çš„åŸæ–‡ï¼‰
                                translation_texts = []
                                for text_key, text_value in all_text.items():
                                    if text_key in translated_texts:
                                        translation_texts.append(translated_texts[text_key].strip())
                                    elif text_value and text_value.strip():
                                        # å¦‚æœæ²¡æœ‰ç¿»è¯‘æˆåŠŸï¼Œåˆ™ä½¿ç”¨åŸæ–‡
                                        translation_texts.append(text_value.strip())
                                
                                if translation_texts:
                                    result["translation_text_combined"] = '\n'.join(translation_texts)
                                    logger.info(f"[Markdown OCR] å›¾ç‰‡ç¿»è¯‘æˆåŠŸ: {os.path.basename(result['image_path'])}, ç¿»è¯‘æ–‡æœ¬é•¿åº¦: {len(result['translation_text_combined'])}")
                                
                            except Exception as e:
                                logger.error(f"[Markdown OCR] ç¿»è¯‘å›¾ç‰‡æ–‡æœ¬æ—¶å‡ºé”™: {e}")
                                logger.exception("[Markdown OCR] ç¿»è¯‘é”™è¯¯è¯¦æƒ…")
                        
                        results.append(result)
                
                logger.info(f"[Markdown OCR] å¤„ç†å®Œæˆï¼Œå…±å¤„ç† {len(results)} ä¸ªå›¾ç‰‡")
                for i, result in enumerate(results):
                    logger.info(f"[Markdown OCR] ç»“æœ {i+1}: æˆåŠŸ={result['success']}, "
                              f"å›¾ç‰‡={os.path.basename(result['image_path'])}, "
                              f"OCRæ–‡æœ¬é•¿åº¦={len(result['ocr_text_combined'])}, "
                              f"ç¿»è¯‘æ–‡æœ¬é•¿åº¦={len(result['translation_text_combined'])}")
            else:
                logger.error(f"[Markdown OCR] OCRç»“æœæ–‡ä»¶ä¸å­˜åœ¨: {json_path}")
            
            return results
            
    except Exception as e:
        logger.error(f"[Markdown OCR] å¤„ç†Markdownå›¾ç‰‡OCRå’Œç¿»è¯‘æ—¶å‡ºé”™: {str(e)}")
        logger.exception("[Markdown OCR] è¯¦ç»†é”™è¯¯ä¿¡æ¯")
        return results


class TextLineSplitter:
    """JSONæ–‡æœ¬è¡Œåˆ†å‰²å™¨ - å°†åŒ…å«æ¢è¡Œç¬¦çš„æ–‡æœ¬æ‹†åˆ†æˆå¤šä¸ªtextå­—æ®µ"""
    
    def __init__(self):
        self.processed_count = 0
        self.split_count = 0
    
    def process_json_file(self, json_file_path: str) -> bool:
        """
        å¤„ç†JSONæ–‡ä»¶ä¸­çš„æ–‡æœ¬è¡Œåˆ†å‰²
        
        Args:
            json_file_path: JSONæ–‡ä»¶è·¯å¾„
            
        Returns:
            æ˜¯å¦å¤„ç†æˆåŠŸ
        """
        try:
            # æ£€æŸ¥æ–‡ä»¶æ˜¯å¦å­˜åœ¨
            if not os.path.exists(json_file_path):
                logger.error(f"JSONæ–‡ä»¶ä¸å­˜åœ¨: {json_file_path}")
                return False
            
            # è¯»å–JSONæ–‡ä»¶
            with open(json_file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # å¤„ç†æ•°æ®
            self._process_mapping_data(data)
            
            # ä¿å­˜å¤„ç†åçš„æ•°æ®
            with open(json_file_path, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
            
            logger.info(f"âœ… æ–‡æœ¬è¡Œåˆ†å‰²å®Œæˆï¼å¤„ç†å›¾ç‰‡æ•°: {self.processed_count}, åˆ†å‰²æ–‡æœ¬æ•°: {self.split_count}")
            
            return True
            
        except Exception as e:
            logger.error(f"å¤„ç†JSONæ–‡ä»¶æ—¶å‡ºé”™: {str(e)}")
            return False
    
    def _process_mapping_data(self, data: Dict[str, Any]):
        """å¤„ç†æ˜ å°„æ•°æ®ä¸­çš„æ‰€æœ‰æ–‡æœ¬"""
        for slide_key, slide_data in data.items():
            if 'images' not in slide_data:
                continue
            
            for image_info in slide_data['images']:
                filename = image_info.get('filename', 'æœªçŸ¥æ–‡ä»¶')
                
                # å¤„ç†all_textå­—æ®µ
                if 'all_text' in image_info and image_info['all_text']:
                    original_all_text = image_info['all_text'].copy()
                    new_all_text = self._split_text_dict(original_all_text)
                    
                    if new_all_text != original_all_text:
                        image_info['all_text'] = new_all_text
                        self.split_count += 1
                        logger.info(f"ğŸ”„ å·²åˆ†å‰² {filename} çš„æ–‡æœ¬è¡Œ")
                
                self.processed_count += 1
    
    def _split_text_dict(self, text_dict: Dict[str, str]) -> Dict[str, str]:
        """
        åˆ†å‰²æ–‡æœ¬å­—å…¸ä¸­çš„æ–‡æœ¬è¡Œ
        
        Args:
            text_dict: åŒ…å«æ–‡æœ¬çš„å­—å…¸
            
        Returns:
            åˆ†å‰²åçš„æ–‡æœ¬å­—å…¸
        """
        new_dict = {}
        text_counter = 1
        
        for key, text_value in text_dict.items():
            if not text_value:
                # ç©ºæ–‡æœ¬ç›´æ¥è·³è¿‡
                continue
            
            # æ£€æŸ¥æ˜¯å¦åŒ…å«æ¢è¡Œç¬¦
            if '\n' in text_value:
                # æŒ‰æ¢è¡Œç¬¦åˆ†å‰²
                lines = text_value.split('\n')
                
                logger.info(f"ğŸ“ å‘ç°å¤šè¡Œæ–‡æœ¬ {key}: {len(lines)} è¡Œ")
                
                # ä¸ºæ¯ä¸€è¡Œåˆ›å»ºæ–°çš„textå­—æ®µ
                for line in lines:
                    new_key = f"text{text_counter}"
                    new_dict[new_key] = line  # ä¿ç•™ç©ºè¡Œï¼Œä¸è¿›è¡Œstrip()
                    text_counter += 1
                    
            else:
                # å•è¡Œæ–‡æœ¬ï¼Œç›´æ¥ä½¿ç”¨æ–°çš„ç¼–å·
                new_key = f"text{text_counter}"
                new_dict[new_key] = text_value
                text_counter += 1
        
        return new_dict


class PPTImageExtractor:
    """PPTå›¾ç‰‡æå–å™¨"""
    
    def __init__(self, temp_dir: str = None):
        self.temp_dir = temp_dir or tempfile.mkdtemp(prefix="ppt_ocr_")
        self.image_mapping = {}
        self.image_counter = 0
    
    def extract_images_from_slides(self, presentation_path: str, 
                                 selected_pages: Optional[List[int]] = None) -> Tuple[str, Dict]:
        """
        ä»PPTæŒ‡å®šé¡µé¢æå–å›¾ç‰‡
        
        Args:
            presentation_path: PPTæ–‡ä»¶è·¯å¾„
            selected_pages: é€‰æ‹©çš„é¡µé¢åˆ—è¡¨ï¼ŒNoneè¡¨ç¤ºå…¨é€‰
            
        Returns:
            Tuple[ä¸´æ—¶æ–‡ä»¶å¤¹è·¯å¾„, å›¾ç‰‡æ˜ å°„å­—å…¸]
        """
        try:
            prs = Presentation(presentation_path)
            total_slides = len(prs.slides)
            
            # ç¡®å®šè¦å¤„ç†çš„é¡µé¢
            if selected_pages is None:
                pages_to_process = list(range(total_slides))
            else:
                pages_to_process = [p for p in selected_pages if 0 <= p < total_slides]
            
            logger.info(f"æ­£åœ¨å¤„ç† {len(pages_to_process)} ä¸ªé¡µé¢çš„å›¾ç‰‡...")
            
            for slide_idx in pages_to_process:
                slide = prs.slides[slide_idx]
                self._extract_images_from_slide(slide, slide_idx)
            
            # ä¿å­˜æ˜ å°„å…³ç³»åˆ°JSONæ–‡ä»¶
            mapping_file = os.path.join(self.temp_dir, "image_mapping.json")
            with open(mapping_file, 'w', encoding='utf-8') as f:
                json.dump(self.image_mapping, f, ensure_ascii=False, indent=2)
            
            logger.info(f"æå–å®Œæˆï¼Œå…±æå– {self.image_counter} å¼ å›¾ç‰‡")
            logger.info(f"ä¸´æ—¶æ–‡ä»¶å¤¹: {self.temp_dir}")
            
            return self.temp_dir, self.image_mapping
            
        except Exception as e:
            logger.error(f"æå–å›¾ç‰‡æ—¶å‘ç”Ÿé”™è¯¯: {str(e)}")
            raise
    
    def _extract_images_from_slide(self, slide, slide_idx: int):
        """ä»å•ä¸ªå¹»ç¯ç‰‡æå–å›¾ç‰‡"""
        slide_images = []
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_info = self._save_image_from_shape(shape, slide_idx)
                    if image_info:
                        slide_images.append(image_info)
                except Exception as e:
                    logger.warning(f"æå–ç¬¬{slide_idx+1}é¡µå›¾ç‰‡æ—¶å‡ºé”™: {str(e)}")
        
        if slide_images:
            self.image_mapping[f"slide_{slide_idx}"] = {
                "slide_number": slide_idx + 1,
                "images": slide_images
            }
    
    def _save_image_from_shape(self, shape: Picture, slide_idx: int) -> Optional[Dict]:
        """ä»å›¾ç‰‡å½¢çŠ¶ä¿å­˜å›¾ç‰‡æ–‡ä»¶"""
        try:
            image = shape.image
            image_bytes = image.blob
            
            if not image_bytes:
                logger.warning(f"ç¬¬{slide_idx+1}é¡µå›¾ç‰‡æ•°æ®ä¸ºç©ºï¼Œè·³è¿‡")
                return None
            
            # æ£€æµ‹å›¾ç‰‡çš„çœŸå®æ ¼å¼
            actual_format, content_type = self._detect_image_format(image_bytes)
            
            # æ ¹æ®çœŸå®æ ¼å¼ç¡®å®šæ–‡ä»¶æ‰©å±•å
            ext_map = {
                'image/jpeg': '.jpg',
                'image/png': '.png', 
                'image/gif': '.gif',
                'image/bmp': '.bmp',
                'image/tiff': '.tiff',
                'image/webp': '.webp',
                'image/x-emf': '.emf',
                'image/x-wmf': '.wmf',
                'unknown': '.bin'  # æœªçŸ¥æ ¼å¼ä¿å­˜ä¸ºäºŒè¿›åˆ¶æ–‡ä»¶
            }
            
            ext = ext_map.get(content_type, '.bin')
            
            # ç”Ÿæˆå”¯ä¸€çš„æ–‡ä»¶å
            self.image_counter += 1
            filename = f"image_{self.image_counter:04d}{ext}"
            filepath = os.path.join(self.temp_dir, filename)
            
            # å¦‚æœæ˜¯ç‰¹æ®Šæ ¼å¼ï¼Œå°è¯•è½¬æ¢ä¸ºPNG
            if content_type in ['image/x-emf', 'image/x-wmf', 'unknown']:
                converted_filepath = self._try_convert_to_png(image_bytes, filename, filepath)
                if converted_filepath:
                    filename = os.path.basename(converted_filepath)
                    filepath = converted_filepath
                    content_type = 'image/png'
                    ext = '.png'
                else:
                    # ä¿å­˜åŸå§‹å›¾ç‰‡æ–‡ä»¶
                    with open(filepath, 'wb') as f:
                        f.write(image_bytes)
            else:
                # ä¿å­˜åŸå§‹å›¾ç‰‡æ–‡ä»¶
                with open(filepath, 'wb') as f:
                    f.write(image_bytes)
            
            # éªŒè¯ä¿å­˜çš„æ–‡ä»¶æ˜¯å¦æœ‰æ•ˆ
            if not self._validate_image_file(filepath):
                logger.warning(f"è­¦å‘Š: å›¾ç‰‡æ–‡ä»¶å¯èƒ½æŸå - {filename}")
                # ä»ç„¶è¿”å›ä¿¡æ¯ï¼Œä½†æ ‡è®°ä¸ºå¯èƒ½æœ‰é—®é¢˜
            
            logger.info(f"å·²ä¿å­˜å›¾ç‰‡: {filename} (æ ¼å¼: {actual_format}, å¤§å°: {len(image_bytes)} bytes)")
            
            # è·å–å›¾ç‰‡åœ¨å¹»ç¯ç‰‡ä¸­çš„ä½ç½®å’Œå¤§å°ä¿¡æ¯
            return {
                "filename": filename,
                "filepath": filepath,
                "shape_id": shape.shape_id,
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
                "content_type": content_type,
                "original_format": actual_format,
                "file_size": len(image_bytes)
            }
            
        except Exception as e:
            logger.error(f"ä¿å­˜å›¾ç‰‡æ—¶å‡ºé”™: {str(e)}")
            return None

    def _detect_image_format(self, image_bytes: bytes) -> tuple:
        """æ£€æµ‹å›¾ç‰‡çš„çœŸå®æ ¼å¼"""
        try:
            # æ£€æŸ¥æ–‡ä»¶å¤´é­”æœ¯å­—èŠ‚
            if len(image_bytes) < 8:
                return "unknown", "unknown"
            
            header = image_bytes[:8]
            
            # PNG: 89 50 4E 47 0D 0A 1A 0A
            if header.startswith(b'\x89PNG\r\n\x1a\n'):
                return "PNG", "image/png"
            
            # JPEG: FF D8 FF
            elif header.startswith(b'\xff\xd8\xff'):
                return "JPEG", "image/jpeg"
            
            # GIF: 47 49 46 38 (GIF8)
            elif header.startswith(b'GIF8'):
                return "GIF", "image/gif"
            
            # BMP: 42 4D (BM)
            elif header.startswith(b'BM'):
                return "BMP", "image/bmp"
            
            # TIFF: 4D 4D æˆ– 49 49
            elif header.startswith(b'MM') or header.startswith(b'II'):
                return "TIFF", "image/tiff"
            
            # WebP: 52 49 46 46 ... 57 45 42 50
            elif header.startswith(b'RIFF') and b'WEBP' in image_bytes[:12]:
                return "WebP", "image/webp"
            
            # EMF: 01 00 00 00 ... (Enhanced Metafile)
            elif len(image_bytes) > 40 and image_bytes[40:44] == b' EMF':
                return "EMF", "image/x-emf"
            
            # WMF: D7 CD C6 9A æˆ–å…¶ä»–WMFæ ‡è¯†
            elif header.startswith(b'\xd7\xcd\xc6\x9a') or header.startswith(b'\x01\x00\x09\x00'):
                return "WMF", "image/x-wmf"
            
            else:
                return "unknown", "unknown"
                
        except Exception as e:
            logger.error(f"æ£€æµ‹å›¾ç‰‡æ ¼å¼æ—¶å‡ºé”™: {str(e)}")
            return "unknown", "unknown"

    def _try_convert_to_png(self, image_bytes: bytes, original_filename: str, original_filepath: str) -> Optional[str]:
        """å°è¯•å°†ç‰¹æ®Šæ ¼å¼è½¬æ¢ä¸ºPNG"""
        try:
            # å°è¯•ä½¿ç”¨PILè½¬æ¢
            try:
                from PIL import Image
                import io
                
                # å…ˆä¿å­˜åŸå§‹æ–‡ä»¶
                with open(original_filepath, 'wb') as f:
                    f.write(image_bytes)
                
                # å°è¯•ç”¨PILæ‰“å¼€å¹¶è½¬æ¢
                img = Image.open(io.BytesIO(image_bytes))
                
                # è½¬æ¢ä¸ºPNG
                png_filename = original_filename.replace(os.path.splitext(original_filename)[1], '_converted.png')
                png_filepath = os.path.join(os.path.dirname(original_filepath), png_filename)
                
                img.save(png_filepath, 'PNG')
                
                logger.info(f"å·²è½¬æ¢ç‰¹æ®Šæ ¼å¼å›¾ç‰‡ä¸ºPNG: {png_filename}")
                return png_filepath
                
            except ImportError:
                logger.warning("è­¦å‘Š: æœªå®‰è£…PIL/Pillowï¼Œæ— æ³•è½¬æ¢ç‰¹æ®Šæ ¼å¼å›¾ç‰‡")
                return None
            except Exception as e:
                logger.error(f"PILè½¬æ¢å¤±è´¥: {str(e)}")
                return None
                
        except Exception as e:
            logger.error(f"è½¬æ¢å›¾ç‰‡æ ¼å¼æ—¶å‡ºé”™: {str(e)}")
            return None

    def _validate_image_file(self, filepath: str) -> bool:
        """éªŒè¯å›¾ç‰‡æ–‡ä»¶æ˜¯å¦æœ‰æ•ˆ"""
        try:
            if not os.path.exists(filepath):
                return False
            
            file_size = os.path.getsize(filepath)
            if file_size == 0:
                return False
            
            # å°è¯•è¯»å–æ–‡ä»¶å¤´
            with open(filepath, 'rb') as f:
                header = f.read(16)
                if len(header) < 4:
                    return False
            
            # å¦‚æœå®‰è£…äº†PILï¼Œå°è¯•æ‰“å¼€å›¾ç‰‡éªŒè¯
            try:
                from PIL import Image
                img = Image.open(filepath)
                img.verify()  # éªŒè¯å›¾ç‰‡å®Œæ•´æ€§
                return True
            except ImportError:
                # æ²¡æœ‰PILï¼Œåªè¿›è¡ŒåŸºæœ¬æ£€æŸ¥
                return True
            except Exception:
                return False
                
        except Exception:
            return False
    
    def cleanup(self):
        """æ¸…ç†ä¸´æ—¶æ–‡ä»¶"""
        if os.path.exists(self.temp_dir):
            shutil.rmtree(self.temp_dir)
            logger.info(f"å·²æ¸…ç†ä¸´æ—¶æ–‡ä»¶å¤¹: {self.temp_dir}")


class PPTImageReplacer:
    """PPTå›¾ç‰‡æ–‡æœ¬æ·»åŠ å™¨"""
    
    @staticmethod
    def add_ocr_text_to_slides(presentation_path: str, 
                             image_mapping: Dict,
                             output_path: str = None,
                             show_translation: bool = True):
        """
        åœ¨PPTé¡µé¢å³ä¾§æ·»åŠ OCRè¯†åˆ«çš„æ–‡æœ¬å’Œç¿»è¯‘ï¼Œä¿ç•™åŸå›¾ç‰‡
        
        Args:
            presentation_path: åŸPPTè·¯å¾„
            image_mapping: åŒ…å«all_textå’Œtranslated_textå­—æ®µçš„å›¾ç‰‡æ˜ å°„å…³ç³»
            output_path: è¾“å‡ºPPTè·¯å¾„ï¼ŒNoneåˆ™è¦†ç›–åŸæ–‡ä»¶
            show_translation: æ˜¯å¦æ˜¾ç¤ºç¿»è¯‘ç»“æœ
        """
        try:
            prs = Presentation(presentation_path)
            
            # éå†æ¯ä¸ªå¹»ç¯ç‰‡
            for slide_key, slide_info in image_mapping.items():
                slide_idx = slide_info["slide_number"] - 1
                slide = prs.slides[slide_idx]
                
                # æ”¶é›†è¯¥é¡µæ‰€æœ‰å›¾ç‰‡çš„OCRæ–‡æœ¬å’Œç¿»è¯‘ï¼ˆæ”¹è¿›ç‰ˆï¼‰
                slide_ocr_data = []
                image_count = 0
                
                for image_info in slide_info["images"]:
                    image_count += 1
                    filename = image_info.get("filename", f"å›¾ç‰‡{image_count}")
                    
                    # è·å–åŸæ–‡å’Œç¿»è¯‘çš„å­—å…¸
                    original_texts = image_info.get('all_text', {})
                    translated_texts = image_info.get('translated_text', {}) if show_translation else {}
                    
                    # æ„å»ºæ–‡æœ¬å¯¹åˆ—è¡¨ï¼ˆåŸæ–‡-è¯‘æ–‡é…å¯¹ï¼‰
                    text_pairs = []
                    
                    if original_texts:
                        for key, original_text in original_texts.items():
                            # ç¡®ä¿æ–‡æœ¬ä¸æ˜¯ç©ºæˆ–"0"
                            cleaned_text = original_text.strip() if original_text else ""
                            if cleaned_text and cleaned_text != "0":
                                pair = {
                                    'key': key,
                                    'original': cleaned_text,
                                    'translated': ''
                                }
                                
                                # æŸ¥æ‰¾å¯¹åº”çš„ç¿»è¯‘ï¼Œç¡®ä¿ç¿»è¯‘æ–‡æœ¬ä¹Ÿä¸æ˜¯ç©ºæˆ–"0"
                                if key in translated_texts:
                                    cleaned_translation = translated_texts[key].strip() if translated_texts[key] else ""
                                    if cleaned_translation and cleaned_translation != "0":
                                        pair['translated'] = cleaned_translation
                                
                                text_pairs.append(pair)
                    
                    # å¦‚æœæœ‰æ–‡æœ¬å¯¹ï¼Œæ·»åŠ åˆ°OCRæ•°æ®ä¸­
                    if text_pairs:
                        ocr_data = {
                            'filename': filename,
                            'text_pairs': text_pairs
                        }
                        slide_ocr_data.append(ocr_data)
                
                # å¦‚æœè¯¥é¡µæœ‰OCRæ•°æ®ï¼Œåˆ™æ·»åŠ æ–‡æœ¬æ¡†
                if slide_ocr_data:
                    PPTImageReplacer._add_paired_text_box_to_slide(
                        slide, slide_ocr_data, slide_info["slide_number"], show_translation
                    )
            
            # ä¿å­˜PPT
            save_path = output_path or presentation_path
            prs.save(save_path)
            logger.info(f"OCRç»“æœå’Œç¿»è¯‘å·²æ·»åŠ åˆ°PPT: {save_path}")
            
        except Exception as e:
            logger.error(f"æ·»åŠ OCRæ–‡æœ¬å’Œç¿»è¯‘æ—¶å‘ç”Ÿé”™è¯¯: {str(e)}")
            raise
    
    @staticmethod
    def _add_paired_text_box_to_slide(slide, ocr_data_list: List[Dict], slide_number: int, show_translation: bool):
        """åœ¨å¹»ç¯ç‰‡å³ä¾§æ·»åŠ æˆå¯¹æ˜¾ç¤ºçš„OCRæ–‡æœ¬æ¡†"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        try:
            # è·å–å¹»ç¯ç‰‡å°ºå¯¸
            try:
                presentation = slide.part.package.presentation_part.presentation
                slide_width = presentation.slide_width
                slide_height = presentation.slide_height
                
                # è½¬æ¢ä¸ºè‹±å¯¸
                slide_width_inches = slide_width / 914400
                slide_height_inches = slide_height / 914400
                
                logger.info(f"âœ… æˆåŠŸè·å–å¹»ç¯ç‰‡å°ºå¯¸: {slide_width_inches:.2f} x {slide_height_inches:.2f} è‹±å¯¸")
                
            except Exception as size_error:
                logger.warning(f"âš ï¸ æ— æ³•è·å–å¹»ç¯ç‰‡å°ºå¯¸ï¼Œä½¿ç”¨æ ‡å‡†å°ºå¯¸: {str(size_error)}")
                # ä½¿ç”¨æ ‡å‡†çš„16:9å¹»ç¯ç‰‡å°ºå¯¸ä½œä¸ºå¤‡ç”¨
                slide_width_inches = 13.33  # æ ‡å‡†å®½å±PPTå®½åº¦
                slide_height_inches = 7.5   # æ ‡å‡†å®½å±PPTé«˜åº¦
            
            # è®¡ç®—æ–‡æœ¬æ¡†ä½ç½®å’Œå¤§å°
            if show_translation:
                # æœ‰ç¿»è¯‘æ—¶ï¼Œæ–‡æœ¬æ¡†éœ€è¦æ›´å®½
                text_box_width = Inches(slide_width_inches * 0.45)
            else:
                # åªæœ‰åŸæ–‡æ—¶ï¼Œæ–‡æœ¬æ¡†ç¨çª„
                text_box_width = Inches(slide_width_inches * 0.35)
            
            # æ–‡æœ¬æ¡†é«˜åº¦ï¼šçº¦ä¸ºå¹»ç¯ç‰‡é«˜åº¦çš„85%
            text_box_height = Inches(slide_height_inches * 0.85)
            
            # æ–‡æœ¬æ¡†ä½ç½®ï¼šåœ¨å¹»ç¯ç‰‡å³ä¾§å¤–éƒ¨
            text_box_left = Inches(slide_width_inches + 0.1)
            text_box_top = Inches(0.3)
            
            # æ·»åŠ æ–‡æœ¬æ¡†
            textbox = slide.shapes.add_textbox(
                text_box_left, 
                text_box_top, 
                text_box_width, 
                text_box_height
            )
            
            # è®¾ç½®æ–‡æœ¬å†…å®¹
            text_frame = textbox.text_frame
            text_frame.clear()  # æ¸…é™¤é»˜è®¤æ®µè½
            
            # æ·»åŠ æ ‡é¢˜æ®µè½
            title_text = f"ç¬¬ {slide_number} é¡µ OCRè¯†åˆ«ç»“æœ"
            if show_translation:
                title_text += " (é€è¡Œå¯¹ç…§)"
            
            title_paragraph = text_frame.paragraphs[0]
            title_paragraph.text = title_text
            title_paragraph.font.size = Pt(14)
            title_paragraph.font.bold = True
            title_paragraph.font.color.rgb = RGBColor(0, 0, 0)  # é»‘è‰²
            title_paragraph.alignment = PP_ALIGN.CENTER
            
            # æ·»åŠ åˆ†éš”çº¿æ®µè½
            separator_paragraph = text_frame.add_paragraph()
            separator_paragraph.text = "=" * 40
            separator_paragraph.font.size = Pt(10)
            separator_paragraph.font.color.rgb = RGBColor(0, 0, 0)
            separator_paragraph.alignment = PP_ALIGN.CENTER
            
            # æ·»åŠ OCRå†…å®¹æ®µè½ï¼ˆæ–°çš„é€è¡Œå¯¹ç…§æ˜¾ç¤ºé€»è¾‘ï¼‰
            for i, ocr_data in enumerate(ocr_data_list, 1):
                filename = ocr_data['filename']
                text_pairs = ocr_data['text_pairs']
                
                # å›¾ç‰‡æ ‡é¢˜æ®µè½
                img_title_paragraph = text_frame.add_paragraph()
                img_title_paragraph.text = f"\nğŸ–¼ï¸ å›¾ç‰‡ {i}: {filename}"
                img_title_paragraph.font.size = Pt(12)
                img_title_paragraph.font.bold = True
                img_title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
                
                # éå†æ¯ä¸ªæ–‡æœ¬å¯¹ï¼Œäº¤é”™æ˜¾ç¤ºåŸæ–‡å’Œè¯‘æ–‡
                for j, text_pair in enumerate(text_pairs, 1):
                    original_text = text_pair['original']
                    translated_text = text_pair['translated']
                    
                    # åŸæ–‡éƒ¨åˆ†
                    if original_text:
                        # åŸæ–‡æ ‡ç­¾
                        original_label_paragraph = text_frame.add_paragraph()
                        original_label_paragraph.text = "ğŸ”¤ åŸæ–‡:"
                        original_label_paragraph.font.size = Pt(9)
                        original_label_paragraph.font.bold = True
                        original_label_paragraph.font.color.rgb = RGBColor(0, 102, 204)  # è“è‰²
                        
                        # åŸæ–‡å†…å®¹
                        original_content_paragraph = text_frame.add_paragraph()
                        original_content_paragraph.text = original_text
                        original_content_paragraph.font.size = Pt(9)
                        original_content_paragraph.font.color.rgb = RGBColor(0, 0, 0)
                    
                    # ç¿»è¯‘éƒ¨åˆ†ï¼ˆç´§è·Ÿåœ¨å¯¹åº”åŸæ–‡åé¢ï¼‰
                    if show_translation and translated_text:
                        # ç¿»è¯‘æ ‡ç­¾
                        translated_label_paragraph = text_frame.add_paragraph()
                        translated_label_paragraph.text = "ğŸŒ è¯‘æ–‡:"
                        translated_label_paragraph.font.size = Pt(9)
                        translated_label_paragraph.font.bold = True
                        translated_label_paragraph.font.color.rgb = RGBColor(204, 102, 0)  # æ©™è‰²
                        
                        # ç¿»è¯‘å†…å®¹
                        translated_content_paragraph = text_frame.add_paragraph()
                        translated_content_paragraph.text = translated_text
                        translated_content_paragraph.font.size = Pt(9)
                        translated_content_paragraph.font.color.rgb = RGBColor(0, 0, 0)
                    
                    # å¦‚æœä¸æ˜¯æœ€åä¸€ä¸ªæ–‡æœ¬å¯¹ï¼Œæ·»åŠ å°åˆ†éš”çº¿
                    if j < len(text_pairs):
                        small_sep_paragraph = text_frame.add_paragraph()
                        small_sep_paragraph.text = "- - -"
                        small_sep_paragraph.font.size = Pt(7)
                        small_sep_paragraph.font.color.rgb = RGBColor(180, 180, 180)  # æ›´æµ…çš„ç°è‰²
                        small_sep_paragraph.alignment = PP_ALIGN.CENTER
                
                # å›¾ç‰‡ä¹‹é—´çš„åˆ†éš”çº¿
                if i < len(ocr_data_list):  # ä¸æ˜¯æœ€åä¸€ä¸ªå›¾ç‰‡åˆ™æ·»åŠ åˆ†éš”çº¿
                    sep_paragraph = text_frame.add_paragraph()
                    sep_paragraph.text = "â”€" * 35
                    sep_paragraph.font.size = Pt(10)
                    sep_paragraph.font.color.rgb = RGBColor(128, 128, 128)  # ç°è‰²åˆ†éš”çº¿
                    sep_paragraph.alignment = PP_ALIGN.CENTER
            
            # è®¾ç½®æ–‡æœ¬æ¡†è¾¹æ¡†å’ŒèƒŒæ™¯
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = RGBColor(248, 248, 248)  # æµ…ç°è‰²èƒŒæ™¯
            
            # è®¾ç½®è¾¹æ¡†
            textbox.line.color.rgb = RGBColor(200, 200, 200)  # æµ…ç°è‰²è¾¹æ¡†
            textbox.line.width = Pt(1)
            
            # è®¾ç½®æ–‡æœ¬æ¡†è¾¹è·
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            
            # è®¾ç½®è‡ªåŠ¨æ¢è¡Œ
            text_frame.word_wrap = True
            
            # ç»Ÿè®¡æ˜¾ç¤ºçš„æ–‡æœ¬å¯¹æ•°é‡
            total_pairs = sum(len(ocr_data['text_pairs']) for ocr_data in ocr_data_list)
            translation_info = "å’Œç¿»è¯‘" if show_translation else ""
            logger.info(f"âœ… å·²åœ¨ç¬¬{slide_number}é¡µå³ä¾§æ·»åŠ OCRæ–‡æœ¬æ¡†{translation_info}")
            logger.info(f"   ğŸ“Š åŒ…å«{len(ocr_data_list)}å¼ å›¾ç‰‡ï¼Œå…±{total_pairs}ä¸ªæ–‡æœ¬å¯¹")
            
        except Exception as e:
            logger.error(f"åœ¨ç¬¬{slide_number}é¡µæ·»åŠ æˆå¯¹æ–‡æœ¬æ¡†æ—¶å‡ºé”™: {str(e)}")
            # å¤‡ç”¨ç®€å•æ–¹æ¡ˆ
            try:
                # ä½¿ç”¨å›ºå®šä½ç½®ä½œä¸ºå¤‡ç”¨æ–¹æ¡ˆ
                textbox = slide.shapes.add_textbox(
                    Inches(11),  # å›ºå®šä½ç½®ï¼šå³ä¾§
                    Inches(1),   # å›ºå®šä½ç½®ï¼šé¡¶éƒ¨
                    Inches(4),   # å›ºå®šå®½åº¦
                    Inches(6)    # å›ºå®šé«˜åº¦
                )
                
                text_frame = textbox.text_frame
                simple_content = f"ç¬¬{slide_number}é¡µOCRç»“æœ (é€è¡Œå¯¹ç…§):\n\n"
                
                for i, ocr_data in enumerate(ocr_data_list, 1):
                    simple_content += f"{i}. {ocr_data['filename']}\n"
                    for j, text_pair in enumerate(ocr_data['text_pairs'], 1):
                        simple_content += f"  {j}. åŸæ–‡: {text_pair['original']}\n"
                        if show_translation and text_pair['translated']:
                            simple_content += f"     è¯‘æ–‡: {text_pair['translated']}\n"
                        simple_content += "\n"
                
                text_frame.text = simple_content
                text_frame.paragraphs[0].font.size = Pt(9)
                text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                
                logger.warning(f"âš ï¸ ä½¿ç”¨å¤‡ç”¨æ–¹æ¡ˆåœ¨ç¬¬{slide_number}é¡µæ·»åŠ OCRæ–‡æœ¬")
                
            except Exception as backup_error:
                logger.error(f"âŒ å¤‡ç”¨æ–¹æ¡ˆä¹Ÿå¤±è´¥: {str(backup_error)}")


def ocr_controller(presentation_path: str, 
                  selected_pages: Optional[List[int]] = None, 
                  output_path: str = None,
                  enable_translation: bool = True,
                  target_language: str = "ä¸­æ–‡",
                  source_language: str = "è‹±æ–‡",
                  enable_text_splitting: str = "False") -> str:
    """
    OCRä¸»æ§åˆ¶å™¨ï¼šæå–å›¾ç‰‡ã€OCRè¯†åˆ«ã€æ–‡æœ¬è¡Œåˆ†å‰²ã€ç¿»è¯‘ã€å†™å›PPT
    
    Args:
        presentation_path: PPTæ–‡ä»¶è·¯å¾„
        selected_pages: çœŸå®é¡µç ï¼ˆç¬¬ä¸€é¡µä¸º1ï¼‰ï¼ŒNoneè¡¨ç¤ºå…¨ç¯‡å¤„ç†
        output_path: è¾“å‡ºæ–‡ä»¶è·¯å¾„
        enable_translation: æ˜¯å¦å¯ç”¨ç¿»è¯‘åŠŸèƒ½
        target_language: ç›®æ ‡è¯­è¨€
        source_language: æºè¯­è¨€
        enable_text_splitting: æ˜¯å¦å¯ç”¨æ–‡æœ¬è¡Œåˆ†å‰²å¤„ç†
        
    Returns:
        å¤„ç†åçš„PPTæ–‡ä»¶è·¯å¾„
    """
    extractor = None
    ocr_processor = None
    try:
        # éªŒè¯è¾“å…¥æ–‡ä»¶
        if not os.path.exists(presentation_path):
            raise FileNotFoundError(f"PPTæ–‡ä»¶ä¸å­˜åœ¨: {presentation_path}")

        # ä¿®æ­£selected_pagesä¸º0-basedç´¢å¼•
        prs = Presentation(presentation_path)
        total_slides = len(prs.slides)
        if selected_pages is not None:
            selected_pages = [p-1 for p in selected_pages if 1 <= p <= total_slides]
            if not selected_pages:
                logger.warning("selected_pageså‚æ•°æ— æœ‰æ•ˆé¡µç ï¼Œå°†å¤„ç†å…¨éƒ¨é¡µé¢")
                selected_pages = None

        # 1. æå–å›¾ç‰‡
        logger.info("=" * 50)
        logger.info("ğŸ” ç¬¬ä¸€æ­¥ï¼šæå–PPTä¸­çš„å›¾ç‰‡")
        logger.info("=" * 50)
        extractor = PPTImageExtractor()
        temp_dir, image_mapping = extractor.extract_images_from_slides(
            presentation_path, selected_pages
        )
        if not image_mapping:
            logger.warning("æœªæ‰¾åˆ°éœ€è¦å¤„ç†çš„å›¾ç‰‡")
            return presentation_path
        logger.info(f"âœ… å›¾ç‰‡æå–å®Œæˆï¼Œä¸´æ—¶ç›®å½•: {temp_dir}")

        # 2. è°ƒç”¨qwen-vl-ocrçš„apiè¿›è¡Œå›¾ç‰‡çš„æ–‡å­—æå–
        logger.info("\n" + "=" * 50)
        logger.info("ğŸ¤– ç¬¬äºŒæ­¥ï¼šè°ƒç”¨OCR QWEN APIè¿›è¡Œæ–‡æœ¬è¯†åˆ«")
        logger.info("=" * 50)

        folder_path = temp_dir  # æ›¿æ¢ä¸ºä½ çš„å›¾ç‰‡æ–‡ä»¶å¤¹è·¯å¾„
        json_path = os.path.join(temp_dir, "image_mapping.json")  # ä½¿ç”¨temp_dirä½œä¸ºæ–‡ä»¶å¤¹è·¯å¾„
        API_KEY = os.getenv("QWEN_API_KEY")

        # æ‰§è¡Œæ‰¹é‡å¤„ç†
        process_folder_with_mapping(folder_path, json_path, API_KEY)

        # 3. æ–‡æœ¬è¡Œåˆ†å‰²å¤„ç†ï¼ˆå¯é€‰ï¼‰
        if enable_text_splitting == "True_spliting":
            logger.info("\n" + "=" * 50)
            logger.info("âœ‚ï¸ ç¬¬ä¸‰æ­¥ï¼šæ–‡æœ¬è¡Œåˆ†å‰²å¤„ç†")
            logger.info("=" * 50)
            
            splitter = TextLineSplitter()
            split_success = splitter.process_json_file(json_path)
            
            if split_success:
                logger.info("âœ… æ–‡æœ¬è¡Œåˆ†å‰²å®Œæˆ")
            else:
                logger.warning("âš ï¸ æ–‡æœ¬è¡Œåˆ†å‰²å¤±è´¥ï¼Œå°†ä½¿ç”¨åŸå§‹æ–‡æœ¬ç»§ç»­å¤„ç†")
        else:
            logger.info("\n" + "=" * 50)
            logger.info("â­ï¸ ç¬¬ä¸‰æ­¥ï¼šè·³è¿‡æ–‡æœ¬è¡Œåˆ†å‰²å¤„ç†")
            logger.info("=" * 50)
            logger.info("âœ… ä¿æŒåŸå§‹æ–‡æœ¬æ ¼å¼")

        # 4. ç¿»è¯‘OCRè¯†åˆ«ç»“æœ
        step_num = 4 if enable_text_splitting != "False" else 3
        if enable_translation:
            logger.info("\n" + "=" * 50)
            logger.info(f"ğŸŒ ç¬¬{step_num}æ­¥ï¼šç¿»è¯‘è¯†åˆ«ç»“æœ ({source_language} â†’ {target_language})")
            logger.info("=" * 50)
            
            translation_success = TranslationManager.translate_ocr_results(
                temp_dir=temp_dir,
                target_language=target_language,
                source_language=source_language
            )
            
            if translation_success:
                logger.info(f"âœ… ç¿»è¯‘å®Œæˆ")
                
                # æ˜¾ç¤ºç¿»è¯‘æ‘˜è¦
                mapping_file = os.path.join(temp_dir, "image_mapping.json")
                summary = TranslationManager.get_translation_summary(mapping_file)
                if summary:
                    logger.info(f"ğŸ“Š ç¿»è¯‘æ‘˜è¦:")
                    logger.info(f"   - æ€»å›¾ç‰‡æ•°: {summary.get('total_images', 0)}")
                    logger.info(f"   - åŒ…å«æ–‡æœ¬çš„å›¾ç‰‡: {summary.get('images_with_text', 0)}")
                    logger.info(f"   - åŒ…å«ç¿»è¯‘çš„å›¾ç‰‡: {summary.get('images_with_translation', 0)}")
                    logger.info(f"   - ç¿»è¯‘æˆåŠŸç‡: {summary.get('translation_success_rate', 0):.1f}%")
            else:
                logger.warning("âš ï¸ ç¿»è¯‘å¤±è´¥ï¼Œå°†åªæ˜¾ç¤ºåŸæ–‡")
                enable_translation = False

        # 5. è¯»å–æ›´æ–°åçš„æ˜ å°„æ–‡ä»¶
        step_num = 5 if enable_translation else (4 if enable_text_splitting != "False" else 3)
        logger.info(f"\n" + "=" * 50)
        logger.info(f"ğŸ“– ç¬¬{step_num}æ­¥ï¼šè¯»å–å¤„ç†ç»“æœ")
        logger.info("=" * 50)
        mapping_file = os.path.join(temp_dir, "image_mapping.json")
        if not os.path.exists(mapping_file):
            raise Exception(f"æ˜ å°„æ–‡ä»¶ä¸å­˜åœ¨: {mapping_file}")
        with open(mapping_file, 'r', encoding='utf-8') as f:
            updated_mapping = json.load(f)
        logger.info("âœ… å¤„ç†ç»“æœè¯»å–å®Œæˆ")
        
        # ç»Ÿè®¡ç»“æœ
        ocr_count = 0
        translation_count = 0
        for slide_info in updated_mapping.values():
            for image_info in slide_info.get("images", []):
                if "all_text" in image_info and image_info["all_text"]:
                    ocr_count += 1
                    filename = image_info.get("filename", "æœªçŸ¥æ–‡ä»¶")
                    text_preview = str(list(image_info["all_text"].values())[0])[:50] + "..."
                    logger.info(f"   ğŸ“„ {filename}: {text_preview}")
                    
                    if enable_translation and "translated_text" in image_info and image_info["translated_text"]:
                        translation_count += 1
                        trans_preview = str(list(image_info["translated_text"].values())[0])[:50] + "..."
                        logger.info(f"   ğŸŒ ç¿»è¯‘: {trans_preview}")
        
        logger.info(f"ğŸ“Š å…±è¯†åˆ«å‡º {ocr_count} å¼ åŒ…å«æ–‡æœ¬çš„å›¾ç‰‡")
        if enable_translation:
            logger.info(f"ğŸ“Š å…±ç¿»è¯‘äº† {translation_count} å¼ å›¾ç‰‡çš„æ–‡æœ¬")

        # 6. å°†OCRç»“æœå’Œç¿»è¯‘æ·»åŠ åˆ°PPTå³ä¾§
        step_num = 6 if enable_translation else (5 if enable_text_splitting != "False" else 4)
        logger.info(f"\n" + "=" * 50)
        content_desc = "OCRè¯†åˆ«ç»“æœå’Œç¿»è¯‘" if enable_translation else "OCRè¯†åˆ«ç»“æœ"
        logger.info(f"ğŸ¨ ç¬¬{step_num}æ­¥ï¼šåœ¨PPTå³ä¾§æ·»åŠ {content_desc}")
        logger.info("=" * 50)
        
        PPTImageReplacer.add_ocr_text_to_slides(
            presentation_path=presentation_path,
            image_mapping=updated_mapping,
            output_path=output_path,
            show_translation=enable_translation
        )
        
        success_desc = "OCRç»“æœå’Œç¿»è¯‘" if enable_translation else "OCRç»“æœ"
        logger.info(f"âœ… {success_desc}å·²æ·»åŠ åˆ°PPTå³ä¾§")
        logger.info("\n" + "=" * 50)
        logger.info("ğŸ‰ å¤„ç†å®Œæˆï¼")
        logger.info("=" * 50)
        return output_path or presentation_path
        
    except Exception as e:
        error_msg = f"OCRæ§åˆ¶å™¨å¤„ç†å¤±è´¥: {str(e)}"
        logger.error(f"âŒ {error_msg}")
        return presentation_path
    finally:
        if extractor:
            logger.info("ğŸ§¹ æ¸…ç†ä¸´æ—¶æ–‡ä»¶...")
            extractor.cleanup()


# ä½¿ç”¨ç¤ºä¾‹
if __name__ == "__main__":
    # OCR API Token (æ›¿æ¢ä¸ºä½ çš„å®é™…token)
    OCR_TOKEN = "ä½ çš„OCR_API_TOKEN"
    QWEN_API_KEY = "ä½ çš„QWEN_API_KEY"
    
    # è®¾ç½®ç¯å¢ƒå˜é‡
    os.environ["OCR_TOKEN"] = OCR_TOKEN
    os.environ["QWEN_API_KEY"] = QWEN_API_KEY
    
    # PPTæ–‡ä»¶è·¯å¾„
    ppt_path = "example.pptx"
    
    # é€‰æ‹©è¦å¤„ç†çš„é¡µé¢ (1-basedç´¢å¼•ï¼ŒNoneè¡¨ç¤ºå…¨éƒ¨é¡µé¢)
    selected_pages = [1, 2, 3]  # å¤„ç†å‰3é¡µ
    # selected_pages = None  # å¤„ç†å…¨éƒ¨é¡µé¢
    
    # è¾“å‡ºæ–‡ä»¶è·¯å¾„
    output_path = "ocr_translated_result.pptx"
    
    # æ‰§è¡ŒOCRå’Œç¿»è¯‘å¤„ç†
    success = ocr_controller(
        presentation_path=ppt_path,
        selected_pages=selected_pages,
        output_path=output_path,
        enable_translation=True,  # å¯ç”¨ç¿»è¯‘
        target_language="è‹±æ–‡",   # ç›®æ ‡è¯­è¨€
        source_language="ä¸­æ–‡",   # æºè¯­è¨€
        enable_text_splitting=True  # å¯ç”¨æ–‡æœ¬è¡Œåˆ†å‰²
    )
    
    if success:
        logger.info("ğŸ‰ PPT OCRå’Œç¿»è¯‘å¤„ç†æˆåŠŸï¼")
        logger.info(f"ğŸ”— ç»“æœæ–‡ä»¶: {output_path}")
    else:
        logger.error("âŒ å¤„ç†å¤±è´¥ï¼Œè¯·æ£€æŸ¥é”™è¯¯ä¿¡æ¯ã€‚")
