#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
测试 PowerPoint 文本上下标格式
验证 run.font._element.set('baseline', 30000) 的效果
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def test_baseline_formatting():
    """测试 baseline 属性对文本格式的影响"""
    
    # 创建演示文稿
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 创建标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "上下标格式测试"
    
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.alignment = PP_ALIGN.CENTER
    title_run = title_paragraph.runs[0]
    title_run.font.size = Pt(24)
    title_run.font.bold = True
    
    # 测试不同的 baseline 值
    test_values = [
        (0, "正常文本 (baseline=0)"),
        (30000, "上标文本 (baseline=30000)"),
        (-30000, "下标文本 (baseline=-30000)"),
        (50000, "更大上标 (baseline=50000)"),
        (-50000, "更大下标 (baseline=-50000)"),
        (10000, "小上标 (baseline=10000)"),
        (-10000, "小下标 (baseline=-10000)")
    ]
    
    # 创建测试内容
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    
    # 添加说明
    p = content_frame.paragraphs[0]
    p.text = "测试不同的 baseline 值对文本位置的影响："
    p.font.size = Pt(16)
    p.font.bold = True
    
    # 添加测试文本
    for baseline_value, description in test_values:
        p = content_frame.add_paragraph()
        p.text = f"{description}: H2O + CO2 = H2CO3"
        p.font.size = Pt(14)
        
        # 设置 baseline 属性
        for run in p.runs:
            try:
                run.font._element.set('baseline', str(baseline_value))
                print(f"成功设置 baseline={baseline_value} 为: {description}")
            except Exception as e:
                print(f"设置 baseline={baseline_value} 失败: {e}")
    
    # 添加数学公式测试
    p = content_frame.add_paragraph()
    p.text = ""
    
    p = content_frame.add_paragraph()
    p.text = "数学公式测试："
    p.font.size = Pt(16)
    p.font.bold = True
    
    # 创建数学公式
    math_tests = [
        ("x² + y² = z²", 30000, "平方"),
        ("H₂O", -30000, "下标"),
        ("E = mc²", 30000, "平方"),
        ("log₁₀(x)", -30000, "下标"),
        ("2³ = 8", 30000, "立方")
    ]
    
    for formula, baseline, desc in math_tests:
        p = content_frame.add_paragraph()
        p.text = f"{formula} ({desc})"
        p.font.size = Pt(14)
        
        for run in p.runs:
            try:
                run.font._element.set('baseline', str(baseline))
                print(f"数学公式设置 baseline={baseline}: {formula}")
            except Exception as e:
                print(f"数学公式设置失败: {e}")
    
    # 保存文件
    output_file = "baseline_test.pptx"
    prs.save(output_file)
    print(f"\n测试文件已保存为: {output_file}")
    
    return output_file

def test_alternative_methods():
    """测试其他可能的上下标方法"""
    
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 创建标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "其他上下标方法测试"
    
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.alignment = PP_ALIGN.CENTER
    title_run = title_paragraph.runs[0]
    title_run.font.size = Pt(24)
    title_run.font.bold = True
    
    # 测试内容
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    
    # 测试方法1：使用 font.superscript 和 font.subscript
    p = content_frame.paragraphs[0]
    p.text = "方法1 - 使用 font.superscript/subscript:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = content_frame.add_paragraph()
    p.text = "正常文本"
    p.font.size = Pt(14)
    
    p = content_frame.add_paragraph()
    p.text = "上标文本"
    p.font.size = Pt(14)
    try:
        p.runs[0].font.superscript = True
        print("成功设置 font.superscript = True")
    except Exception as e:
        print(f"设置 font.superscript 失败: {e}")
    
    p = content_frame.add_paragraph()
    p.text = "下标文本"
    p.font.size = Pt(14)
    try:
        p.runs[0].font.subscript = True
        print("成功设置 font.subscript = True")
    except Exception as e:
        print(f"设置 font.subscript 失败: {e}")
    
    # 测试方法2：使用不同的 baseline 值
    p = content_frame.add_paragraph()
    p.text = ""
    
    p = content_frame.add_paragraph()
    p.text = "方法2 - 使用不同的 baseline 值:"
    p.font.size = Pt(16)
    p.font.bold = True
    
    baseline_tests = [
        (1000, "微调上标"),
        (5000, "小上标"),
        (10000, "中上标"),
        (30000, "大上标"),
        (-1000, "微调下标"),
        (-5000, "小下标"),
        (-10000, "中下标"),
        (-30000, "大下标")
    ]
    
    for baseline, desc in baseline_tests:
        p = content_frame.add_paragraph()
        p.text = f"{desc} (baseline={baseline})"
        p.font.size = Pt(12)
        
        for run in p.runs:
            try:
                run.font._element.set('baseline', str(baseline))
                print(f"设置 baseline={baseline}: {desc}")
            except Exception as e:
                print(f"设置 baseline={baseline} 失败: {e}")
    
    # 保存文件
    output_file = "alternative_methods_test.pptx"
    prs.save(output_file)
    print(f"\n替代方法测试文件已保存为: {output_file}")
    
    return output_file

def main():
    """主函数"""
    try:
        print("=== PowerPoint 上下标格式测试 ===")
        print("测试 run.font._element.set('baseline', 30000) 的效果")
        print()
        
        # 测试 baseline 方法
        print("1. 测试 baseline 属性方法...")
        file1 = test_baseline_formatting()
        
        print("\n" + "="*50)
        
        # 测试其他方法
        print("2. 测试其他上下标方法...")
        file2 = test_alternative_methods()
        
        print("\n" + "="*50)
        print("测试完成！")
        print(f"生成的文件:")
        print(f"  - {file1} (baseline 方法测试)")
        print(f"  - {file2} (其他方法测试)")
        print("\n请使用 PowerPoint 打开这些文件查看效果。")
        print("注意观察文本的垂直位置是否发生变化。")
        
    except ImportError:
        print("错误: 请先安装 python-pptx 库")
        print("安装命令: pip install python-pptx")
    except Exception as e:
        print(f"测试过程中发生错误: {e}")

if __name__ == "__main__":
    main()
