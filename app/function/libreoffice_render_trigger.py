"""
ä½¿ç”¨LibreOfficeå‘½ä»¤è¡Œè§¦å‘PPTæ–‡æœ¬æ¡†è‡ªé€‚åº”æ¸²æŸ“
é€šè¿‡è½¬æ¢PPTä¸ºPDFçš„è¿‡ç¨‹è§¦å‘å®Œæ•´æ¸²æŸ“ï¼Œä½¿æ–‡æœ¬æ¡†è‡ªé€‚åº”è®¾ç½®çœŸæ­£ç”Ÿæ•ˆ
"""
import os
import logging
import platform
import subprocess
import tempfile
import time
import shutil
from typing import Optional, Dict, Any

try:
    from pptx import Presentation
    from pptx.enum.text import MSO_AUTO_SIZE
    from pptx.util import Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.error("python-pptxä¸å¯ç”¨")

logger = logging.getLogger(__name__)


class LibreOfficeRenderTrigger:
    """LibreOfficeæ¸²æŸ“è§¦å‘å™¨"""

    def __init__(self):
        """åˆå§‹åŒ–LibreOfficeæ¸²æŸ“è§¦å‘å™¨"""
        # æ£€æŸ¥LibreOfficeæ˜¯å¦å¯ç”¨
        self.libreoffice_available = self._check_libreoffice()
        
        # åˆå§‹åŒ–ç»Ÿè®¡ä¿¡æ¯
        self.stats = {
            'total_textboxes': 0,  # æ€»æ–‡æœ¬æ¡†æ•°
            'autofit_set': 0,      # è®¾ç½®äº†è‡ªé€‚åº”çš„æ–‡æœ¬æ¡†æ•°
            'render_triggered': 0, # è§¦å‘æ¸²æŸ“æ¬¡æ•°
            'pdf_generated': 0,    # ç”Ÿæˆçš„PDFæ–‡ä»¶æ•°
            'pdf_deleted': 0,      # åˆ é™¤çš„PDFæ–‡ä»¶æ•°
            'pptx_generated': 0    # ç”Ÿæˆçš„PPTXæ–‡ä»¶æ•°
        }

    def _check_libreoffice(self) -> bool:
        """æ£€æŸ¥LibreOfficeæ˜¯å¦å¯ç”¨"""
        try:
            import platform

            # æ ¹æ®æ“ä½œç³»ç»Ÿå®šä¹‰å¯èƒ½çš„LibreOfficeå‘½ä»¤è·¯å¾„
            if platform.system() == "Windows":
                # Windowsä¸‹çš„å¸¸è§LibreOfficeå®‰è£…è·¯å¾„
                commands = [
                    r"C:\Program Files\LibreOffice\program\soffice.exe",
                    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                    r"soffice.exe",  # å¦‚æœåœ¨PATHä¸­
                    r"libreoffice.exe"
                ]
            else:
                # Linux/macOSä¸‹çš„å¸¸è§è·¯å¾„
                commands = [
                    'libreoffice',
                    'soffice',
                    '/usr/bin/libreoffice',
                    '/opt/libreoffice/program/soffice',
                    '/Applications/LibreOffice.app/Contents/MacOS/soffice'  # macOS
                ]

            for cmd in commands:
                try:
                    result = subprocess.run(
                        [cmd, '--version'],
                        capture_output=True,
                        text=True,
                        timeout=15
                    )
                    if result.returncode == 0:
                        self.libreoffice_cmd = cmd
                        version_info = result.stdout.strip()
                        logger.info(f"LibreOfficeå¯ç”¨: {cmd}")
                        logger.info(f"ç‰ˆæœ¬ä¿¡æ¯: {version_info}")
                        return True
                except (FileNotFoundError, subprocess.TimeoutExpired, OSError):
                    continue

            logger.warning("LibreOfficeä¸å¯ç”¨")
            logger.info("è¯·ç¡®ä¿LibreOfficeå·²æ­£ç¡®å®‰è£…")
            return False

        except Exception as e:
            logger.warning(f"æ£€æŸ¥LibreOfficeæ—¶å‡ºé”™: {e}")
            return False

    def process_ppt_with_render_trigger(self, ppt_path: str) -> bool:
        """
        å¤„ç†PPTå¹¶é€šè¿‡LibreOfficeè§¦å‘æ¸²æŸ“

        Args:
            ppt_path: PPTæ–‡ä»¶è·¯å¾„

        Returns:
            bool: å¤„ç†æ˜¯å¦æˆåŠŸ
        """
        if not os.path.exists(ppt_path):
            logger.error(f"æ–‡ä»¶ä¸å­˜åœ¨: {ppt_path}")
            return False

        if not self.libreoffice_available:
            logger.error("LibreOfficeä¸å¯ç”¨ï¼Œæ— æ³•è§¦å‘æ¸²æŸ“")
            return False

        logger.info(f"å¼€å§‹å¤„ç†PPTå¹¶è§¦å‘æ¸²æŸ“: {os.path.basename(ppt_path)}")

        try:
            # æ­¥éª¤1: ä½¿ç”¨python-pptxè®¾ç½®è‡ªé€‚åº”å±æ€§
            if not self._set_autofit_properties(ppt_path):
                logger.error("è®¾ç½®è‡ªé€‚åº”å±æ€§å¤±è´¥")
                return False

            # æ­¥éª¤2: ä½¿ç”¨LibreOfficeè½¬æ¢PDFè§¦å‘æ¸²æŸ“
            if not self._trigger_render_via_pdf_conversion(ppt_path):
                logger.error("LibreOfficeæ¸²æŸ“è§¦å‘å¤±è´¥")
                return False

            logger.info("âœ… PPTæ–‡æœ¬æ¡†è‡ªé€‚åº”å¤„ç†å®Œæˆï¼ˆåŒ…å«æ¸²æŸ“è§¦å‘ï¼‰")
            self._log_stats()
            return True

        except Exception as e:
            logger.error(f"å¤„ç†PPTæ—¶å‡ºé”™: {e}")
            return False

    def _set_autofit_properties(self, ppt_path: str) -> bool:
        """è®¾ç½®æ–‡æœ¬æ¡†è‡ªé€‚åº”å±æ€§"""
        if not PPTX_AVAILABLE:
            logger.error("python-pptxåº“ä¸å¯ç”¨")
            return False

        try:
            logger.info("æ­¥éª¤1: è®¾ç½®æ–‡æœ¬æ¡†è‡ªé€‚åº”å±æ€§")
            prs = Presentation(ppt_path)

            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        self.stats['total_textboxes'] += 1
                        text_frame = shape.text_frame

                        # æ£€æŸ¥æ˜¯å¦æœ‰å®é™…å†…å®¹
                        has_content = any(
                            paragraph.text.strip()
                            for paragraph in text_frame.paragraphs
                        )

                        if has_content:
                            # è®¾ç½®æ–‡æœ¬å¤§å°é€‚åº”æ–‡æœ¬æ¡†
                            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                            text_frame.word_wrap = True

                            # ä¼˜åŒ–è¾¹è·
                            text_frame.margin_left = Pt(3)
                            text_frame.margin_right = Pt(3)
                            text_frame.margin_top = Pt(2)
                            text_frame.margin_bottom = Pt(2)

                            self.stats['autofit_set'] += 1

                    elif shape.has_table:
                        # å¤„ç†è¡¨æ ¼ä¸­çš„æ–‡æœ¬æ¡†
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                self.stats['total_textboxes'] += 1
                                text_frame = cell.text_frame

                                has_content = any(
                                    paragraph.text.strip()
                                    for paragraph in text_frame.paragraphs
                                )

                                if has_content:
                                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                                    text_frame.word_wrap = True
                                    self.stats['autofit_set'] += 1

            # ä¿å­˜è®¾ç½®
            prs.save(ppt_path)
            logger.info(f"è®¾ç½®äº† {self.stats['autofit_set']} ä¸ªæ–‡æœ¬æ¡†çš„è‡ªé€‚åº”å±æ€§")
            return True

        except Exception as e:
            logger.error(f"è®¾ç½®è‡ªé€‚åº”å±æ€§æ—¶å‡ºé”™: {e}")
            return False

    def _trigger_render_via_pdf_conversion(self, ppt_path: str) -> bool:
        """é€šè¿‡ODPä¸­è½¬è½¬æ¢ä¸ºPDFè§¦å‘æ¸²æŸ“"""
        try:
            logger.info("æ­¥éª¤2: é€šè¿‡LibreOffice ODPä¸­è½¬è½¬æ¢ä¸ºPDFè§¦å‘æ¸²æŸ“")

            # åˆ›å»ºä¸´æ—¶ç›®å½•ç”¨äºä¸­é—´æ–‡ä»¶å’ŒPDFè¾“å‡º
            with tempfile.TemporaryDirectory() as temp_dir:
                # ç¡®ä¿è·¯å¾„æ ¼å¼æ­£ç¡®
                temp_dir = os.path.abspath(temp_dir)
                ppt_path = os.path.abspath(ppt_path)
                ppt_filename = os.path.basename(ppt_path)
                ppt_name = os.path.splitext(ppt_filename)[0]

                logger.debug(f"è¾“å…¥æ–‡ä»¶: {ppt_path}")
                logger.debug(f"è¾“å‡ºç›®å½•: {temp_dir}")

                # æ­¥éª¤2.1: å…ˆè½¬æ¢ä¸ºODPæ ¼å¼
                cmd_to_odp = [
                    self.libreoffice_cmd,
                    '--headless',           # æ— å¤´æ¨¡å¼
                    '--invisible',          # ä¸å¯è§
                    '--nodefault',          # ä¸ä½¿ç”¨é»˜è®¤è®¾ç½®
                    '--nolockcheck',        # ä¸æ£€æŸ¥é”å®š
                    '--nologo',             # ä¸æ˜¾ç¤ºlogo
                    '--norestore',          # ä¸æ¢å¤
                    '--convert-to', 'odp',  # è½¬æ¢ä¸ºODP
                    '--outdir', temp_dir,   # è¾“å‡ºç›®å½•
                    ppt_path                # è¾“å…¥æ–‡ä»¶
                ]

                # åœ¨Windowsä¸‹ï¼Œå¦‚æœè·¯å¾„åŒ…å«ç©ºæ ¼ï¼Œéœ€è¦ç‰¹æ®Šå¤„ç†
                if platform.system() == "Windows":
                    # å¯¹äºWindowsï¼Œç¡®ä¿è·¯å¾„è¢«æ­£ç¡®å¼•ç”¨
                    cmd_str_odp = f'"{self.libreoffice_cmd}" --headless --invisible --nodefault --nolockcheck --nologo --norestore --convert-to odp --outdir "{temp_dir}" "{ppt_path}"'
                    logger.debug(f"Windows ODPè½¬æ¢å‘½ä»¤: {cmd_str_odp}")
                    
                    # Windowsä¸‹ä½¿ç”¨shell=Trueæ‰§è¡Œå‘½ä»¤å­—ç¬¦ä¸²
                    result_odp = subprocess.run(
                        cmd_str_odp,
                        shell=True,
                        capture_output=True,
                        text=True,
                        timeout=120,  # 2åˆ†é’Ÿè¶…æ—¶
                        cwd=temp_dir  # è®¾ç½®å·¥ä½œç›®å½•
                    )
                else:
                    # Linux/macOSä¸‹ä½¿ç”¨åˆ—è¡¨å½¢å¼
                    logger.debug(f"æ‰§è¡ŒODPè½¬æ¢å‘½ä»¤: {' '.join(cmd_to_odp)}")
                    result_odp = subprocess.run(
                        cmd_to_odp,
                        capture_output=True,
                        text=True,
                        timeout=120,
                        cwd=temp_dir
                    )
                
                # æ£€æŸ¥ODPè½¬æ¢æ˜¯å¦æˆåŠŸ
                if result_odp.returncode != 0:
                    logger.error(f"è½¬æ¢ä¸ºODPå¤±è´¥: {result_odp.stderr}")
                    return False
                
                # æ£€æŸ¥ODPæ–‡ä»¶æ˜¯å¦ç”Ÿæˆ
                odp_file = os.path.join(temp_dir, f"{ppt_name}.odp")
                if not os.path.exists(odp_file):
                    logger.error("æœªèƒ½ç”ŸæˆODPæ–‡ä»¶")
                    return False
                
                logger.info(f"æˆåŠŸè½¬æ¢ä¸ºODPæ ¼å¼: {odp_file}")

                # æ­¥éª¤2.2: ä»ODPè½¬æ¢ä¸ºPDFæ ¼å¼
                cmd_to_pdf = [
                    self.libreoffice_cmd,
                    '--headless',           # æ— å¤´æ¨¡å¼
                    '--invisible',          # ä¸å¯è§
                    '--nodefault',          # ä¸ä½¿ç”¨é»˜è®¤è®¾ç½®
                    '--nolockcheck',        # ä¸æ£€æŸ¥é”å®š
                    '--nologo',             # ä¸æ˜¾ç¤ºlogo
                    '--norestore',          # ä¸æ¢å¤
                    '--convert-to', 'pdf:writer_pdf_Export',  # è½¬æ¢ä¸ºPDFå¹¶æŒ‡å®šå¯¼å‡ºè¿‡æ»¤å™¨
                    '--outdir', temp_dir,   # è¾“å‡ºç›®å½•
                    odp_file                # è¾“å…¥æ–‡ä»¶
                ]

                # åœ¨Windowsä¸‹ï¼Œå¦‚æœè·¯å¾„åŒ…å«ç©ºæ ¼ï¼Œéœ€è¦ç‰¹æ®Šå¤„ç†
                if platform.system() == "Windows":
                    # å¯¹äºWindowsï¼Œç¡®ä¿è·¯å¾„è¢«æ­£ç¡®å¼•ç”¨
                    cmd_str_pdf = f'"{self.libreoffice_cmd}" --headless --invisible --nodefault --nolockcheck --nologo --norestore --convert-to pdf:writer_pdf_Export --outdir "{temp_dir}" "{odp_file}"'
                    logger.debug(f"Windows PDFè½¬æ¢å‘½ä»¤: {cmd_str_pdf}")
                    
                    # æ‰§è¡Œè½¬æ¢å‘½ä»¤
                    start_time = time.time()
                    
                    # Windowsä¸‹ä½¿ç”¨shell=Trueæ‰§è¡Œå‘½ä»¤å­—ç¬¦ä¸²
                    result = subprocess.run(
                        cmd_str_pdf,
                        shell=True,
                        capture_output=True,
                        text=True,
                        timeout=120,  # 2åˆ†é’Ÿè¶…æ—¶
                        cwd=temp_dir  # è®¾ç½®å·¥ä½œç›®å½•
                    )
                else:
                    # Linux/macOSä¸‹ä½¿ç”¨åˆ—è¡¨å½¢å¼
                    logger.debug(f"æ‰§è¡ŒPDFè½¬æ¢å‘½ä»¤: {' '.join(cmd_to_pdf)}")
                    
                    # æ‰§è¡Œè½¬æ¢å‘½ä»¤
                    start_time = time.time()
                    
                    result = subprocess.run(
                        cmd_to_pdf,
                        capture_output=True,
                        text=True,
                        timeout=120,
                        cwd=temp_dir
                    )

                end_time = time.time()
                conversion_time = end_time - start_time

                logger.debug(f"LibreOfficeè½¬æ¢è€—æ—¶: {conversion_time:.2f}ç§’")
                logger.debug(f"è¿”å›ç : {result.returncode}")

                if result.stdout:
                    logger.debug(f"æ ‡å‡†è¾“å‡º: {result.stdout}")
                if result.stderr:
                    logger.debug(f"æ ‡å‡†é”™è¯¯: {result.stderr}")

                # æ£€æŸ¥è½¬æ¢æ˜¯å¦æˆåŠŸ
                if result.returncode == 0:
                    # æŸ¥æ‰¾ç”Ÿæˆçš„PDFæ–‡ä»¶
                    try:
                        pdf_file = os.path.join(temp_dir, f"{ppt_name}.pdf")
                        if os.path.exists(pdf_file):
                            pdf_size = os.path.getsize(pdf_file)

                            logger.info(f"âœ… PDFè½¬æ¢æˆåŠŸ: {os.path.basename(pdf_file)} ({pdf_size} bytes)")
                            logger.info("ğŸ¯ PPTå·²è¢«LibreOfficeå®Œæ•´æ¸²æŸ“ï¼Œæ–‡æœ¬æ¡†è‡ªé€‚åº”è®¾ç½®å·²ç”Ÿæ•ˆ")

                            self.stats['render_triggered'] = 1
                            self.stats['pdf_generated'] = 1
                            self.stats['pdf_deleted'] = 1  # PDFåœ¨ä¸´æ—¶ç›®å½•ä¸­ä¼šè‡ªåŠ¨åˆ é™¤

                            # æ­¥éª¤2.3: ä»ODPè½¬æ¢å›PPTXæ ¼å¼
                            logger.info("æ­¥éª¤2.3: ä»ODPè½¬æ¢å›PPTXæ ¼å¼å¹¶è¦†ç›–åŸæ–‡ä»¶")
                            pptx_result = self._convert_odp_to_pptx(odp_file, temp_dir, ppt_path)
                            if pptx_result:
                                logger.info("âœ… å®Œæ•´æ¸²æŸ“æµç¨‹æˆåŠŸ: ODP -> PDF -> PPTX")
                                self.stats['pptx_generated'] = 1
                            else:
                                logger.warning("âš ï¸ PDFæ¸²æŸ“æˆåŠŸï¼Œä½†PPTXè½¬æ¢å¤±è´¥")

                            return True
                        else:
                            logger.warning("PDFè½¬æ¢å‘½ä»¤æˆåŠŸä½†æœªæ‰¾åˆ°PDFæ–‡ä»¶")
                            logger.debug(f"ä¸´æ—¶ç›®å½•å†…å®¹: {os.listdir(temp_dir)}")
                            return False
                    except Exception as e:
                        logger.error(f"æ£€æŸ¥PDFæ–‡ä»¶æ—¶å‡ºé”™: {e}")
                        return False
                else:
                    logger.error(f"LibreOfficeè½¬æ¢å¤±è´¥ï¼Œè¿”å›ç : {result.returncode}")
                    if result.stderr:
                        logger.error(f"é”™è¯¯ä¿¡æ¯: {result.stderr}")
                    if result.stdout:
                        logger.error(f"è¾“å‡ºä¿¡æ¯: {result.stdout}")
                    return False

        except subprocess.TimeoutExpired:
            logger.error("LibreOfficeè½¬æ¢è¶…æ—¶ï¼ˆ120ç§’ï¼‰")
            return False
        except Exception as e:
            logger.error(f"PDFè½¬æ¢è§¦å‘æ¸²æŸ“å¤±è´¥: {e}")
            import traceback
            logger.debug(f"è¯¦ç»†é”™è¯¯: {traceback.format_exc()}")
            return False

    def _convert_odp_to_pptx(self, odp_file: str, output_dir: str, original_pptx_path: str = None) -> str:
        """
        å°†ODPæ–‡ä»¶è½¬æ¢å›PPTXæ ¼å¼
        
        Args:
            odp_file: ODPæ–‡ä»¶è·¯å¾„
            output_dir: è¾“å‡ºç›®å½•
            original_pptx_path: åŸå§‹PPTXæ–‡ä»¶è·¯å¾„ï¼Œå¦‚æœæä¾›åˆ™ä¼šè¦†ç›–è¯¥æ–‡ä»¶
            
        Returns:
            str: è½¬æ¢åçš„PPTXæ–‡ä»¶è·¯å¾„ï¼Œå¤±è´¥åˆ™è¿”å›ç©ºå­—ç¬¦ä¸²
        """
        try:
            logger.info(f"å°†ODPè½¬æ¢å›PPTXæ ¼å¼: {os.path.basename(odp_file)}")
            
            # æ„å»ºè½¬æ¢å‘½ä»¤
            cmd = [
                self.libreoffice_cmd,
                '--headless',
                '--invisible',
                '--nodefault',
                '--nolockcheck',
                '--nologo',
                '--norestore',
                '--convert-to', 'pptx:Impress MS PowerPoint 2007 XML',  # æ˜¾å¼æŒ‡å®šå¯¼å‡ºæ ¼å¼
                '--outdir', output_dir,
                odp_file
            ]
            
            # åœ¨Windowsä¸‹ï¼Œå¦‚æœè·¯å¾„åŒ…å«ç©ºæ ¼ï¼Œéœ€è¦ç‰¹æ®Šå¤„ç†
            if platform.system() == "Windows":
                cmd_str = f'"{self.libreoffice_cmd}" --headless --invisible --nodefault --nolockcheck --nologo --norestore --convert-to "pptx:Impress MS PowerPoint 2007 XML" --outdir "{output_dir}" "{odp_file}"'
                logger.debug(f"Windows PPTXè½¬æ¢å‘½ä»¤: {cmd_str}")
                
                # Windowsä¸‹ä½¿ç”¨shell=Trueæ‰§è¡Œå‘½ä»¤å­—ç¬¦ä¸²
                result = subprocess.run(
                    cmd_str,
                    shell=True,
                    capture_output=True,
                    text=True,
                    timeout=120,
                    cwd=output_dir
                )
            else:
                # Linux/macOSä¸‹ä½¿ç”¨åˆ—è¡¨å½¢å¼
                logger.debug(f"æ‰§è¡ŒPPTXè½¬æ¢å‘½ä»¤: {' '.join(cmd)}")
                result = subprocess.run(
                    cmd,
                    capture_output=True,
                    text=True,
                    timeout=120,
                    cwd=output_dir
                )
            
            if result.returncode == 0:
                # æŸ¥æ‰¾ç”Ÿæˆçš„PPTXæ–‡ä»¶
                ppt_name = os.path.splitext(os.path.basename(odp_file))[0]
                pptx_file = os.path.join(output_dir, f"{ppt_name}.pptx")
                
                if os.path.exists(pptx_file):
                    logger.info(f"ODPæˆåŠŸè½¬æ¢ä¸ºPPTX: {pptx_file}")
                    
                    # å¦‚æœæä¾›äº†åŸå§‹PPTXè·¯å¾„ï¼Œåˆ™è¦†ç›–å®ƒ
                    if original_pptx_path:
                        try:
                            shutil.copyfile(pptx_file, original_pptx_path)
                            logger.info(f"æ¸²æŸ“åçš„PPTXå·²è¦†ç›–åŸæ–‡ä»¶: {original_pptx_path}")
                            return original_pptx_path
                        except Exception as e:
                            logger.error(f"æ— æ³•è¦†ç›–åŸPPTXæ–‡ä»¶: {e}")
                            return pptx_file
                    
                    return pptx_file
                else:
                    logger.warning("PPTXè½¬æ¢å‘½ä»¤æˆåŠŸä½†æœªæ‰¾åˆ°PPTXæ–‡ä»¶")
                    logger.debug(f"ç›®å½•å†…å®¹: {os.listdir(output_dir)}")
                    return ""
            else:
                logger.error(f"PPTXè½¬æ¢å¤±è´¥ï¼Œè¿”å›ç : {result.returncode}")
                if result.stderr:
                    logger.error(f"é”™è¯¯ä¿¡æ¯: {result.stderr}")
                return ""
                
        except Exception as e:
            logger.error(f"ODPè½¬æ¢ä¸ºPPTXå‡ºé”™: {e}")
            import traceback
            logger.debug(f"è¯¦ç»†é”™è¯¯: {traceback.format_exc()}")
            return ""

    def _log_stats(self):
        """è¾“å‡ºç»Ÿè®¡ä¿¡æ¯"""
        logger.info("=" * 50)
        logger.info("LibreOfficeæ¸²æŸ“è§¦å‘ç»Ÿè®¡")
        logger.info("=" * 50)
        logger.info(f"æ€»æ–‡æœ¬æ¡†æ•°: {self.stats['total_textboxes']}")
        logger.info(f"è®¾ç½®è‡ªé€‚åº”: {self.stats['autofit_set']}")
        logger.info(f"æ¸²æŸ“è§¦å‘: {self.stats['render_triggered']}")
        logger.info(f"PDFç”Ÿæˆ: {self.stats['pdf_generated']}")
        logger.info(f"PDFåˆ é™¤: {self.stats['pdf_deleted']}")

        if self.stats['total_textboxes'] > 0:
            success_rate = (self.stats['autofit_set'] / self.stats['total_textboxes']) * 100
            logger.info(f"å¤„ç†æˆåŠŸç‡: {success_rate:.1f}%")

    def get_stats(self) -> Dict[str, Any]:
        """è·å–ç»Ÿè®¡ä¿¡æ¯"""
        return {
            'total_textboxes': self.stats.get('total_textboxes', 0),
            'autofit_set': self.stats.get('autofit_set', 0),
            'render_triggered': self.stats.get('render_triggered', 0),
            'pdf_generated': self.stats.get('pdf_generated', 0),
            'pdf_deleted': self.stats.get('pdf_deleted', 0),
            'pptx_generated': self.stats.get('pptx_generated', 0),
            'libreoffice_available': self.libreoffice_available,
            'libreoffice_cmd': getattr(self, 'libreoffice_cmd', None)
        }


def libreoffice_trigger_ppt_autofit(ppt_path: str) -> bool:
    """
    ä½¿ç”¨LibreOfficeè§¦å‘PPTæ–‡æœ¬æ¡†è‡ªé€‚åº”æ¸²æŸ“

    Args:
        ppt_path: PPTæ–‡ä»¶è·¯å¾„

    Returns:
        bool: å¤„ç†æ˜¯å¦æˆåŠŸ
    """
    trigger = LibreOfficeRenderTrigger()
    return trigger.process_ppt_with_render_trigger(ppt_path)


def install_libreoffice_instructions():
    """LibreOfficeå®‰è£…è¯´æ˜"""
    instructions = """
    LibreOfficeå®‰è£…è¯´æ˜:

    Ubuntu/Debian:
        sudo apt-get update
        sudo apt-get install -y libreoffice --no-install-recommends

    CentOS/RHEL:
        sudo yum install -y libreoffice-headless
        # æˆ–è€…
        sudo dnf install -y libreoffice-headless

    Alpine Linux:
        apk add --no-cache libreoffice

    Docker Dockerfileç¤ºä¾‹:
        FROM python:3.9-slim
        RUN apt-get update && \\
            apt-get install -y libreoffice --no-install-recommends && \\
            rm -rf /var/lib/apt/lists/*

    éªŒè¯å®‰è£…:
        libreoffice --version
        # æˆ–è€…
        soffice --version

    æ³¨æ„äº‹é¡¹:
        1. ç¡®ä¿æœ‰è¶³å¤Ÿçš„ç£ç›˜ç©ºé—´ç”¨äºä¸´æ—¶PDFæ–‡ä»¶
        2. è½¬æ¢è¿‡ç¨‹å¯èƒ½éœ€è¦1-2åˆ†é’Ÿï¼Œå–å†³äºPPTå¤§å°
        3. æ— å¤´æ¨¡å¼ä¸éœ€è¦å›¾å½¢ç•Œé¢
        4. è½¬æ¢å®ŒæˆåPDFæ–‡ä»¶ä¼šè‡ªåŠ¨åˆ é™¤
    """
    return instructions


# æµ‹è¯•å‡½æ•°
def test_libreoffice_availability():
    """æµ‹è¯•LibreOfficeå¯ç”¨æ€§"""
    trigger = LibreOfficeRenderTrigger()

    if trigger.libreoffice_available:
        print(f"âœ… LibreOfficeå¯ç”¨: {trigger.libreoffice_cmd}")
        return True
    else:
        print("âŒ LibreOfficeä¸å¯ç”¨")
        print("\nå®‰è£…è¯´æ˜:")
        print(install_libreoffice_instructions())
        return False


if __name__ == "__main__":
    # æµ‹è¯•LibreOfficeå¯ç”¨æ€§
    test_libreoffice_availability()
