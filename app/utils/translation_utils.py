"""
翻译相关工具函数
整合了项目中重复的翻译工具函数
"""
import json
import re
import logging
from typing import Dict, List, Any

logger = logging.getLogger(__name__)


def build_map(data: List[Dict[str, str]]) -> Dict[str, str]:
    """
    构建翻译映射字典
    
    Args:
        data: 解析后的翻译数据列表，每个元素包含source_language和target_language
        
    Returns:
        翻译映射字典 {原文: 译文}
    """
    result = {}
    for item in data:
        if isinstance(item, dict) and 'source_language' in item and 'target_language' in item:
            source = item['source_language']
            target = item['target_language']
            result[source] = target
            # 添加调试信息
            # print(f"构建映射: {repr(source)} -> {repr(target)}")
    return result


def parse_formatted_text(text: str) -> List[Dict[str, str]]:
    """
    解析JSON格式的翻译结果
    
    Args:
        text: JSON格式的翻译结果字符串
        
    Returns:
        解析后的翻译数据列表
    """
    try:
        # 尝试直接解析JSON
        return json.loads(text)
    except json.JSONDecodeError:
        # 如果直接解析失败，尝试修复格式
        return re_parse_formatted_text(text)


def re_parse_formatted_text(text: str) -> List[Dict[str, str]]:
    """
    重新解析格式化文本（处理格式不规范的JSON）
    
    Args:
        text: 可能格式不规范的JSON字符串
        
    Returns:
        解析后的翻译数据列表
    """
    try:
        # 移除可能的代码块标记
        text = re.sub(r'``json\s*', '', text)
        text = re.sub(r'```\s*$', '', text)
        
        # 尝试修复常见的JSON格式问题
        text = text.strip()
        
        # 如果不是以[开头，尝试添加
        if not text.startswith('['):
            text = '[' + text
        
        # 如果不是以]结尾，尝试添加
        if not text.endswith(']'):
            text = text + ']'
        
        # 尝试解析
        result = json.loads(text)
        
        # 验证结果格式
        if isinstance(result, list):
            valid_items = []
            for item in result:
                if isinstance(item, dict) and 'source_language' in item and 'target_language' in item:
                    valid_items.append(item)
            return valid_items
        
        return []
        
    except Exception as e:
        logger.error(f"重新解析格式化文本失败: {str(e)}")
        logger.error(f"原始文本: {text[:200]}...")
        return []


def build_english_to_chinese_map(data: List[Dict[str, str]]) -> Dict[str, str]:
    """
    构建英文到中文的翻译映射
    
    Args:
        data: 翻译数据列表
        
    Returns:
        英文到中文的映射字典
    """
    result = {}
    for item in data:
        if isinstance(item, dict):
            # 检查是否包含英文和中文字段
            english_text = item.get('source_language', '')
            chinese_text = item.get('target_language', '')
            
            if english_text and chinese_text:
                result[english_text] = chinese_text
    
    return result


def extract_text_from_pptx(file_path: str) -> str:
    """
    从PPTX文件中提取文本
    
    Args:
        file_path: PPTX文件路径
        
    Returns:
        提取的文本内容
    """
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        text_content = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
        
        return '\n'.join(text_content)
        
    except Exception as e:
        logger.error(f"从PPTX文件提取文本失败: {str(e)}")
        return ""


def clean_translation_text(text: str) -> str:
    """
    清理翻译文本中的特殊字符
    
    Args:
        text: 原始翻译文本
        
    Returns:
        清理后的文本
    """
    if not text:
        return text
    
    # 移除退格符和垂直制表符
    text = text.replace('\b', '')
    text = text.replace('\x0b', '')
    
    # 移除其他控制字符
    text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', text)
    
    return text.strip()


def validate_translation_result(result: Any) -> bool:
    """
    验证翻译结果的格式是否正确
    
    Args:
        result: 翻译结果
        
    Returns:
        是否为有效的翻译结果
    """
    if not isinstance(result, (list, dict)):
        return False
    
    if isinstance(result, dict):
        # 如果是字典，检查是否包含必要的键
        return len(result) > 0
    
    if isinstance(result, list):
        # 如果是列表，检查每个元素是否包含必要的键
        for item in result:
            if not isinstance(item, dict):
                return False
            if 'source_language' not in item or 'target_language' not in item:
                return False
        return len(result) > 0
    
    return False


def merge_translation_results(*results: Dict[str, str]) -> Dict[str, str]:
    """
    合并多个翻译结果字典
    
    Args:
        *results: 多个翻译结果字典
        
    Returns:
        合并后的翻译字典
    """
    merged = {}
    for result in results:
        if isinstance(result, dict):
            merged.update(result)
    return merged


def filter_translation_by_length(translation_dict: Dict[str, str], 
                                min_length: int = 1, 
                                max_length: int = 1000) -> Dict[str, str]:
    """
    根据文本长度过滤翻译结果
    
    Args:
        translation_dict: 翻译字典
        min_length: 最小长度
        max_length: 最大长度
        
    Returns:
        过滤后的翻译字典
    """
    filtered = {}
    for source, target in translation_dict.items():
        if min_length <= len(source) <= max_length and min_length <= len(target) <= max_length:
            filtered[source] = target
    return filtered


def get_translation_statistics(translation_dict: Dict[str, str]) -> Dict[str, Any]:
    """
    获取翻译结果的统计信息
    
    Args:
        translation_dict: 翻译字典
        
    Returns:
        统计信息字典
    """
    if not translation_dict:
        return {
            'total_pairs': 0,
            'avg_source_length': 0,
            'avg_target_length': 0,
            'max_source_length': 0,
            'max_target_length': 0,
            'min_source_length': 0,
            'min_target_length': 0
        }
    
    source_lengths = [len(source) for source in translation_dict.keys()]
    target_lengths = [len(target) for target in translation_dict.values()]
    
    return {
        'total_pairs': len(translation_dict),
        'avg_source_length': sum(source_lengths) / len(source_lengths),
        'avg_target_length': sum(target_lengths) / len(target_lengths),
        'max_source_length': max(source_lengths),
        'max_target_length': max(target_lengths),
        'min_source_length': min(source_lengths),
        'min_target_length': min(target_lengths)
    }
