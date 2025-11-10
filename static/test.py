from pptx import Presentation
from pptx.util import Pt

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])
tx_box = slide.shapes.add_textbox(left=Pt(50), top=Pt(50), width=Pt(400), height=Pt(100))
p      = tx_box.text_frame.paragraphs[0]

run1 = p.add_run()               # 普通文本
run1.text = "H"
run2 = p.add_run()               # 上标
run2.text = "2"
run2.font.size = Pt(14)
run2.font._element.set('baseline', '30000')  # +30 % = superscript
run3 = p.add_run()
run3.text = "O"

prs.save("superscript_demo.pptx")
